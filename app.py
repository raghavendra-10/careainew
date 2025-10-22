import os
import numpy as np
import pandas as pd
import json
import time
import uuid
import threading
import traceback
import requests
import re
from datetime import datetime, timezone
from flask import Flask, request, jsonify, Response, send_file
from flask_cors import CORS
from dotenv import load_dotenv
from threading import Thread, Lock
import asyncio
from waitress import serve
from openai import OpenAI
from tenacity import retry, stop_after_attempt, wait_exponential, retry_if_exception_type
import zipfile
import io

# Load environment variables
load_dotenv()

# Phase 1: Import Redis for progress tracking only
try:
    from redis_manager import redis_manager
    REDIS_AVAILABLE = redis_manager.is_connected()
    print("✅ Redis progress tracking initialized" if REDIS_AVAILABLE else "⚠️ Redis unavailable, using memory progress")
except ImportError as e:
    REDIS_AVAILABLE = False
    print(f"⚠️ Redis not available: {e}")

# Global variables
FIREBASE_AVAILABLE = False
OPENAI_AVAILABLE = False
VERTEX_AVAILABLE = False
db = None
openai_client = None

# Server startup time for uptime tracking
start_time = time.time()

# Keep-alive service variables
keep_alive_thread = None
keep_alive_running = False

# Flask Setup
app = Flask(__name__)
CORS(app, origins="*", methods=["GET", "POST", "DELETE", "OPTIONS", "PUT"], 
     allow_headers=["Content-Type", "Authorization"])

# Global progress tracking with thread safety
upload_progress = {}
progress_lock = Lock()

# Global embedding cache for duplicate content optimization
embedding_cache = {}
cache_lock = Lock()
CACHE_MAX_SIZE = 1000  # Limit cache size to prevent memory issues

# Create uploads folder with absolute path
UPLOAD_FOLDER = os.path.abspath("uploads")
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

# ================================
# DEPENDENCY CHECKING
# ================================

def check_file_processing_dependencies():
    """Check what file types are currently supported and provide installation guidance"""
    supported = {"always": ["txt", "json"], "conditional": []}
    missing = []
    
    print("🔍 Checking file processing capabilities...")
    
    try:
        import openpyxl
        supported["conditional"].append("xlsx")
        print("✅ openpyxl available - Excel .xlsx files supported")
    except ImportError:
        missing.append("openpyxl")
        print("❌ openpyxl missing - Excel .xlsx files NOT supported")
    
    try:
        import xlrd
        supported["conditional"].append("xls") 
        print("✅ xlrd available - Excel .xls files supported")
    except ImportError:
        missing.append("xlrd")
        print("❌ xlrd missing - Excel .xls files NOT supported")
    
    try:
        from PyPDF2 import PdfReader
        supported["conditional"].append("pdf")
        print("✅ PyPDF2 available - PDF files supported")
    except ImportError:
        missing.append("PyPDF2")
        print("❌ PyPDF2 missing - PDF files NOT supported")
    
    try:
        import docx
        supported["conditional"].append("docx")
        print("✅ python-docx available - Word files supported")
    except ImportError:
        missing.append("python-docx")
        print("❌ python-docx missing - Word files NOT supported")
    
    try:
        from pptx import Presentation
        supported["conditional"].append("pptx")
        print("✅ python-pptx available - PowerPoint files supported")
    except ImportError:
        missing.append("python-pptx")
        print("⚠️ python-pptx missing - PowerPoint files NOT supported")
    
    # Always supported with pandas
    supported["conditional"].extend(["csv"])
    print("✅ CSV files always supported (pandas)")
    
    if missing:
        print(f"\n⚠️ Missing dependencies for: {', '.join(missing)}")
        print(f"Install with: pip install {' '.join(missing)}")
        print("For all file support: pip install openpyxl xlrd PyPDF2 python-docx python-pptx")
    
    all_supported = supported["always"] + supported["conditional"]
    print(f"\n✅ Currently supported file types: {', '.join(sorted(all_supported))}")
    
    return supported, missing

# ================================
# RERANKING FUNCTIONALITY
# ================================

# Global variable to cache the cross-encoder model
_cross_encoder_model = None

def get_cross_encoder_model():
    """Lazy load the cross-encoder model for reranking"""
    global _cross_encoder_model
    if _cross_encoder_model is None:
        try:
            from sentence_transformers import CrossEncoder
            print("🔄 Loading cross-encoder model for reranking...")
            
            # Try to load the model with timeout and retry logic
            try:
                # Check for HF token to avoid rate limiting
                import os
                hf_token = os.getenv('HF_TOKEN') or os.getenv('HUGGINGFACE_TOKEN')
                if hf_token:
                    print("🔑 Using Hugging Face token to avoid rate limits")
                else:
                    print("⚠️ No HF_TOKEN found - may hit rate limits")
                
                # Try to load model with device handling
                try:
                    import torch
                    device = "cpu"  # Force CPU to avoid tensor issues
                    print(f"🔧 Loading cross-encoder on device: {device}")
                    _cross_encoder_model = CrossEncoder('cross-encoder/ms-marco-MiniLM-L-2-v2', device=device)
                except Exception as device_error:
                    # Fallback without explicit device
                    print(f"🔧 Device-specific loading failed: {device_error}")
                    print("🔧 Loading cross-encoder with auto device detection")
                    try:
                        _cross_encoder_model = CrossEncoder('cross-encoder/ms-marco-MiniLM-L-2-v2')
                    except Exception as fallback_error:
                        print(f"🔧 Auto device loading failed: {fallback_error}")
                        # Try with trust_remote_code for newer models
                        _cross_encoder_model = CrossEncoder('cross-encoder/ms-marco-MiniLM-L-2-v2', trust_remote_code=True)
                
                print("✅ Cross-encoder model loaded successfully")
            except Exception as download_error:
                error_str = str(download_error)
                if "429" in error_str or "rate limit" in error_str.lower():
                    print("⚠️ Hugging Face rate limit hit - reranking temporarily disabled")
                    print("💡 Fallback: Using hybrid search + similarity ranking instead")
                elif "connection" in error_str.lower() or "timeout" in error_str.lower():
                    print("⚠️ Network connection issue - reranking temporarily disabled")  
                    print("💡 Fallback: Using hybrid search + similarity ranking instead")
                else:
                    print(f"⚠️ Error loading cross-encoder model: {error_str}")
                return None
                
        except ImportError:
            print("⚠️ sentence-transformers not available - reranking disabled")
            return None
        except Exception as e:
            print(f"⚠️ Error initializing cross-encoder: {str(e)}")
            return None
    return _cross_encoder_model

def rerank_documents(query, documents, top_k=5):
    """
    Rerank documents using cross-encoder model
    
    Args:
        query (str): The search query
        documents (list): List of document dictionaries with 'content' and 'score' keys
        top_k (int): Number of top documents to return after reranking
    
    Returns:
        list: Reranked documents with updated scores
    """
    try:
        model = get_cross_encoder_model()
        if model is None:
            print("⚠️ Cross-encoder model not available - falling back to similarity ranking")
            return sorted(documents, key=lambda x: x.get("score", 0), reverse=True)[:top_k]
        
        if not documents:
            return []
        
        print(f"🔄 Reranking {len(documents)} documents with cross-encoder...")
        
        # Prepare query-document pairs for cross-encoder
        pairs = []
        for doc in documents:
            content = doc.get("content", "")
            # Truncate content to avoid token limits (cross-encoder typically handles ~512 tokens)
            truncated_content = content[:2000] if len(content) > 2000 else content
            pairs.append([query, truncated_content])
        
        # Get cross-encoder scores with error handling
        try:
            cross_encoder_scores = model.predict(pairs)
        except Exception as predict_error:
            error_str = str(predict_error)
            if "meta tensor" in error_str or "torch.nn.Module.to_empty" in error_str:
                print("⚠️ PyTorch tensor loading issue - falling back to similarity ranking")
            elif "cuda" in error_str.lower() or "device" in error_str.lower():
                print("⚠️ Device/CUDA issue - falling back to similarity ranking")
            else:
                print(f"⚠️ Cross-encoder prediction error: {error_str}")
            # Return similarity-ranked results as fallback
            return sorted(documents, key=lambda x: x.get("score", 0), reverse=True)[:top_k]
        
        # Update documents with new scores
        reranked_docs = []
        for i, doc in enumerate(documents):
            doc_copy = doc.copy()
            doc_copy["cross_encoder_score"] = float(cross_encoder_scores[i])
            doc_copy["original_score"] = doc.get("score", 0)
            reranked_docs.append(doc_copy)
        
        # Sort by cross-encoder score and return top_k
        reranked_docs.sort(key=lambda x: x["cross_encoder_score"], reverse=True)
        top_reranked = reranked_docs[:top_k]
        
        print(f"✅ Reranking complete - returned top {len(top_reranked)} documents")
        return top_reranked
        
    except Exception as e:
        print(f"❌ Error during reranking: {str(e)}")
        print("🔄 Falling back to similarity-based ranking...")
        return sorted(documents, key=lambda x: x.get("score", 0), reverse=True)[:top_k]

# ================================
# HYBRID SEARCH (DENSE + SPARSE)
# ================================

# Global variables for BM25 models and text preprocessing
_bm25_models = {}  # Cache BM25 models per organization/project
_text_preprocessor = None

def get_text_preprocessor():
    """Initialize and cache text preprocessor"""
    global _text_preprocessor
    if _text_preprocessor is None:
        try:
            import nltk
            from nltk.tokenize import word_tokenize
            from nltk.corpus import stopwords
            from nltk.stem import PorterStemmer
            import string
            import re
            
            # Download required NLTK data (only once) - handles Cloud Run deployment
            print("🔄 Checking and downloading NLTK data...")
            
            # Download punkt_tab (newer requirement)
            try:
                nltk.data.find('tokenizers/punkt_tab')
                print("✅ punkt_tab already available")
            except LookupError:
                try:
                    print("📥 Downloading NLTK punkt_tab tokenizer...")
                    nltk.download('punkt_tab', quiet=False)
                    print("✅ punkt_tab downloaded successfully")
                except Exception as e:
                    print(f"⚠️ Failed to download punkt_tab: {e}")
            
            # Download punkt (fallback)
            try:
                nltk.data.find('tokenizers/punkt')
                print("✅ punkt already available")
            except LookupError:
                try:
                    print("📥 Downloading NLTK punkt tokenizer...")
                    nltk.download('punkt', quiet=False)
                    print("✅ punkt downloaded successfully")
                except Exception as e:
                    print(f"⚠️ Failed to download punkt: {e}")
            
            # Download stopwords
            try:
                nltk.data.find('corpora/stopwords')
                print("✅ stopwords already available")
            except LookupError:
                try:
                    print("📥 Downloading NLTK stopwords...")
                    nltk.download('stopwords', quiet=False)
                    print("✅ stopwords downloaded successfully")
                except Exception as e:
                    print(f"⚠️ Failed to download stopwords: {e}")
            
            try:
                stemmer = PorterStemmer()
            except:
                stemmer = None
                
            try:
                stop_words = set(stopwords.words('english'))
            except:
                # Fallback to basic English stopwords
                stop_words = set(['the', 'a', 'an', 'and', 'or', 'but', 'in', 'on', 'at', 'to', 'for', 'of', 'with', 'by', 'is', 'are', 'was', 'were', 'be', 'been', 'being', 'have', 'has', 'had', 'do', 'does', 'did', 'will', 'would', 'should', 'could', 'can', 'may', 'might', 'must', 'shall', 'this', 'that', 'these', 'those', 'i', 'you', 'he', 'she', 'it', 'we', 'they', 'me', 'him', 'her', 'us', 'them'])
            
            def preprocess_text(text):
                if not text or not isinstance(text, str):
                    return []
                
                try:
                    # Convert to lowercase and remove extra whitespace
                    text = text.lower().strip()
                    
                    # Remove special characters but keep alphanumeric and spaces
                    text = re.sub(r'[^\w\s]', ' ', text)
                    
                    # Try NLTK tokenization first
                    try:
                        tokens = word_tokenize(text)
                    except:
                        # Fallback to simple split if NLTK fails
                        tokens = text.split()
                    
                    # Remove stopwords, punctuation, and stem
                    processed_tokens = []
                    for token in tokens:
                        if (token not in stop_words and 
                            token not in string.punctuation and 
                            len(token) > 2 and
                            token.isalpha()):
                            if stemmer:
                                try:
                                    processed_tokens.append(stemmer.stem(token))
                                except:
                                    # If stemming fails, use original token
                                    processed_tokens.append(token)
                            else:
                                # No stemmer available, use original token
                                processed_tokens.append(token)
                    
                    return processed_tokens
                    
                except Exception as e:
                    print(f"⚠️ Text preprocessing error: {e}")
                    # Ultimate fallback - simple split
                    return text.lower().split() if isinstance(text, str) else []
            
            _text_preprocessor = preprocess_text
            print("✅ Text preprocessor initialized")
            
        except ImportError as e:
            print(f"⚠️ NLTK not available: {e}")
            # Fallback simple preprocessor
            def simple_preprocess(text):
                if not text or not isinstance(text, str):
                    return []
                return text.lower().split()
            _text_preprocessor = simple_preprocess
            
    return _text_preprocessor

def create_bm25_index(documents, index_key):
    """
    Create BM25 index for a set of documents
    
    Args:
        documents (list): List of document dictionaries with 'content' field
        index_key (str): Unique key for caching this index
    
    Returns:
        BM25Okapi: Initialized BM25 model
    """
    global _bm25_models
    
    if index_key in _bm25_models:
        return _bm25_models[index_key]
    
    try:
        from rank_bm25 import BM25Okapi
        
        print(f"🔧 Creating BM25 index for {len(documents)} documents (key: {index_key})")
        
        preprocessor = get_text_preprocessor()
        
        # Preprocess documents
        tokenized_docs = []
        for doc in documents:
            content = doc.get('content', '')
            tokens = preprocessor(content)
            tokenized_docs.append(tokens)
        
        # Create BM25 index
        bm25 = BM25Okapi(tokenized_docs)
        _bm25_models[index_key] = bm25
        
        print(f"✅ BM25 index created for {index_key}")
        return bm25
        
    except ImportError:
        print("⚠️ rank-bm25 not available - BM25 search disabled")
        return None
    except Exception as e:
        print(f"❌ Error creating BM25 index: {str(e)}")
        return None

def bm25_search(query, documents, index_key, top_k=10):
    """
    Perform BM25 search on documents
    
    Args:
        query (str): Search query
        documents (list): List of document dictionaries
        index_key (str): Unique key for BM25 index
        top_k (int): Number of top results to return
    
    Returns:
        list: Documents with BM25 scores
    """
    try:
        bm25 = create_bm25_index(documents, index_key)
        if bm25 is None:
            return []
        
        preprocessor = get_text_preprocessor()
        query_tokens = preprocessor(query)
        
        if not query_tokens:
            return []
        
        # Get BM25 scores
        scores = bm25.get_scores(query_tokens)
        
        # Combine documents with scores
        scored_docs = []
        for i, doc in enumerate(documents):
            if i < len(scores):
                doc_copy = doc.copy()
                doc_copy['bm25_score'] = float(scores[i])
                scored_docs.append(doc_copy)
        
        # Sort by BM25 score and return top_k
        scored_docs.sort(key=lambda x: x['bm25_score'], reverse=True)
        return scored_docs[:top_k]
        
    except Exception as e:
        print(f"❌ Error in BM25 search: {str(e)}")
        return []

def hybrid_search(query, documents, index_key, top_k=10, dense_weight=0.7, sparse_weight=0.3):
    """
    Perform hybrid search combining dense (semantic) and sparse (BM25) results
    
    Args:
        query (str): Search query
        documents (list): List of document dictionaries with 'content' and 'score' (semantic)
        index_key (str): Unique key for BM25 indexing
        top_k (int): Number of top results to return
        dense_weight (float): Weight for dense/semantic scores (0.0-1.0)
        sparse_weight (float): Weight for sparse/BM25 scores (0.0-1.0)
    
    Returns:
        list: Hybrid ranked documents with combined scores
    """
    try:
        print(f"🔄 Performing hybrid search: dense({dense_weight}) + sparse({sparse_weight})")
        
        # Step 1: Get BM25 results
        bm25_results = bm25_search(query, documents, index_key, top_k=min(50, len(documents)))
        
        # Create lookup for BM25 scores
        bm25_lookup = {}
        for doc in bm25_results:
            # Use content hash as key for matching
            content_key = hash(doc.get('content', ''))
            bm25_lookup[content_key] = doc.get('bm25_score', 0.0)
        
        # Step 2: Normalize scores to 0-1 range
        if documents:
            # Normalize semantic scores
            semantic_scores = [doc.get('score', 0.0) for doc in documents]
            if semantic_scores:
                max_semantic = max(semantic_scores)
                min_semantic = min(semantic_scores)
                semantic_range = max_semantic - min_semantic if max_semantic > min_semantic else 1.0
        
        if bm25_results:
            # Normalize BM25 scores
            bm25_scores = [doc.get('bm25_score', 0.0) for doc in bm25_results]
            max_bm25 = max(bm25_scores) if bm25_scores else 1.0
            min_bm25 = min(bm25_scores) if bm25_scores else 0.0
            bm25_range = max_bm25 - min_bm25 if max_bm25 > min_bm25 else 1.0
        else:
            max_bm25 = min_bm25 = bm25_range = 1.0
        
        # Step 3: Combine scores
        hybrid_docs = []
        for doc in documents:
            content_key = hash(doc.get('content', ''))
            
            # Normalize semantic score
            semantic_score = doc.get('score', 0.0)
            normalized_semantic = (semantic_score - min_semantic) / semantic_range if 'semantic_range' in locals() else semantic_score
            
            # Get and normalize BM25 score
            bm25_score = bm25_lookup.get(content_key, 0.0)
            normalized_bm25 = (bm25_score - min_bm25) / bm25_range if bm25_range > 0 else 0.0
            
            # Calculate hybrid score
            hybrid_score = (dense_weight * normalized_semantic) + (sparse_weight * normalized_bm25)
            
            # Create new document with hybrid score
            hybrid_doc = doc.copy()
            hybrid_doc['hybrid_score'] = float(hybrid_score)
            hybrid_doc['semantic_score'] = float(semantic_score)
            hybrid_doc['bm25_score'] = float(bm25_score)
            hybrid_doc['normalized_semantic'] = float(normalized_semantic)
            hybrid_doc['normalized_bm25'] = float(normalized_bm25)
            
            hybrid_docs.append(hybrid_doc)
        
        # Step 4: Sort by hybrid score and return top_k
        hybrid_docs.sort(key=lambda x: x['hybrid_score'], reverse=True)
        top_hybrid = hybrid_docs[:top_k]
        
        print(f"✅ Hybrid search complete - returned {len(top_hybrid)} results")
        return top_hybrid
        
    except Exception as e:
        print(f"❌ Error in hybrid search: {str(e)}")
        print("🔄 Falling back to semantic search only...")
        return sorted(documents, key=lambda x: x.get('score', 0), reverse=True)[:top_k]

# ================================
# QUERY ENHANCEMENT & EXPANSION
# ================================

def expand_query(original_query, context_type="general", max_variations=3):
    """
    Generate related query variations using LLM to improve search recall
    
    Args:
        original_query (str): The original user query
        context_type (str): Type of context ("general", "rfp", "technical", "compliance")
        max_variations (int): Maximum number of query variations to generate
    
    Returns:
        list: List of expanded queries including the original
    """
    try:
        if not VERTEX_AVAILABLE:
            print("⚠️ Vertex AI not available - using original query only")
            return [original_query]
        
        from vertexai.generative_models import GenerativeModel, Part
        
        print(f"🔄 Expanding query: '{original_query}' (context: {context_type})")
        
        # Context-specific expansion prompts
        if context_type == "rfp":
            expansion_prompt = f"""
You are helping to improve search for RFP (Request for Proposal) documents. Generate {max_variations} alternative ways to search for the same information.

Original query: "{original_query}"

Generate {max_variations} related search queries that would find the same or similar information in RFP documents. Focus on:
- Different terminology (e.g., "requirements" vs "specifications" vs "criteria")
- Technical vs business language variations
- Specific vs general phrasing
- Common RFP section names and structures

Return only the queries, one per line, without numbering or explanations.
"""
        elif context_type == "technical":
            expansion_prompt = f"""
You are helping to improve search for technical documentation. Generate {max_variations} alternative ways to search for the same information.

Original query: "{original_query}"

Generate {max_variations} related search queries that would find the same or similar technical information. Focus on:
- Synonyms and technical terms
- Different technical approaches or implementations
- Related technologies or concepts
- Specific vs general technical language

Return only the queries, one per line, without numbering or explanations.
"""
        elif context_type == "compliance":
            expansion_prompt = f"""
You are helping to improve search for compliance and regulatory documents. Generate {max_variations} alternative ways to search for the same information.

Original query: "{original_query}"

Generate {max_variations} related search queries that would find the same or similar compliance information. Focus on:
- Regulatory terminology and standards
- Compliance frameworks and requirements
- Legal and policy language variations
- Audit and certification terms

Return only the queries, one per line, without numbering or explanations.
"""
        else:  # general
            expansion_prompt = f"""
You are helping to improve document search. Generate {max_variations} alternative ways to search for the same information.

Original query: "{original_query}"

Generate {max_variations} related search queries that would find the same or similar information. Focus on:
- Synonyms and alternative phrasing
- More specific or more general versions
- Different perspectives on the same topic
- Related concepts and terms

Return only the queries, one per line, without numbering or explanations.
"""
        
        model = GenerativeModel("gemini-2.0-flash")
        response = model.generate_content(
            [Part.from_text(expansion_prompt)],
            generation_config={
                "max_output_tokens": 500,
                "temperature": 0.7,  # Some creativity for variations
                "top_p": 0.9,
            }
        )
        
        expanded_text = response.text.strip()
        
        # Parse the response into separate queries
        expanded_queries = []
        for line in expanded_text.split('\n'):
            query = line.strip()
            # Remove numbering, bullets, or other formatting
            query = query.lstrip('123456789.- •*')
            query = query.strip()
            if query and len(query) > 3:  # Valid query
                expanded_queries.append(query)
        
        # Combine original with expansions (limit to max_variations + 1)
        all_queries = [original_query] + expanded_queries[:max_variations]
        
        print(f"✅ Generated {len(all_queries)} query variations")
        for i, q in enumerate(all_queries):
            print(f"  {i+1}. {q}")
        
        return all_queries
        
    except Exception as e:
        print(f"❌ Error expanding query: {str(e)}")
        return [original_query]

def multi_query_search(queries, documents, search_func, **search_kwargs):
    """
    Perform search with multiple query variations and aggregate results
    
    Args:
        queries (list): List of query variations
        documents (list): Documents to search through
        search_func (function): Search function to use (semantic, hybrid, etc.)
        **search_kwargs: Additional arguments for search function
    
    Returns:
        list: Aggregated and deduplicated search results
    """
    try:
        print(f"🔄 Performing multi-query search with {len(queries)} variations...")
        
        all_results = {}  # Use dict to deduplicate by content
        query_weights = {}  # Track which queries found each document
        
        for i, query in enumerate(queries):
            print(f"  🔍 Searching variation {i+1}: '{query[:50]}...'")
            
            # Perform search for this query variation
            if 'index_key' in search_kwargs:
                # Update index key to include query variation
                base_key = search_kwargs['index_key']
                search_kwargs['index_key'] = f"{base_key}_q{i}"
            
            results = search_func(query, documents, **search_kwargs)
            
            # Aggregate results with weighting
            weight = 1.0 - (i * 0.1)  # Original query gets highest weight
            weight = max(weight, 0.5)  # Minimum weight of 0.5
            
            for doc in results:
                content_key = hash(doc.get('content', ''))
                
                if content_key in all_results:
                    # Combine scores from multiple queries
                    existing_doc = all_results[content_key]
                    
                    # Weighted average of scores
                    if 'hybrid_score' in doc:
                        current_weight = query_weights[content_key]
                        new_weight = current_weight + weight
                        existing_score = existing_doc.get('hybrid_score', 0)
                        new_score = doc.get('hybrid_score', 0)
                        combined_score = ((existing_score * current_weight) + (new_score * weight)) / new_weight
                        existing_doc['hybrid_score'] = combined_score
                        query_weights[content_key] = new_weight
                    elif 'score' in doc:
                        current_weight = query_weights[content_key]
                        new_weight = current_weight + weight
                        existing_score = existing_doc.get('score', 0)
                        new_score = doc.get('score', 0)
                        combined_score = ((existing_score * current_weight) + (new_score * weight)) / new_weight
                        existing_doc['score'] = combined_score
                        query_weights[content_key] = new_weight
                    
                    # Track which queries found this document
                    if 'matched_queries' not in existing_doc:
                        existing_doc['matched_queries'] = []
                    existing_doc['matched_queries'].append(f"Q{i+1}: {query[:30]}...")
                    
                else:
                    # New document
                    doc_copy = doc.copy()
                    doc_copy['matched_queries'] = [f"Q{i+1}: {query[:30]}..."]
                    
                    # Apply initial weight
                    if 'hybrid_score' in doc_copy:
                        doc_copy['hybrid_score'] *= weight
                    elif 'score' in doc_copy:
                        doc_copy['score'] *= weight
                    
                    all_results[content_key] = doc_copy
                    query_weights[content_key] = weight
        
        # Convert back to list and sort
        aggregated_results = list(all_results.values())
        
        # Sort by the best available score
        if aggregated_results and 'hybrid_score' in aggregated_results[0]:
            aggregated_results.sort(key=lambda x: x.get('hybrid_score', 0), reverse=True)
        else:
            aggregated_results.sort(key=lambda x: x.get('score', 0), reverse=True)
        
        print(f"✅ Multi-query search complete: {len(aggregated_results)} unique results")
        
        return aggregated_results
        
    except Exception as e:
        print(f"❌ Error in multi-query search: {str(e)}")
        print("🔄 Falling back to single query search...")
        return search_func(queries[0] if queries else "", documents, **search_kwargs)

def enhanced_search(query, documents, index_key, context_type="general", enable_query_expansion=True, 
                   max_query_variations=3, search_method="hybrid", **search_kwargs):
    """
    Enhanced search with query expansion and multi-query aggregation
    
    Args:
        query (str): Original search query
        documents (list): Documents to search
        index_key (str): Index key for caching
        context_type (str): Type of context for query expansion
        enable_query_expansion (bool): Whether to enable query expansion
        max_query_variations (int): Maximum query variations to generate
        search_method (str): "semantic", "hybrid", or "bm25"
        **search_kwargs: Additional arguments for search function
    
    Returns:
        list: Enhanced search results
    """
    try:
        # Step 1: Query expansion (if enabled)
        if enable_query_expansion:
            expanded_queries = expand_query(query, context_type, max_query_variations)
        else:
            expanded_queries = [query]
        
        # Step 2: Choose search function
        if search_method == "hybrid":
            search_func = hybrid_search
        elif search_method == "bm25":
            search_func = bm25_search
        else:  # semantic (fallback)
            def semantic_search(q, docs, key, **kwargs):
                # Simple semantic search fallback
                return sorted(docs, key=lambda x: x.get('score', 0), reverse=True)[:kwargs.get('top_k', 10)]
            search_func = semantic_search
        
        # Step 3: Multi-query search
        search_kwargs['index_key'] = index_key
        results = multi_query_search(expanded_queries, documents, search_func, **search_kwargs)
        
        return results
        
    except Exception as e:
        print(f"❌ Error in enhanced search: {str(e)}")
        print("🔄 Falling back to basic search...")
        # Fallback to basic search
        if search_method == "hybrid":
            return hybrid_search(query, documents, index_key, **search_kwargs)
        else:
            return sorted(documents, key=lambda x: x.get('score', 0), reverse=True)[:search_kwargs.get('top_k', 10)]

# ================================
# ENHANCED FILE EXTRACTION
# ================================

def extract_pdf_text(file_path):
    """Extract text from PDF files"""
    try:
        from PyPDF2 import PdfReader
        reader = PdfReader(file_path)
        text = ""
        
        for i, page in enumerate(reader.pages):
            try:
                page_text = page.extract_text()
                if page_text:
                    text += f"\n--- Page {i+1} ---\n{page_text}\n"
            except Exception as page_error:
                print(f"⚠️ Error extracting page {i+1}: {page_error}")
                continue
        
        return text
    except ImportError:
        return "ERROR: PyPDF2 not installed. Install with: pip install PyPDF2"
    except Exception as e:
        return f"PDF extraction error: {str(e)}"

def extract_docx_text(file_path):
    """Extract text from Word documents"""
    try:
        import docx
        doc = docx.Document(file_path)
        text = ""
        
        # Extract paragraphs
        for para in doc.paragraphs:
            if para.text.strip():
                text += para.text + "\n"
        
        # Extract tables
        for i, table in enumerate(doc.tables):
            text += f"\n--- Table {i+1} ---\n"
            for row in table.rows:
                row_text = " | ".join([cell.text.strip() for cell in row.cells])
                if row_text.strip():
                    text += row_text + "\n"
        
        return text
    except ImportError:
        return "ERROR: python-docx not installed. Install with: pip install python-docx"
    except Exception as e:
        return f"DOCX extraction error: {str(e)}"

def extract_excel_text(file_path, file_type):
    """Extract text from Excel files with comprehensive error handling"""
    try:
        engines = {"xlsx": "openpyxl", "xls": "xlrd"}
        engine = engines.get(file_type, "openpyxl")
        
        try:
            # Read the Excel file
            excel_file = pd.ExcelFile(file_path, engine=engine)
            sheet_names = excel_file.sheet_names
            
            text = f"Excel File: {len(sheet_names)} sheet(s) - {', '.join(sheet_names)}\n\n"
            
            # Process all sheets (limit to reasonable number for performance)
            max_sheets = min(len(sheet_names), 20)  # Process up to 20 sheets
            
            if len(sheet_names) > 20:
                text += f"⚠️ Note: Processing first {max_sheets} sheets out of {len(sheet_names)} total sheets for performance.\n\n"
            for sheet_name in sheet_names[:max_sheets]:
                try:
                    # Read more rows but still limit for performance
                    df = pd.read_excel(file_path, sheet_name=sheet_name, engine=engine, nrows=2000)
                    
                    text += f"=== Sheet: {sheet_name} ===\n"
                    text += f"Dimensions: {len(df)} rows × {len(df.columns)} columns\n"
                    text += f"Columns: {', '.join(df.columns.astype(str))}\n\n"
                    
                    # Show more sample data for better context
                    sample_rows = min(50, len(df))  # Show up to 50 rows
                    text += f"Sample Data (first {sample_rows} rows):\n"
                    text += df.head(sample_rows).to_string(index=False) + "\n\n"
                    
                    # Add summary statistics for numeric columns
                    numeric_cols = df.select_dtypes(include=[np.number]).columns
                    if len(numeric_cols) > 0:
                        text += f"Numeric Summary:\n"
                        text += df[numeric_cols].describe().to_string() + "\n\n"
                    
                except Exception as sheet_error:
                    text += f"Error reading sheet '{sheet_name}': {str(sheet_error)}\n\n"
            
            return text
            
        except ImportError as ie:
            error_msg = str(ie)
            if "openpyxl" in error_msg:
                return "ERROR: openpyxl not installed. Install with: pip install openpyxl"
            elif "xlrd" in error_msg:
                return "ERROR: xlrd not installed. Install with: pip install xlrd"
            else:
                return f"ERROR: Missing Excel dependency: {error_msg}"
                
    except Exception as e:
        error_msg = str(e)
        if "openpyxl" in error_msg or "xlrd" in error_msg:
            return "ERROR: Missing Excel dependencies. Install with: pip install openpyxl xlrd"
        return f"Excel processing error: {error_msg}"
    

def download_single_file_as_zip(backend_api_url, file_id, filename, headers):
    """Download a single file by requesting it as a ZIP and extracting it"""
    try:
        # Use the existing ZIP endpoint but with a single file ID
        zip_download_url = f"{backend_api_url}/api/files/org/files/download"
        
        payload = {"fileIds": [file_id]}
        
        print(f"📦 Downloading {filename} as ZIP from: {zip_download_url}")
        print(f"📋 Payload: {payload}")
        
        zip_response = requests.post(
            zip_download_url,
            json=payload,
            headers=headers,
            timeout=120  # Increased timeout for ZIP downloads
        )
        
        if zip_response.status_code != 200:
            error_text = zip_response.text[:500] if zip_response.text else "No response text"
            print(f"❌ ZIP download failed with status {zip_response.status_code}")
            print(f"Response: {error_text}")
            raise Exception(f"Failed to download ZIP: {zip_response.status_code} - {error_text}")
        
        # Validate ZIP content
        content_type = zip_response.headers.get('Content-Type', '').lower()
        content_length = zip_response.headers.get('Content-Length', 'Unknown')
        
        print(f"📊 ZIP Download Info:")
        print(f"  Content-Type: {content_type}")
        print(f"  Content-Length: {content_length}")
        print(f"  Actual Size: {len(zip_response.content)} bytes")
        
        # Check if we got a ZIP file
        if 'application/zip' not in content_type and 'application/octet-stream' not in content_type:
            # Sometimes the content-type might not be set correctly, check by magic bytes
            if len(zip_response.content) < 4 or not zip_response.content.startswith(b'PK'):
                raise Exception(f"Expected ZIP file, got content-type: {content_type}, content preview: {zip_response.content[:100]}")
        
        # Validate ZIP file by magic bytes
        if len(zip_response.content) < 4:
            raise Exception("Response too short to be a valid ZIP file")
            
        if not zip_response.content.startswith(b'PK'):
            # Try to see if it's an error message
            try:
                error_preview = zip_response.content.decode('utf-8', errors='ignore')[:200]
                if any(keyword in error_preview.lower() for keyword in ['error', 'exception', 'not found', '<html>']):
                    raise Exception(f"Received error response instead of ZIP: {error_preview}")
            except:
                pass
            raise Exception("Response does not have ZIP file signature (PK)")
        
        print(f"✅ Valid ZIP file received ({len(zip_response.content)} bytes)")
        
        # Extract the file from the ZIP
        try:
            with zipfile.ZipFile(io.BytesIO(zip_response.content)) as zip_file:
                # List files in the ZIP
                file_list = zip_file.namelist()
                print(f"📁 ZIP contains {len(file_list)} file(s): {file_list}")
                
                if not file_list:
                    raise Exception("ZIP file is empty")
                
                # Find the target file (by name or take the first one)
                target_file = None
                
                # Strategy 1: Exact filename match
                for zip_filename in file_list:
                    if zip_filename == filename:
                        target_file = zip_filename
                        print(f"✅ Found exact filename match: {target_file}")
                        break
                
                # Strategy 2: Partial filename match (ends with)
                if not target_file:
                    for zip_filename in file_list:
                        if zip_filename.endswith(filename) or filename.endswith(zip_filename):
                            target_file = zip_filename
                            print(f"✅ Found partial filename match: {target_file}")
                            break
                
                # Strategy 3: Same file extension
                if not target_file:
                    file_ext = filename.split('.')[-1].lower() if '.' in filename else ''
                    for zip_filename in file_list:
                        if zip_filename.split('.')[-1].lower() == file_ext:
                            target_file = zip_filename
                            print(f"✅ Found file with same extension: {target_file}")
                            break
                
                # Strategy 4: Take the first file (fallback)
                if not target_file:
                    target_file = file_list[0]
                    print(f"⚠️ No exact match found, using first file: {target_file}")
                
                # Validate the target file in ZIP
                file_info = zip_file.getinfo(target_file)
                print(f"📄 Target file info: {target_file} (size: {file_info.file_size} bytes, compressed: {file_info.compress_size} bytes)")
                
                if file_info.file_size == 0:
                    raise Exception(f"Target file {target_file} is empty in ZIP")
                
                # Extract the file content
                file_content = zip_file.read(target_file)
                
                if len(file_content) == 0:
                    raise Exception(f"Extracted content from {target_file} is empty")
                
                print(f"✅ Successfully extracted {len(file_content)} bytes from ZIP file {target_file}")
                
                # Additional validation for known file types
                file_ext = target_file.split('.')[-1].lower() if '.' in target_file else ''
                
                if file_ext == 'xlsx' and not file_content.startswith(b'PK'):
                    raise Exception(f"Extracted XLSX file doesn't have valid format (should start with PK)")
                elif file_ext == 'pdf' and not file_content.startswith(b'%PDF'):
                    raise Exception(f"Extracted PDF file doesn't have valid format (should start with %PDF)")
                elif file_ext in ['txt', 'csv', 'json']:
                    # Try to decode text files to validate
                    try:
                        decoded = file_content.decode('utf-8', errors='ignore')
                        if len(decoded.strip()) == 0:
                            raise Exception(f"Extracted {file_ext} file appears to be empty or invalid")
                    except Exception as decode_error:
                        print(f"⚠️ Warning: Could not validate {file_ext} file: {decode_error}")
                
                return file_content
                
        except zipfile.BadZipFile as zip_error:
            print(f"❌ Invalid ZIP file: {zip_error}")
            # Try to see what we actually got
            content_preview = zip_response.content[:200]
            try:
                text_preview = content_preview.decode('utf-8', errors='ignore')
                print(f"Content preview: {text_preview}")
            except:
                print(f"Binary content preview: {content_preview}")
            raise Exception(f"Invalid ZIP file format: {zip_error}")
            
    except Exception as e:
        print(f"❌ Error downloading file as ZIP: {str(e)}")
        raise


def extract_csv_text(file_path):
    """Extract text from CSV files with smart delimiter detection"""
    try:
        # Try different separators and encodings
        encodings = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']
        separators = [',', ';', '\t', '|']
        
        df = None
        used_encoding = None
        used_separator = None
        
        for encoding in encodings:
            for sep in separators:
                try:
                    df = pd.read_csv(file_path, encoding=encoding, sep=sep, nrows=10000)  # Read more rows
                    if len(df.columns) > 1 and len(df) > 0:  # Valid data
                        used_encoding = encoding
                        used_separator = sep
                        break
                except:
                    continue
            if df is not None and len(df.columns) > 1:
                break
        
        if df is not None and len(df) > 0:
            text = f"CSV File Analysis:\n"
            text += f"Encoding: {used_encoding}, Separator: '{used_separator}'\n"
            text += f"Dimensions: {len(df)} rows × {len(df.columns)} columns\n"
            text += f"Columns: {', '.join(df.columns.astype(str))}\n\n"
            
            # Clean column names (remove extra whitespace)
            df.columns = df.columns.str.strip()
            
            # Show more sample data for better context
            sample_rows = min(100, len(df))  # Show up to 100 rows for CSV
            text += f"Sample Data (first {sample_rows} rows):\n"
            text += df.head(sample_rows).to_string(index=False) + "\n\n"
            
            # Add summary statistics
            numeric_cols = df.select_dtypes(include=[np.number]).columns
            if len(numeric_cols) > 0:
                text += "Numeric Summary:\n"
                text += df[numeric_cols].describe().to_string() + "\n\n"
            
            # Show unique values for categorical columns (first few)
            categorical_cols = df.select_dtypes(include=['object']).columns[:3]
            for col in categorical_cols:
                unique_vals = df[col].value_counts().head(5)
                if len(unique_vals) > 0:
                    text += f"Top values in '{col}':\n{unique_vals.to_string()}\n\n"
            
            return text
        else:
            # Fallback to plain text
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                content = f.read()
                return f"CSV file (read as text):\n{content[:2000]}..." if len(content) > 2000 else content
                
    except Exception as e:
        # Ultimate fallback
        try:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                content = f.read()
                return f"CSV parsing failed, read as text: {str(e)}\n\nContent:\n{content[:1000]}..."
        except:
            return f"CSV processing error: {str(e)}"

def extract_json_text(file_path):
    """Extract and format JSON data - supports both JSON and JSONL formats"""
    try:
        # First, try to parse as standard JSON
        with open(file_path, "r", encoding="utf-8") as f:
            try:
                data = json.load(f)
                
                # Pretty format JSON for AI understanding
                formatted = json.dumps(data, indent=2, ensure_ascii=False)
                
                # Add summary information
                text = f"JSON File Analysis:\n"
                text += f"Structure: {type(data).__name__}\n"
                
                if isinstance(data, dict):
                    text += f"Top-level keys: {', '.join(list(data.keys())[:10])}\n"
                elif isinstance(data, list):
                    text += f"Array length: {len(data)}\n"
                    if len(data) > 0:
                        text += f"First item type: {type(data[0]).__name__}\n"
                
                text += f"\nFormatted Content:\n{formatted[:5000]}"
                if len(formatted) > 5000:
                    text += "\n... (truncated)"
                
                return text
                
            except json.JSONDecodeError:
                # If standard JSON parsing fails, try JSONL (JSON Lines) format
                print(f"📄 Standard JSON parsing failed, trying JSONL format...")
                f.seek(0)  # Reset file pointer
                lines = f.readlines()
                
                json_objects = []
                parsed_count = 0
                max_objects = 1000  # Limit for performance
                
                for line_num, line in enumerate(lines[:max_objects], 1):
                    line = line.strip()
                    if line:  # Skip empty lines
                        try:
                            obj = json.loads(line)
                            json_objects.append(obj)
                            parsed_count += 1
                        except json.JSONDecodeError as line_error:
                            print(f"⚠️ Error parsing line {line_num}: {line_error}")
                            continue
                
                if parsed_count > 0:
                    # Successfully parsed as JSONL
                    text = f"JSONL (JSON Lines) File Analysis:\n"
                    text += f"Total lines processed: {len(lines)}\n"
                    text += f"Valid JSON objects: {parsed_count}\n"
                    
                    if parsed_count > max_objects:
                        text += f"⚠️ Note: Processed first {max_objects} objects out of {len(lines)} total lines\n"
                    
                    # Show structure of first object
                    if json_objects:
                        first_obj = json_objects[0]
                        text += f"Object structure: {type(first_obj).__name__}\n"
                        if isinstance(first_obj, dict):
                            text += f"Object keys: {', '.join(list(first_obj.keys())[:10])}\n"
                    
                    text += f"\nSample Objects (first {min(10, len(json_objects))}):\n"
                    for i, obj in enumerate(json_objects[:10]):
                        text += f"Object {i+1}: {json.dumps(obj, ensure_ascii=False)}\n"
                    
                    # Include more objects in a compact format
                    if len(json_objects) > 10:
                        text += f"\n... and {len(json_objects) - 10} more objects\n"
                        
                        # Add a summary of all objects
                        all_formatted = "\n".join([json.dumps(obj, ensure_ascii=False) for obj in json_objects])
                        if len(all_formatted) <= 10000:  # Include all if reasonable size
                            text += f"\nAll Objects:\n{all_formatted}"
                        else:
                            text += f"\nAdditional Objects (truncated):\n{all_formatted[:5000]}..."
                    
                    return text
                else:
                    # Neither JSON nor JSONL worked, fall back to text
                    raise json.JSONDecodeError("No valid JSON objects found", "", 0)
        
    except json.JSONDecodeError as e:
        # Final fallback: read as regular text
        try:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                content = f.read()
                return f"JSON parsing failed: {str(e)}\nFile content (as text):\n{content[:2000]}..."
        except:
            return f"JSON processing error: {str(e)}"
    except Exception as e:
        return f"JSON file error: {str(e)}"

def extract_pptx_text(file_path):
    """Extract text from PowerPoint files"""
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        text = f"PowerPoint Presentation: {len(prs.slides)} slides\n\n"
        
        for i, slide in enumerate(prs.slides):
            text += f"--- Slide {i+1} ---\n"
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text += shape.text + "\n"
            
            text += "\n"
        
        return text
    except ImportError:
        return "ERROR: python-pptx not installed. Install with: pip install python-pptx"
    except Exception as e:
        return f"PowerPoint extraction error: {str(e)}"

def extract_xml_text(file_path):
    """Extract text from XML files"""
    try:
        import xml.etree.ElementTree as ET
        
        tree = ET.parse(file_path)
        root = tree.getroot()
        
        def extract_text_from_xml(element, level=0):
            result = ""
            indent = "  " * level
            
            # Add element name
            result += f"{indent}<{element.tag}>\n"
            
            # Add text content
            if element.text and element.text.strip():
                result += f"{indent}  {element.text.strip()}\n"
            
            # Process children
            for child in element:
                result += extract_text_from_xml(child, level + 1)
            
            return result
        
        text = f"XML File: Root element '{root.tag}'\n\n"
        text += extract_text_from_xml(root)
        
        return text[:5000] + ("..." if len(text) > 5000 else "")
        
    except Exception as e:
        # Fallback to text reading
        try:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                content = f.read()
                return f"XML parsing failed: {str(e)}\nContent:\n{content[:1000]}..."
        except:
            return f"XML processing error: {str(e)}"

def extract_text_with_encoding_detection(file_path):
    """Extract text with automatic encoding detection"""
    try:
        # Try to detect encoding
        import chardet
        with open(file_path, "rb") as f:
            raw_data = f.read()
            result = chardet.detect(raw_data)
            encoding = result['encoding']
        
        with open(file_path, "r", encoding=encoding) as f:
            return f.read()
    except ImportError:
        # Fallback without chardet
        encodings = ['utf-8', 'utf-16', 'latin-1', 'cp1252', 'iso-8859-1']
        for encoding in encodings:
            try:
                with open(file_path, "r", encoding=encoding) as f:
                    return f.read()
            except UnicodeDecodeError:
                continue
        
        # Ultimate fallback
        with open(file_path, "rb") as f:
            return f.read().decode('utf-8', errors='ignore')
    except Exception as e:
        return f"Text extraction error: {str(e)}"

def is_ai_error_response(text):
    """Check if the AI response is an error message rather than extracted content"""
    if not text:
        return False
    
    text_lower = text.lower().strip()
    
    # Common AI error patterns
    error_patterns = [
        "i am unable to fulfill this request",
        "i cannot fulfill this request", 
        "i'm unable to",
        "i cannot",
        "i am sorry, but i cannot",
        "i apologize, but i cannot",
        "sorry, i cannot",
        "text extraction error",
        "codec can't decode",
        "utf-8 codec",
        "invalid start byte",
        "encoding error",
        "unable to extract",
        "cannot extract",
        "extraction failed",
        "need a document to extract",
        "[object object]",
        "provided \"[object object]\"",
        "not a document",
        "please provide the actual document"
    ]
    
    # Check if the response starts with common error phrases
    for pattern in error_patterns:
        if text_lower.startswith(pattern) or pattern in text_lower[:200]:
            return True
    
    # Check if it's suspiciously short for a document (likely an error)
    if len(text.strip()) < 50 and any(word in text_lower for word in ["error", "unable", "cannot", "failed"]):
        return True
        
    return False

def extract_text_with_ai_direct(file_path, filename):
    """Extract text content directly from any file using Vertex AI multimodal capabilities"""
    if not VERTEX_AVAILABLE:
        print("❌ Vertex AI not available, falling back to traditional extraction")
        return None
        
    try:
        # Verify file exists and is readable
        print(f"🔍 AI extraction - file_path: {file_path}")
        print(f"🔍 AI extraction - absolute path: {os.path.abspath(file_path)}")
        print(f"🔍 AI extraction - current working directory: {os.getcwd()}")
        
        if not os.path.exists(file_path):
            print(f"❌ File does not exist: {file_path}")
            return None
        
        file_size = os.path.getsize(file_path)
        if file_size == 0:
            print(f"❌ File is empty: {file_path}")
            return None
            
        print(f"🤖 AI-based text extraction for: {filename} ({file_size} bytes)")
        
        # Get file extension
        file_ext = filename.split(".")[-1].lower() if "." in filename else "unknown"
        
        # Customize prompt based on file type
        if file_ext in ["png", "jpg", "jpeg", "gif", "webp"]:
            prompt = f"""
Extract ALL visible text and content from this image. Include:
- All text visible in the image (OCR)
- Descriptions of charts, diagrams, or visual elements
- Any data, numbers, or structured information shown
- Context and meaning of visual elements

Please provide a comprehensive text representation of everything in this image that could be useful for document analysis and search.

Return the extracted content as plain text without any formatting or special characters.
"""
        else:
            prompt = f"""
Extract ALL text content from this document. Include:
- All paragraphs, headings, and body text
- Tables, lists, and structured data
- Any metadata or document information
- Numbers, dates, and specific details
- Do not add your own interpretation, just extract the actual content

Please provide a comprehensive text extraction of the entire document content.

Return the extracted content as plain text without any formatting or special characters.
"""

        from vertexai.generative_models import GenerativeModel, Part, SafetySetting
        model = GenerativeModel("gemini-2.0-flash")
        
        # Prepare the content
        parts = [Part.from_text(prompt)]
        
        # MIME type mapping
        mime_mapping = {
            # Documents
            "pdf": "application/pdf",
            "doc": "application/msword",
            "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "odt": "application/vnd.oasis.opendocument.text",
            "rtf": "application/rtf",
            
            # Spreadsheets
            "xls": "application/vnd.ms-excel",
            "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            "ods": "application/vnd.oasis.opendocument.spreadsheet",
            "numbers": "application/vnd.apple.numbers",
            
            # Presentations
            "ppt": "application/vnd.ms-powerpoint",
            "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "odp": "application/vnd.oasis.opendocument.presentation",
            "key": "application/vnd.apple.keynote",
            
            # Text formats
            "txt": "text/plain",
            "md": "text/markdown",
            "markdown": "text/markdown", 
            "csv": "text/csv",
            "tsv": "text/tab-separated-values",
            "xml": "application/xml",
            "html": "text/html",
            "htm": "text/html",
            "json": "application/json",
            "yaml": "application/x-yaml",
            "yml": "application/x-yaml",
            
            # Code files
            "py": "text/x-python",
            "js": "application/javascript",
            "ts": "application/typescript",
            "jsx": "text/jsx",
            "tsx": "text/tsx",
            "java": "text/x-java-source",
            "cpp": "text/x-c++src",
            "c": "text/x-csrc",
            "cs": "text/x-csharp",
            "php": "application/x-php",
            "rb": "text/x-ruby",
            "go": "text/x-go",
            "rs": "text/x-rust",
            "swift": "text/x-swift",
            "kt": "text/x-kotlin",
            "scala": "text/x-scala",
            "sql": "application/sql",
            "sh": "application/x-sh",
            "ps1": "text/x-powershell",
            
            # Images
            "png": "image/png",
            "jpg": "image/jpeg",
            "jpeg": "image/jpeg",
            "gif": "image/gif",
            "webp": "image/webp",
            "bmp": "image/bmp",
            "tiff": "image/tiff",
            "tif": "image/tiff",
            "svg": "image/svg+xml",
            "ico": "image/x-icon",
            "heic": "image/heic",
            "heif": "image/heif",
            
            # Audio
            "mp3": "audio/mpeg",
            "wav": "audio/wav",
            "m4a": "audio/m4a",
            "aac": "audio/aac",
            "ogg": "audio/ogg",
            "flac": "audio/flac",
            
            # Video  
            "mp4": "video/mp4",
            "mov": "video/quicktime",
            "avi": "video/x-msvideo",
            "wmv": "video/x-ms-wmv",
            "flv": "video/x-flv",
            "webm": "video/webm",
            "mkv": "video/x-matroska",
            "m4v": "video/x-m4v",
            
            # Archives
            "zip": "application/zip",
            "rar": "application/vnd.rar",
            "7z": "application/x-7z-compressed",
            "tar": "application/x-tar",
            "gz": "application/gzip",
            "bz2": "application/x-bzip2",
            "xz": "application/x-xz",
            
            # eBooks
            "epub": "application/epub+zip",
            "mobi": "application/x-mobipocket-ebook",
            "azw": "application/vnd.amazon.ebook",
            "azw3": "application/vnd.amazon.ebook",
            
            # CAD/Design
            "dwg": "application/acad",
            "dxf": "application/dxf",
            "ai": "application/postscript",
            "psd": "application/vnd.adobe.photoshop",
            "sketch": "application/sketch",
            
            # Other common formats
            "log": "text/plain",
            "cfg": "text/plain",
            "conf": "text/plain",
            "ini": "text/plain",
            "properties": "text/plain",
            "env": "text/plain",
            "lock": "text/plain",
            "gitignore": "text/plain",
            "dockerfile": "text/plain",
        }
        
        # Files that Gemini can handle directly
        gemini_supported_files = [
            # Documents
            "pdf", "txt", "csv", "md", "markdown", "html", "htm", "xml", "yaml", "yml",
            
            # Images  
            "png", "jpg", "jpeg", "gif", "webp", "bmp", "tiff", "tif", "svg", "heic", "heif",
            
            # Code files (text-based)
            "py", "js", "ts", "jsx", "tsx", "java", "cpp", "c", "cs", "php", "rb", "go", "rs", 
            "swift", "kt", "scala", "sql", "sh", "ps1",
            
            # Configuration files
            "log", "cfg", "conf", "ini", "properties", "env", "lock", "gitignore", "dockerfile",
            
            # Audio (Gemini 2.0 supports audio)
            "mp3", "wav", "m4a", "aac", "ogg", "flac",
            
            # Video (Gemini 2.0 supports video)
            "mp4", "mov", "avi", "wmv", "webm", "mkv", "m4v"
        ]
        
        if file_ext in gemini_supported_files:
            try:
                # Read file as bytes for direct upload
                with open(file_path, "rb") as f:
                    file_data = f.read()
                
                mime_type = mime_mapping.get(file_ext, "application/octet-stream")
                
                # Check file size (Gemini has limits)
                max_size = 20 * 1024 * 1024  # 20MB limit for safety
                if len(file_data) > max_size:
                    print(f"⚠️ File too large for AI extraction ({len(file_data)} bytes), falling back")
                    return None
                
                # Add file part directly
                file_part = Part.from_data(data=file_data, mime_type=mime_type)
                parts.append(file_part)
                
                print(f"✅ Using AI multimodal extraction for {filename} ({mime_type}, {len(file_data)} bytes)")
                
            except Exception as e:
                print(f"⚠️ Direct AI upload failed for {filename}: {e}")
                return None
        else:
            # For unsupported file types, try to extract text using appropriate method
            try:
                print(f"📄 Attempting text extraction for unsupported type: {file_ext}")
                
                # Use proper extraction method based on file type
                if file_ext == "docx":
                    content = extract_docx_text(file_path)
                elif file_ext in ["xlsx", "xls"]:
                    content = extract_excel_text(file_path, file_ext)
                elif file_ext == "csv":
                    content = extract_csv_text(file_path)
                else:
                    # For truly unsupported types, try encoding detection
                    content = extract_text_with_encoding_detection(file_path)
                
                # Check if extraction returned an error
                if content and content.startswith("ERROR:"):
                    print(f"❌ Extraction failed: {content}")
                    return None
                
                if content and len(content.strip()) > 0:
                    # Limit content length for AI processing
                    if len(content) > 50000:  # 50k char limit
                        content = content[:50000] + "...[truncated]"
                    parts.append(Part.from_text(f"\nDocument Content to extract from:\n{content}"))
                else:
                    print(f"❌ No text content found in {filename}")
                    return None
            except Exception as e:
                print(f"❌ Failed to extract text: {e}")
                return None
        
        # Generate response with AI
        print(f"🤖 Sending to AI for text extraction...")
        
        try:
            response = model.generate_content(
                parts,
                generation_config={
                    "max_output_tokens": 8192,
                    "temperature": 0.1,  # Low temperature for accurate extraction
                    "top_p": 0.95,
                },
                safety_settings=[
                    SafetySetting(
                        category=SafetySetting.HarmCategory.HARM_CATEGORY_HATE_SPEECH,
                        threshold=SafetySetting.HarmBlockThreshold.BLOCK_NONE,
                    ),
                    SafetySetting(
                        category=SafetySetting.HarmCategory.HARM_CATEGORY_DANGEROUS_CONTENT,
                        threshold=SafetySetting.HarmBlockThreshold.BLOCK_NONE,
                    ),
                    SafetySetting(
                        category=SafetySetting.HarmCategory.HARM_CATEGORY_SEXUALLY_EXPLICIT,
                        threshold=SafetySetting.HarmBlockThreshold.BLOCK_NONE,
                    ),
                    SafetySetting(
                        category=SafetySetting.HarmCategory.HARM_CATEGORY_HARASSMENT,
                        threshold=SafetySetting.HarmBlockThreshold.BLOCK_NONE,
                    ),
                ]
            )
            
            if response and response.text:
                extracted_text = response.text.strip()
                
                # Check if the AI response is an error message
                if is_ai_error_response(extracted_text):
                    print(f"❌ AI returned error response for {filename}: {extracted_text[:100]}...")
                    return None
                
                print(f"✅ AI extracted {len(extracted_text)} characters from {filename}")
                return extracted_text
            else:
                print(f"❌ Empty response from AI for {filename}")
                return None
                
        except Exception as e:
            print(f"❌ AI generation failed for {filename}: {str(e)}")
            return None
            
    except Exception as e:
        print(f"❌ AI text extraction error for {filename}: {str(e)}")
        return None

def extract_text_from_file(file_path, file_type):
    """Enhanced text extraction supporting many file types - MAIN FUNCTION"""
    try:
        file_size = os.path.getsize(file_path)
        print(f"📄 Extracting from {file_path} (type: {file_type}, size: {file_size} bytes)")
        
        if file_size == 0:
            return "ERROR: File is empty"
        
        # Route to appropriate extraction function
        if file_type == "pdf":
            text = extract_pdf_text(file_path)
        elif file_type == "docx":
            text = extract_docx_text(file_path)
        elif file_type in ["xlsx", "xls"]:
            text = extract_excel_text(file_path, file_type)
        elif file_type == "csv":
            text = extract_csv_text(file_path)
        elif file_type == "json":
            text = extract_json_text(file_path)
        elif file_type == "pptx":
            text = extract_pptx_text(file_path)
        elif file_type == "xml":
            text = extract_xml_text(file_path)
        elif file_type == "txt":
            text = extract_text_with_encoding_detection(file_path)
        elif file_type in ["rtf"]:
            # Basic RTF handling
            content = extract_text_with_encoding_detection(file_path)
            # Strip RTF formatting codes
            text = re.sub(r'\{.*?\}', '', content)
            text = re.sub(r'\\[a-z]+\d*', '', text)
            text = ' '.join(text.split())
        else:
            # Unknown file type - try as text
            text = extract_text_with_encoding_detection(file_path)
            print(f"⚠️ Unknown file type {file_type}, treated as text")
        
        if not text or text.strip() == "":
            print(f"⚠️ No text extracted from {file_path}")
            return ""
        
        # Check for error messages in extracted text
        if text.startswith("ERROR:"):
            print(f"❌ {text}")
            return text
        
        print(f"✅ Successfully extracted {len(text)} characters from {file_path}")
        return text
        
    except Exception as e:
        error_msg = f"❌ Error extracting text from {file_path}: {str(e)}"
        print(error_msg)
        traceback.print_exc()
        return f"ERROR: {error_msg}"

# ================================
# FIREBASE AND AI INITIALIZATION
# ================================

# Initialize Firebase/Firestore
try:
    from firebase_admin import credentials, firestore, initialize_app
    from google.cloud import aiplatform
    from google.oauth2 import service_account
    
    cred_path = os.environ.get("GOOGLE_APPLICATION_CREDENTIALS", "fire.json")
    
    # For Firebase/Firestore, always use fire.json with buildnblog project
    if os.path.exists("fire.json"):
        # Use fire.json for Firebase/Firestore (buildnblog-450618 project)
        firebase_cred = credentials.Certificate("fire.json")
        print("✅ Using fire.json for Firebase/Firestore")
    elif cred_path.startswith("{"):
        # Fallback: parse JSON from environment but this might not have Firestore access
        import json
        cred_dict = json.loads(cred_path)
        firebase_cred = credentials.Certificate(cred_dict)
        print("✅ Using environment credentials for Firebase (may not have Firestore access)")
    else:
        # Fallback to file path
        firebase_cred = credentials.Certificate(cred_path)
        print(f"✅ Using {cred_path} for Firebase")
    
    initialize_app(firebase_cred)
    db = firestore.client()
    FIREBASE_AVAILABLE = True
    print("✅ Successfully initialized Firebase")
        
    try:
        from vertexai import init as vertex_init
        from vertexai.generative_models import GenerativeModel, Part, SafetySetting
        
        # Allow separate project IDs for different services
        project_id = os.environ.get("VERTEX_AI_PROJECT", os.environ.get("GOOGLE_CLOUD_PROJECT", "buildnblog-450618"))
        
        # For Vertex AI, always use fire.json file if it exists, otherwise use env credentials
        if os.path.exists("fire.json"):
            # Use fire.json for Vertex AI (buildnblog-450618 project)
            vertex_credentials = service_account.Credentials.from_service_account_file("fire.json")
            print("✅ Using fire.json for Vertex AI")
        elif cred_path.startswith("{"):
            # Fallback to environment JSON
            vertex_credentials = service_account.Credentials.from_service_account_info(cred_dict)
            print("✅ Using environment credentials for Vertex AI")
        else:
            # Fallback to file path
            vertex_credentials = service_account.Credentials.from_service_account_file(cred_path)
            print(f"✅ Using {cred_path} for Vertex AI")
            
        aiplatform.init(project=project_id, location="us-central1", credentials=vertex_credentials)
        vertex_init(project=project_id, location="us-central1")
        VERTEX_AVAILABLE = True
        print("✅ Successfully initialized Vertex AI")
    except Exception as e:
        print(f"⚠️ Vertex AI initialization failed: {str(e)}")

except Exception as e:
    print(f"⚠️ Firebase initialization failed: {str(e)}")

# Initialize OpenAI client
try:
    openai_api_key = os.environ.get("OPENAI_API_KEY")
    if not openai_api_key:
        print("⚠️ OPENAI_API_KEY environment variable not found. OpenAI embeddings will be disabled.")
    else:
        openai_client = OpenAI(api_key=openai_api_key)
        OPENAI_AVAILABLE = True
        print("✅ Successfully initialized OpenAI client")
except Exception as e:
    print(f"⚠️ OpenAI initialization failed: {str(e)}")

def get_openai_client():
    """Get the OpenAI client, initializing if needed"""
    global openai_client, OPENAI_AVAILABLE
    
    if not OPENAI_AVAILABLE:
        raise ValueError("OpenAI functionality is not available")
        
    if openai_client is None:
        try:
            openai_api_key = os.environ.get("OPENAI_API_KEY")
            if not openai_api_key:
                raise ValueError("OPENAI_API_KEY environment variable not found")
            
            openai_client = OpenAI(api_key=openai_api_key)
            print("✅ OpenAI client initialized successfully")
        except Exception as e:
            print(f"❌ Error initializing OpenAI client: {str(e)}")
            OPENAI_AVAILABLE = False
            raise
            
    return openai_client

# ================================
# UTILITY FUNCTIONS
# ================================

def update_upload_progress(upload_id, status, progress, stage, filename=""):
    """Enhanced upload progress tracking with Redis support"""
    if not upload_id:
        return
    
    # Try Redis first if available
    if REDIS_AVAILABLE:
        try:
            redis_manager.set_progress(upload_id, status, progress, stage, filename)
            print(f"📊 Progress updated in Redis: {upload_id} → {status} ({progress}%)")
            return
        except Exception as e:
            print(f"⚠️ Redis progress failed, using memory: {e}")
    
    # Fallback to memory-based progress
    progress_data = {
        "upload_id": upload_id,
        "status": status,
        "progress": progress,
        "stage": stage,
        "filename": filename,
        "timestamp": time.time()
    }
    
    with progress_lock:
        upload_progress[upload_id] = progress_data
        print(f"📊 Progress updated in memory: {upload_id} → {status} ({progress}%)")
    
    # Clean up completed uploads after a delay
    if status in ['Completed', 'error']:
        def cleanup():
            time.sleep(300)  # Wait 5 minutes before cleanup
            with progress_lock:
                if upload_id in upload_progress:
                    del upload_progress[upload_id]
                    print(f"Cleaned up completed upload: {upload_id}")
        
        cleanup_thread = Thread(target=cleanup)
        cleanup_thread.daemon = True
        cleanup_thread.start()

def update_backend_embedding_status(file_id, org_id, is_from_embedding):
    """Update the isFromEmbedding status in the backend database for both files and links"""
    try:
        BACKEND_API_URL = os.environ.get("BACKEND_API_URL", "http://localhost:5000")
        
        # Method 1: Try files endpoint first (for file uploads)
        files_update_url = f"{BACKEND_API_URL}/api/files/updateFileSource"
        files_payload = {
            "fileId": file_id,
            "isFromEmbedding": is_from_embedding
        }
        
        try:
            files_response = requests.post(
                files_update_url,
                json=files_payload,
                headers={"Content-Type": "application/json"},
                timeout=10
            )
            
            if files_response.status_code == 200:
                print(f"✅ Backend status updated for fileId: {file_id} -> {is_from_embedding}")
                return True
        except Exception as files_error:
            print(f"⚠️ Files endpoint failed: {files_error}")
        
        # Method 2: Try links endpoint (for website links)
        links_update_url = f"{BACKEND_API_URL}/api/links/{file_id}/embedding"
        links_payload = {
            "isFromEmbedding": is_from_embedding
        }
        
        try:
            links_response = requests.put(
                links_update_url,
                json=links_payload,
                headers={"Content-Type": "application/json"},
                timeout=10
            )
            
            if links_response.status_code == 200:
                print(f"✅ Backend status updated for linkId: {file_id} -> {is_from_embedding}")
                return True
            else:
                print(f"⚠️ Links endpoint failed: {links_response.status_code}")
                print(f"Response: {links_response.text}")
                return False
                
        except Exception as links_error:
            print(f"⚠️ Links endpoint error: {links_error}")
            return False
            
    except Exception as e:
        print(f"⚠️ Error updating backend status: {str(e)}")
        return False

@retry(
    retry=retry_if_exception_type((Exception)),
    stop=stop_after_attempt(3),
    wait=wait_exponential(multiplier=1, min=2, max=10),
    reraise=True
)
def embed_chunks(chunks, upload_id=None, org_id=None, filename=None):
    """Embed chunks with OpenAI API with retry logic - OPTIMIZED with caching"""
    try:
        client = get_openai_client()
        
        all_embeddings = []
        total = len(chunks)
        cache_hits = 0
        
        # OPTIMIZATION: Check cache for duplicate content
        print(f"🚀 Processing {total} chunks with caching optimization")
        
        # First pass: separate cached vs new chunks
        chunks_to_embed = []
        chunk_indices = []
        
        with cache_lock:
            for i, chunk in enumerate(chunks):
                chunk_hash = hash(chunk.strip())
                if chunk_hash in embedding_cache:
                    all_embeddings.append(embedding_cache[chunk_hash])
                    cache_hits += 1
                else:
                    chunks_to_embed.append(chunk)
                    chunk_indices.append(i)
        
        if cache_hits > 0:
            print(f"🎯 Cache hits: {cache_hits}/{total} chunks ({cache_hits/total*100:.1f}%)")
        
        # Only embed chunks not in cache
        if chunks_to_embed:
            # OPTIMIZATION: Increased batch size from 20 to 100 for 5x speedup
            batch_size = 100
            new_embeddings = []
            
            print(f"🚀 Embedding {len(chunks_to_embed)} new chunks in batches of {batch_size}")
            
            for i in range(0, len(chunks_to_embed), batch_size):
                batch = chunks_to_embed[i:i + batch_size]
                batch_end = min(i + batch_size, len(chunks_to_embed))
                
                try:
                    response = client.embeddings.create(
                        model="text-embedding-3-small",
                        input=batch
                    )
                    
                    for j, item in enumerate(response.data):
                        embedding = item.embedding
                        new_embeddings.append(embedding)
                        
                        # Cache the embedding
                        chunk_hash = hash(batch[j].strip())
                        with cache_lock:
                            # Manage cache size
                            if len(embedding_cache) >= CACHE_MAX_SIZE:
                                # Remove oldest 10% of cache
                                keys_to_remove = list(embedding_cache.keys())[:CACHE_MAX_SIZE//10]
                                for key in keys_to_remove:
                                    del embedding_cache[key]
                            
                            embedding_cache[chunk_hash] = embedding
                    
                    if upload_id:
                        progress = min(75 + ((batch_end) / len(chunks_to_embed)) * 20, 95)
                        update_upload_progress(upload_id, "Processing", progress, 
                                              f"Generating embeddings ({batch_end}/{len(chunks_to_embed)} new)")
                    
                    print(f"✓ Processed batch {i//batch_size + 1}: {batch_end}/{len(chunks_to_embed)} new chunks")
                    
                    # OPTIMIZATION: Removed fixed 0.5s delay - only sleep on rate limit errors
                    
                except Exception as batch_error:
                    # If we hit rate limits, add a small delay and retry
                    if "rate" in str(batch_error).lower():
                        print(f"⚠️ Rate limit hit, waiting 2s before retry...")
                        time.sleep(2)
                        # Retry the batch
                        response = client.embeddings.create(
                            model="text-embedding-3-small",
                            input=batch
                        )
                        for j, item in enumerate(response.data):
                            embedding = item.embedding
                            new_embeddings.append(embedding)
                    else:
                        raise batch_error
            
            # Merge cached and new embeddings in correct order
            result_embeddings = [None] * total
            
            # Place cached embeddings
            cached_idx = 0
            new_idx = 0
            for i in range(total):
                if i in chunk_indices:
                    result_embeddings[i] = new_embeddings[new_idx]
                    new_idx += 1
                else:
                    result_embeddings[i] = all_embeddings[cached_idx]
                    cached_idx += 1
            
            all_embeddings = result_embeddings
        
        cache_efficiency = f" ({cache_hits}/{total} cached)" if cache_hits > 0 else ""
        print(f"✅ Generated {len(all_embeddings)} embeddings{cache_efficiency} (5x faster with optimizations)")
        return all_embeddings
    except Exception as e:
        print(f"❌ Error generating embeddings: {str(e)}")
        raise

@retry(
    retry=retry_if_exception_type((Exception)),
    stop=stop_after_attempt(3),
    wait=wait_exponential(multiplier=1, min=2, max=10),
    reraise=True
)
def embed_query(query):
    """Embed a single query with OpenAI API with retry logic"""
    try:
        client = get_openai_client()
        
        response = client.embeddings.create(
            model="text-embedding-3-small",
            input=[query]
        )
        
        return response.data[0].embedding
    except Exception as e:
        print(f"❌ Error embedding query: {str(e)}")
        raise

def parse_and_chunk(file_path, file_ext, chunk_size=50, max_chunks=1000):
    """Parse file content into chunks using AI-first extraction"""
    try:
        # Extract filename from path for AI processing
        filename = os.path.basename(file_path)
        
        # Try AI-based extraction first
        print(f"🔍 About to call AI extraction: {file_path} exists={os.path.exists(file_path)}")
        text = extract_text_with_ai_direct(file_path, filename)
        print(f"🔍 After AI extraction call: {file_path} exists={os.path.exists(file_path)}")
        
        # If AI extraction fails, returns None, or returns an error message, fall back to traditional extraction
        if not text or text.strip() == "" or is_ai_error_response(text):
            if text and is_ai_error_response(text):
                print(f"🔄 AI extraction returned error response, falling back to traditional extraction for {filename}")
                print(f"📋 AI error response: {text[:100]}...")
            else:
                print(f"🔄 AI extraction failed, falling back to traditional extraction for {filename}")
            print(f"🔍 Checking file existence: {file_path} -> {os.path.exists(file_path)}")
            text = extract_text_from_file(file_path, file_ext)
        
        if not text or text.strip() == "":
            print(f"❌ No text extracted from {file_path} using any method")
            return []
        
        # Check for extraction errors from traditional method or remaining AI errors
        if text.startswith("ERROR:") or is_ai_error_response(text):
            if is_ai_error_response(text):
                print(f"❌ AI error response detected in final text: {text[:100]}...")
            else:
                print(f"❌ Extraction error: {text}")
            return []
            
        words = text.split()
        
        if len(words) > max_chunks * chunk_size:
            words = words[:max_chunks * chunk_size]
            print(f"⚠️ File truncated to {max_chunks} chunks to avoid memory issues")
        
        chunks = [" ".join(words[i: i + chunk_size]) for i in range(0, len(words), chunk_size)]
        print(f"✅ Successfully chunked file into {len(chunks)} chunks")
        return chunks
    except Exception as e:
        print(f"❌ Error in parse_and_chunk: {str(e)}")
        return []

def delete_collection(collection_ref, batch_size):
    """Enhanced helper function to delete all documents in a collection with counting"""
    if not FIREBASE_AVAILABLE:
        return 0
        
    deleted_count = 0
    
    try:
        docs = collection_ref.limit(batch_size).stream()
        batch_deleted = 0
        
        for doc in docs:
            doc.reference.delete()
            batch_deleted += 1
            deleted_count += 1
        
        if batch_deleted >= batch_size:
            deleted_count += delete_collection(collection_ref, batch_size)
            
    except Exception as e:
        print(f"❌ Error deleting collection batch: {str(e)}")
    
    return deleted_count

def get_knowledge_base_context(query, org_id, project_id, knowledge_base_option="global", rerank=False, enable_hybrid_search=True, dense_weight=0.7, sparse_weight=0.3, enable_query_expansion=True, max_query_variations=2, context_type="general"):
    """Get context from knowledge base based on option (global or specific)"""
    try:
        # Get query embedding
        query_embedding = np.array(embed_query(query))
        retrieved_docs = []
        
        if knowledge_base_option == "global":
            # Search organization-level documents (global knowledge base)
            org_files_ref = db.collection("document_embeddings").document(f"org-{org_id}").collection("files")
            org_files = org_files_ref.stream()
            
            # Process each organization file
            for file_doc in org_files:
                file_data = file_doc.to_dict()
                
                # Get chunks for this file
                chunks_ref = file_doc.reference.collection("chunks")
                chunks = chunks_ref.stream()
                
                # Process each chunk
                for chunk_doc in chunks:
                    chunk_data = chunk_doc.to_dict()
                    
                    # Convert to numpy array
                    chunk_embedding = np.array(chunk_data["embedding"])
                    
                    # Calculate cosine similarity
                    score = np.dot(query_embedding, chunk_embedding) / (
                        np.linalg.norm(query_embedding) * np.linalg.norm(chunk_embedding)
                    )
                    
                    if score >= 0.2:  # Similarity threshold
                        chunk_content = chunk_data.get("content") or chunk_data.get("text", "")
                        retrieved_docs.append({
                            "content": chunk_content,
                            "score": float(score)
                        })
        else:
            # Use specific project support documents
            project_files_ref = (db.collection("org_project_support_embeddings")
                               .document(f"org-{org_id}")
                               .collection("projects")
                               .document(f"project-{project_id}")
                               .collection("files"))
            
            files = project_files_ref.stream()
            
            # Process each file
            for file_doc in files:
                file_data = file_doc.to_dict()
                
                # Get chunks for this file
                chunks_ref = file_doc.reference.collection("chunks")
                chunks = chunks_ref.stream()
                
                # Process each chunk
                for chunk_doc in chunks:
                    chunk_data = chunk_doc.to_dict()
                    
                    # Convert to numpy array
                    chunk_embedding = np.array(chunk_data["embedding"])
                    
                    # Calculate cosine similarity
                    score = np.dot(query_embedding, chunk_embedding) / (
                        np.linalg.norm(query_embedding) * np.linalg.norm(chunk_embedding)
                    )
                    
                    if score >= 0.2:  # Similarity threshold
                        chunk_content = chunk_data.get("content") or chunk_data.get("text", "")
                        retrieved_docs.append({
                            "content": chunk_content,
                            "score": float(score)
                        })
        
        # Apply enhanced search with query expansion, hybrid search, and/or reranking
        if retrieved_docs:
            index_key = f"context_{knowledge_base_option}_{org_id}_{project_id}"
            
            if enable_query_expansion or enable_hybrid_search:
                print(f"🚀 Applying enhanced search to {len(retrieved_docs)} documents...")
                
                enhanced_candidates = enhanced_search(
                    query=query,
                    documents=retrieved_docs,
                    index_key=index_key,
                    context_type=context_type,
                    enable_query_expansion=enable_query_expansion,
                    max_query_variations=max_query_variations,
                    search_method="hybrid" if enable_hybrid_search else "semantic",
                    top_k=10,
                    dense_weight=dense_weight,
                    sparse_weight=sparse_weight
                )
                
                if rerank and enhanced_candidates:
                    print(f"🔄 Applying reranking to {len(enhanced_candidates)} enhanced candidates...")
                    top_chunks = rerank_documents(query, enhanced_candidates, top_k=3)
                else:
                    top_chunks = enhanced_candidates[:3]
                    
            elif rerank:
                print(f"🔄 Applying reranking to {len(retrieved_docs)} documents...")
                # Get more candidates for reranking (top 10-15), then rerank to get top 3
                similarity_candidates = sorted(retrieved_docs, key=lambda x: x["score"], reverse=True)[:15]
                top_chunks = rerank_documents(query, similarity_candidates, top_k=3)
            else:
                # Use traditional similarity ranking
                top_chunks = sorted(retrieved_docs, key=lambda x: x["score"], reverse=True)[:3]
        else:
            top_chunks = []
        
        context_text = "\n\n".join([doc["content"] for doc in top_chunks])
        
        return context_text if context_text.strip() else "No specific context available."
        
    except Exception as e:
        print(f"❌ Error getting knowledge base context: {str(e)}")
        return "No specific context available."

def generate_answer_with_gcp(query, context_chunks, conversation_history=""):
    """Generate answer using Google's Vertex AI"""
    if not VERTEX_AVAILABLE:
        return "Sorry, the AI generation service is currently unavailable."
        
    try:
        context_text = "\n\n".join(context_chunks)
        
        prompt = f"""
You are an intelligent assistant tasked with generating an accurate and comprehensive answer using only the information provided below.
Your response must rely solely on the conversation history and the retrieved context, without adding external knowledge or assumptions.

Conversation History:
{conversation_history}

Retrieved Context:
{context_text}

Question:
{query}
"""

        model = GenerativeModel("gemini-2.0-flash")

        responses = model.generate_content(
            [Part.from_text(prompt)],
            generation_config={
                "max_output_tokens": 8192,
                "temperature": 1,
                "top_p": 0.95,
            },
            safety_settings=[
                SafetySetting(
                    category=SafetySetting.HarmCategory.HARM_CATEGORY_HATE_SPEECH,
                    threshold=SafetySetting.HarmBlockThreshold.OFF
                ),
                SafetySetting(
                    category=SafetySetting.HarmCategory.HARM_CATEGORY_DANGEROUS_CONTENT,
                    threshold=SafetySetting.HarmBlockThreshold.OFF
                ),
                SafetySetting(
                    category=SafetySetting.HarmCategory.HARM_CATEGORY_SEXUALLY_EXPLICIT,
                    threshold=SafetySetting.HarmBlockThreshold.OFF
                ),
                SafetySetting(
                    category=SafetySetting.HarmCategory.HARM_CATEGORY_HARASSMENT,
                    threshold=SafetySetting.HarmBlockThreshold.OFF
                ),
            ],
            stream=False
        )

        return responses.text.strip()
    except Exception as e:
        print(f"❌ Error generating answer: {str(e)}")
        return f"Error generating answer: {str(e)}"

def extract_questions_with_ai_direct(file_path, filename):
    """Extract questions directly from file using Vertex AI multimodal capabilities"""
    if not VERTEX_AVAILABLE:
        print("❌ Vertex AI not available")
        return []
        
    try:
        # Verify file exists and is readable
        if not os.path.exists(file_path):
            print(f"❌ File does not exist: {file_path}")
            return []
        
        file_size = os.path.getsize(file_path)
        if file_size == 0:
            print(f"❌ File is empty: {file_path}")
            return []
            
        print(f"📄 Processing file: {filename} ({file_size} bytes)")
        
        # Get file extension to determine how to handle it
        file_ext = filename.split(".")[-1].lower()
        
        # Customize prompt based on file type
        if file_ext in ["png", "jpg", "jpeg", "gif", "webp"]:
            prompt = """
Analyze the uploaded image/visual content and intelligently organize questions into logical sections with meaningful questions that can be answered based on what you can see in the image.

TASK: Create an intelligent sectioning system that:
1. Analyzes the content and identifies logical topic areas
2. Creates meaningful sections based on the content
3. Generates relevant questions for each section
4. Maps questions to appropriate sections

Focus on extracting questions about:
- Visual elements and content
- Text visible in the image (if any)
- Diagrams, charts, or data visualizations
- Objects, people, or scenes depicted
- Any processes or workflows shown
- Information that can be read or interpreted

Format your response as a valid JSON object with sections and questions, like this:
{{
  "sections": [
    {{
      "section_id": "visual_analysis",
      "section_title": "Visual Analysis",
      "section_description": "Questions about visual elements and overall composition",
      "questions": [
        {{
          "question": "What does this chart show?",
          "description": "Understanding the main data or information presented in the visualization"
        }},
        {{
          "question": "What are the key visual elements?",
          "description": "Identifying main components and their relationships"
        }}
      ]
    }},
    {{
      "section_id": "content_details",
      "section_title": "Content Details",
      "section_description": "Questions about specific content and information",
      "questions": [
        {{
          "question": "What processes are illustrated?",
          "description": "Understanding workflows or procedures shown in the image"
        }}
      ]
    }}
  ]
}}

IMPORTANT: Return ONLY the JSON object, no additional text or formatting.
""".format(filename=filename)
        else:
            prompt = """
Analyze the uploaded document and intelligently organize content into logical sections with meaningful questions based on the information provided.

TASK: Create an intelligent sectioning system that:
1. Analyzes the document content and identifies logical topic areas
2. Creates meaningful sections based on the document structure and content themes
3. Generates relevant questions for each section
4. Maps questions to appropriate sections based on content relevance

Extract questions from the document and organize them by:
- Main topics and themes
- Document structure (if applicable)
- Content areas and subjects
- Functional areas or processes
- Technical aspects vs business aspects
- Strategic vs operational content

Format your response as a valid JSON object with sections and questions, like this:
{{
  "sections": [
    {{
      "section_id": "overview_purpose",
      "section_title": "Overview & Purpose",
      "section_description": "Questions about the main objectives, goals, and overall scope",
      "questions": [
        {{
          "question": "What is the main purpose of this system?",
          "description": "Understanding the primary objective and goals described in the document"
        }},
        {{
          "question": "What are the key benefits mentioned?",
          "description": "Identifying the advantages and value propositions outlined"
        }}
      ]
    }},
    {{
      "section_id": "processes_workflow",
      "section_title": "Processes & Workflow",
      "section_description": "Questions about how things work and operational procedures",
      "questions": [
        {{
          "question": "How does the process work?",
          "description": "Detailed explanation of the workflow and steps involved"
        }},
        {{
          "question": "What are the key steps in implementation?",
          "description": "Understanding the sequence of actions required"
        }}
      ]
    }},
    {{
      "section_id": "technical_requirements",
      "section_title": "Technical Requirements",
      "section_description": "Questions about technical specifications and requirements",
      "questions": [
        {{
          "question": "What are the technical specifications?",
          "description": "Understanding the technical requirements and constraints"
        }}
      ]
    }}
  ]
}}

IMPORTANT: 
- Create 3-6 logical sections based on the actual document content
- Each section should have 2-5 relevant questions
- Section IDs should be snake_case and descriptive
- Return ONLY the JSON object, no additional text or formatting.
""".format(filename=filename)

        model = GenerativeModel("gemini-2.0-flash")
        
        # Prepare the content based on file type
        parts = [Part.from_text(prompt)]
        
        # Define MIME type mapping with proper Excel support
        mime_mapping = {
            "pdf": "application/pdf",
            "txt": "text/plain", 
            "csv": "text/csv",
            "json": "application/json",
            "png": "image/png",
            "jpg": "image/jpeg",
            "jpeg": "image/jpeg", 
            "gif": "image/gif",
            "webp": "image/webp",
        }
        
        # Files that Gemini can handle directly (Excel and Word files are NOT included)
        gemini_supported_files = [
            # Documents
            "pdf", "txt", "csv", "md", "markdown", "html", "htm", "xml", "yaml", "yml",
            
            # Images  
            "png", "jpg", "jpeg", "gif", "webp", "bmp", "tiff", "tif", "svg", "heic", "heif",
            
            # Code files (text-based)
            "py", "js", "ts", "jsx", "tsx", "java", "cpp", "c", "cs", "php", "rb", "go", "rs", 
            "swift", "kt", "scala", "sql", "sh", "ps1",
            
            # Configuration files
            "log", "cfg", "conf", "ini", "properties", "env", "lock", "gitignore", "dockerfile",
            
            # Audio (Gemini 2.0 supports audio)
            "mp3", "wav", "m4a", "aac", "ogg", "flac",
            
            # Video (Gemini 2.0 supports video)
            "mp4", "mov", "avi", "wmv", "webm", "mkv", "m4v"
        ]
        
        if file_ext in gemini_supported_files:
            try:
                # Read file as bytes for direct upload
                with open(file_path, "rb") as f:
                    file_data = f.read()
                
                mime_type = mime_mapping.get(file_ext, "application/octet-stream")
                
                # Validate file size for Gemini (max ~10MB)
                if len(file_data) > 10 * 1024 * 1024:  # 10MB limit
                    print(f"⚠️ File too large for direct upload ({len(file_data)} bytes), falling back to text extraction")
                    if file_ext not in ["png", "jpg", "jpeg", "gif", "webp"]:
                        content = extract_text_from_file(file_path, file_ext)
                        if content and not content.startswith("ERROR:"):
                            if len(content) > 30000:
                                content = content[:30000] + "..."
                            parts.append(Part.from_text(f"\nDocument Content:\n{content}"))
                        else:
                            return []
                    else:
                        return []  # Can't process large image
                else:
                    # Add file part directly
                    file_part = Part.from_data(data=file_data, mime_type=mime_type)
                    parts.append(file_part)
                    
                    print(f"✅ Using direct multimodal upload for {filename} ({mime_type}, {len(file_data)} bytes)")
                
            except Exception as e:
                print(f"⚠️ Direct upload failed for {filename}, falling back to text extraction: {e}")
                # Fallback to text extraction for non-image files
                if file_ext not in ["png", "jpg", "jpeg", "gif", "webp"]:
                    content = extract_text_from_file(file_path, file_ext)
                    if content and not content.startswith("ERROR:"):
                        if len(content) > 30000:
                            content = content[:30000] + "..."
                        parts.append(Part.from_text(f"\nDocument Content:\n{content}"))
                    else:
                        return []
                else:
                    return []  # Can't process image if direct upload fails
        else:
            # For Excel, Word, and other files - extract text first (Gemini doesn't support these directly)
            print(f"📄 Extracting text from {filename} (file type: {file_ext})")
            content = extract_text_from_file(file_path, file_ext)
            if not content or content.startswith("ERROR:"):
                print(f"❌ No content extracted from {filename}: {content}")
                return []
            
            # Limit content length
            if len(content) > 30000:
                content = content[:30000] + "..."
            
            parts.append(Part.from_text(f"\nDocument Content:\n{content}"))
        
        # Generate response
        print(f"🤖 Sending request to Gemini for {filename}")
        response = model.generate_content(
            parts,
            generation_config={
                "max_output_tokens": 4096,
                "temperature": 0.3,
                "top_p": 0.8,
            },
            safety_settings=[
                SafetySetting(
                    category=SafetySetting.HarmCategory.HARM_CATEGORY_HATE_SPEECH,
                    threshold=SafetySetting.HarmBlockThreshold.OFF
                ),
                SafetySetting(
                    category=SafetySetting.HarmCategory.HARM_CATEGORY_DANGEROUS_CONTENT,
                    threshold=SafetySetting.HarmBlockThreshold.OFF
                ),
                SafetySetting(
                    category=SafetySetting.HarmCategory.HARM_CATEGORY_SEXUALLY_EXPLICIT,
                    threshold=SafetySetting.HarmBlockThreshold.OFF
                ),
                SafetySetting(
                    category=SafetySetting.HarmCategory.HARM_CATEGORY_HARASSMENT,
                    threshold=SafetySetting.HarmBlockThreshold.OFF
                ),
            ],
            stream=False
        )
        
        print(f"✅ Received response from Gemini for {filename}")
        
        # Parse the JSON response
        response_text = response.text.strip()
        print(f"📝 Raw response length: {len(response_text)} characters")
        
        # Clean up the response - remove any markdown formatting
        if response_text.startswith("```json"):
            response_text = response_text.replace("```json", "").replace("```", "").strip()
        elif response_text.startswith("```"):
            response_text = response_text.replace("```", "").strip()
        
        # Try to extract JSON from the response
        try:
            parsed_response = json.loads(response_text)
            
            # Check if response has the new sectioned format
            if isinstance(parsed_response, dict) and "sections" in parsed_response:
                # New sectioned format
                sections_data = parsed_response["sections"]
                if isinstance(sections_data, list):
                    valid_sections = []
                    total_questions = 0
                    
                    for section in sections_data:
                        if isinstance(section, dict) and all(key in section for key in ["section_id", "section_title", "questions"]):
                            valid_questions = []
                            questions_list = section.get("questions", [])
                            
                            for q in questions_list:
                                if isinstance(q, dict) and "question" in q and "description" in q:
                                    valid_questions.append({
                                        "question": str(q["question"]).strip(),
                                        "description": str(q["description"]).strip()
                                    })
                                    total_questions += 1
                            
                            if valid_questions:  # Only add sections with valid questions
                                valid_sections.append({
                                    "section_id": str(section["section_id"]).strip(),
                                    "section_title": str(section["section_title"]).strip(),
                                    "section_description": str(section.get("section_description", "")).strip(),
                                    "questions": valid_questions[:5]  # Limit to 5 questions per section
                                })
                    
                    if valid_sections:
                        sectioned_result = {
                            "format": "sectioned",
                            "sections": valid_sections[:6],  # Limit to 6 sections
                            "total_sections": len(valid_sections),
                            "total_questions": total_questions
                        }
                        print(f"✅ Successfully parsed {len(valid_sections)} sections with {total_questions} questions from {filename}")
                        return sectioned_result
                    else:
                        print(f"⚠️ No valid sections found in response for {filename}")
                        return []
                else:
                    print(f"⚠️ Invalid sections structure in response for {filename}")
                    return []
            
            # Fallback: Check if it's the old format (direct list of questions)
            elif isinstance(parsed_response, list):
                valid_questions = []
                for q in parsed_response:
                    if isinstance(q, dict) and "question" in q and "description" in q:
                        valid_questions.append({
                            "question": str(q["question"]).strip(),
                            "description": str(q["description"]).strip()
                        })
                
                if valid_questions:
                    # Convert old format to new sectioned format with a single section
                    sectioned_result = {
                        "format": "sectioned",
                        "sections": [{
                            "section_id": "general_questions",
                            "section_title": "General Questions",
                            "section_description": "Questions extracted from the document",
                            "questions": valid_questions[:15]
                        }],
                        "total_sections": 1,
                        "total_questions": len(valid_questions)
                    }
                    print(f"✅ Successfully converted {len(valid_questions)} questions to sectioned format from {filename}")
                    return sectioned_result
                else:
                    print(f"⚠️ No valid questions found in list format for {filename}")
                    return []
            else:
                print(f"⚠️ Invalid JSON structure from AI for {filename} (unknown format)")
                return []
                
        except json.JSONDecodeError as e:
            print(f"⚠️ JSON parsing error for {filename}: {e}")
            print(f"Raw response preview: {response_text[:500]}...")
            
            # Fallback: try to extract questions manually using regex
            questions = []
            
            # Try to find JSON objects in the text
            json_pattern = r'\{\s*"question":\s*"([^"]+)"\s*,\s*"description":\s*"([^"]+)"\s*\}'
            matches = re.findall(json_pattern, response_text)
            
            for question_text, description_text in matches:
                questions.append({
                    "question": question_text.strip(),
                    "description": description_text.strip()
                })
            
            if questions:
                # Convert to sectioned format
                sectioned_result = {
                    "format": "sectioned",
                    "sections": [{
                        "section_id": "extracted_questions",
                        "section_title": "Extracted Questions",
                        "section_description": "Questions extracted using fallback parsing",
                        "questions": questions[:15]
                    }],
                    "total_sections": 1,
                    "total_questions": len(questions)
                }
                print(f"✅ Extracted {len(questions)} questions using fallback regex for {filename}")
                return sectioned_result
            else:
                print(f"❌ Could not extract any questions from {filename}")
                return []
            
    except Exception as e:
        print(f"❌ Error extracting questions from {filename}: {str(e)}")
        traceback.print_exc()
        return []

# ================================
# FLASK ROUTES
# ================================

@app.route("/", methods=["GET"])
def index():
    """Root endpoint - simple health check with file support information"""
    supported, missing = check_file_processing_dependencies()
    
    return jsonify({
        "status": "online",
        "message": "Welcome to CareAI API - Enhanced Universal File Support",
        "version": "3.0.0",
        "features": [
            "Universal file type support",
            "Enhanced Excel/PDF/Word processing", 
            "Automatic encoding detection",
            "Smart CSV parsing",
            "Multimodal question extraction",
            "Real-time progress tracking",
            "Comprehensive error handling"
        ],
        "file_support": {
            "always_supported": supported["always"],
            "currently_available": supported["conditional"],
            "missing_dependencies": missing
        }
    })
# Add this function after postAgentResultsToBackend
def postQuestionsToBackend(rfpId, questions, auth_token):
    """Post extracted questions to the backend questions API"""
    if not rfpId or not questions:
        return

    try:
        backend_api_url = os.environ.get("BACKEND_API_URL", "http://localhost:8080")
        questions_url = f"{backend_api_url}/api/rfps/questions"
        
        headers = {"Content-Type": "application/json"}
        if auth_token:
            headers["Authorization"] = f"Bearer {auth_token}"
        
        # Format questions for the API
        formatted_questions = []
        for q in questions:
            if isinstance(q, dict) and "question" in q and "description" in q:
                formatted_questions.append({
                    "question": q["question"],
                    "description": q["description"]
                })
        
        if not formatted_questions:
            print("⚠️ No valid questions to post to backend")
            return

        payload = {
            "rfpId": rfpId,
            "questions": formatted_questions
        }

        print(f"📝 Posting {len(formatted_questions)} questions to backend")
        
        response = requests.post(
            questions_url, 
            json=payload, 
            headers=headers, 
            timeout=30
        )

        if response.status_code in [200, 201]:
            print(f"✅ Successfully posted {len(formatted_questions)} questions to backend")
        else:
            print(f"⚠️ Failed to post questions: {response.status_code} - {response.text[:200]}")

    except Exception as e:
        print(f"❌ Error posting questions to backend: {str(e)}")

def postAgentResultsAndQuestions(rfpId, agentResults, auth_token, agent_id=None):
    """Post both agent results and questions to backend"""
    try:
        # Post agent results
        backend_api_url = os.environ.get("BACKEND_API_URL", "http://localhost:8080")
        agent_response_url = f"{backend_api_url}/api/rfps/{rfpId}/agent-response"
        
        headers = {"Content-Type": "application/json"}
        if auth_token:
            headers["Authorization"] = f"Bearer {auth_token}"

        # Extract file extension from processed files
        file_extension = None
        processed_files = agentResults.get("processed_files", [])
        if processed_files and len(processed_files) > 0:
            # Get file extension from the first processed file
            file_extension = processed_files[0].get("file_type", "")
        print(f"📂 Detected file extension: {file_extension}")
        response = requests.post(
            agent_response_url,
            json={
                "agentId": agent_id,
                "questions": agentResults.get("questions", []),
                "fileExtension": file_extension,  # ADD THIS LINE
            },
            headers=headers,
            timeout=30
        )

        if response.status_code in [200, 201]:
            print("✅ Agent results posted to backend successfully")
            
            # **NEW: Also post individual questions**
            questions = agentResults.get("questions", [])
            if questions and len(questions) > 0:
                postQuestionsToBackend(rfpId, questions, auth_token)
        else:
            print(f"⚠️ Failed to post agent results: {response.status_code} - {response.text[:200]}")

    except Exception as e:
        print(f"❌ Error posting agent results: {str(e)}")


@app.route("/health", methods=["GET"])
def health_check():
    """Enhanced health check endpoint with keep-alive status"""
    global keep_alive_thread, keep_alive_running
    
    # Calculate uptime
    uptime_seconds = time.time() - start_time
    uptime_hours = uptime_seconds / 3600
    
    # Check keep-alive status
    keep_alive_status = {
        "enabled": os.environ.get("KEEP_ALIVE_ENABLED", "false").lower() == "true",
        "running": keep_alive_running and keep_alive_thread and keep_alive_thread.is_alive(),
        "thread_alive": keep_alive_thread.is_alive() if keep_alive_thread else False
    }
    
    return jsonify({
        "status": "healthy",
        "timestamp": time.time(),
        "uptime_seconds": round(uptime_seconds, 2),
        "uptime_hours": round(uptime_hours, 2),
        "keep_alive": keep_alive_status,
        "server_info": {
            "port": int(os.environ.get("PORT", "8002")),
            "debug_mode": os.environ.get("DEBUG", "0") == "1",
            "threads": int(os.environ.get("WAITRESS_THREADS", "8"))
        }
    }), 200

@app.route("/status", methods=["GET"])
def service_status():
    """Detailed service status endpoint with file support info"""
    with progress_lock:
        active_uploads = len(upload_progress)
    
    supported, missing = check_file_processing_dependencies()
    
    status = {
        "status": "healthy",
        "timestamp": time.time(),
        "services": {
            "firebase": FIREBASE_AVAILABLE,
            "openai": OPENAI_AVAILABLE,
            "vertex_ai": VERTEX_AVAILABLE
        },
        "active_uploads": active_uploads,
        "version": "3.0.0",
        "file_processing": {
            "supported_types": supported["always"] + supported["conditional"],
            "missing_dependencies": missing,
            "categories": {
                "documents": ["pdf", "docx", "txt", "rtf"],
                "spreadsheets": ["xlsx", "xls", "csv"],
                "data": ["json", "xml"],
                "images": ["png", "jpg", "jpeg", "gif", "webp"],
                "presentations": ["pptx"]
            }
        },
        "capabilities": {
            "multimodal_analysis": VERTEX_AVAILABLE,
            "direct_file_upload": VERTEX_AVAILABLE,
            "image_processing": VERTEX_AVAILABLE,
            "document_embeddings": OPENAI_AVAILABLE and FIREBASE_AVAILABLE,
            "smart_encoding_detection": True,
            "comprehensive_error_handling": True
        }
    }
    
    return jsonify(status), 200

# Keep-alive endpoints for Render server
@app.route("/api/wake-up", methods=["GET"])
def wake_up_server():
    """Wake-up endpoint to prevent Render server from sleeping"""
    return jsonify({
        "status": "awake",
        "message": "Server is active and responding",
        "timestamp": time.time(),
        "uptime_seconds": time.time() - start_time if 'start_time' in globals() else 0,
        "keep_alive_enabled": os.environ.get("KEEP_ALIVE_ENABLED", "false").lower() == "true"
    }), 200

@app.route("/supported-files", methods=["GET"])
def get_supported_file_types():
    """Return detailed information about supported file types"""
    supported, missing = check_file_processing_dependencies()
    
    return jsonify({
        "supported_categories": {
            "documents": {
                "types": ["pdf", "docx", "txt", "rtf"],
                "description": "Text documents and reports"
            },
            "spreadsheets": {
                "types": ["xlsx", "xls", "csv"],
                "description": "Tabular data and calculations"
            },
            "data_formats": {
                "types": ["json", "xml"],
                "description": "Structured data formats"
            },
            "images": {
                "types": ["png", "jpg", "jpeg", "gif", "webp"],
                "description": "Visual content and diagrams"
            },
            "presentations": {
                "types": ["pptx"],
                "description": "Slide presentations"
            }
        },
        "currently_available": supported["always"] + supported["conditional"],
        "missing_dependencies": missing,
        "installation_commands": {
            "essential": "pip install openpyxl xlrd python-docx PyPDF2",
            "extended": "pip install python-pptx xmltodict chardet",
            "complete": "pip install openpyxl xlrd python-docx PyPDF2 python-pptx xmltodict chardet"
        },
        "dependency_status": {
            dep: dep not in [m.split()[0] for m in missing] 
            for dep in ["openpyxl", "xlrd", "python-docx", "PyPDF2", "python-pptx"]
        }
    })

# @app.route("/extract-questions", methods=["POST", "OPTIONS"])
# def extract_questions_from_files():
#     """Extract questions from multiple uploaded files using AI (multimodal)"""
#     if request.method == "OPTIONS":
#         return "", 200
        
#     if not VERTEX_AVAILABLE:
#         return jsonify({"error": "AI generation service is not available"}), 503
        
#     try:
#         org_id = request.args.get("orgId")
#         extraction_id = request.args.get("extractionId", str(uuid.uuid4()))
        
#         print(f"🤖 Starting multimodal question extraction with ID: {extraction_id} for org: {org_id}")
        
#         if not org_id:
#             return jsonify({"error": "orgId is required"}), 400

#         # Check if files were uploaded and validate
#         if not request.files:
#             return jsonify({"error": "No files provided"}), 400

#         uploaded_files = request.files.getlist('files')  # Get multiple files
#         if not uploaded_files:
#             return jsonify({"error": "No files found in request"}), 400
        
#         # Check for empty files
#         valid_files = []
#         for file in uploaded_files:
#             if file.filename and file.filename.strip():
#                 valid_files.append(file)
        
#         if not valid_files:
#             return jsonify({"error": "No valid files found (all files appear to be empty or unnamed)"}), 400

#         print(f"📁 Processing {len(valid_files)} valid files")

#         # Read all files into memory before starting async processing
#         file_data_list = []
#         for file in valid_files:
#             try:
#                 file.seek(0)
#                 file_content = file.read()
                
#                 if len(file_content) == 0:
#                     print(f"⚠️ Skipping empty file: {file.filename}")
#                     continue
                    
#                 file_data_list.append({
#                     'filename': file.filename,
#                     'content': file_content,
#                     'size': len(file_content)
#                 })
#                 print(f"✅ Read {file.filename} ({len(file_content)} bytes)")
                
#             except Exception as e:
#                 print(f"❌ Error reading file {file.filename}: {str(e)}")
#                 return jsonify({"error": f"Failed to read file {file.filename}: {str(e)}"}), 400
#             finally:
#                 try:
#                     file.close()
#                 except:
#                     pass
        
#         if not file_data_list:
#             return jsonify({"error": "No files could be read successfully"}), 400

#         # Initialize progress tracking
#         update_upload_progress(extraction_id, "Processing", 0, "Starting multimodal question extraction", "")
        
#         def process_files_async():
#             try:
#                 results = []
#                 total_files = len(file_data_list)
                
#                 # Get supported file types
#                 supported, missing = check_file_processing_dependencies()
#                 all_supported = supported["always"] + supported["conditional"] + ["png", "jpg", "jpeg", "gif", "webp"]
                
#                 for index, file_data in enumerate(file_data_list):
#                     filename = file_data['filename']
#                     file_content = file_data['content']
#                     file_ext = filename.split(".")[-1].lower()
                    
#                     # Validate file type
#                     if file_ext not in all_supported:
#                         results.append({
#                             "filename": filename,
#                             "status": "error",
#                             "error": f"Unsupported file type: {file_ext}. Supported: {', '.join(sorted(all_supported))}",
#                             "questions": []
#                         })
#                         continue
                    
#                     try:
#                         # Update progress
#                         progress = (index / total_files) * 80
#                         update_upload_progress(extraction_id, "Processing", progress, 
#                                             f"Processing file {index + 1}/{total_files}: {filename}")
                        
#                         # Save file content to temporary file
#                         save_path = os.path.join(UPLOAD_FOLDER, f"{extraction_id}_{index}_{filename}")
#                         with open(save_path, 'wb') as f:
#                             f.write(file_content)
                        
#                         update_upload_progress(extraction_id, "Processing", progress + 2, 
#                                             f"File saved: {filename}")
                        
#                         # Use direct file upload to AI
#                         update_upload_progress(extraction_id, "Processing", progress + 5, 
#                                             f"AI analyzing {filename}")
                        
#                         questions = extract_questions_with_ai_direct(save_path, filename)
                        
#                         if not questions:
#                             results.append({
#                                 "filename": filename,
#                                 "status": "error", 
#                                 "error": "AI could not extract questions from this file",
#                                 "questions": []
#                             })
#                         else:
#                             results.append({
#                                 "filename": filename,
#                                 "status": "success",
#                                 "file_type": file_ext,
#                                 "questions": questions,
#                                 "question_count": len(questions)
#                             })
                            
#                             print(f"✅ Extracted {len(questions)} questions from {filename}")
                        
#                         # Clean up temporary file
#                         try:
#                             os.remove(save_path)
#                         except Exception as e:
#                             print(f"⚠️ Failed to remove temp file {save_path}: {e}")
                            
#                     except Exception as file_error:
#                         print(f"❌ Error processing file {filename}: {str(file_error)}")
#                         traceback.print_exc()
#                         results.append({
#                             "filename": filename,
#                             "status": "error",
#                             "error": str(file_error),
#                             "questions": []
#                         })
                
#                 # Final processing - compile all questions
#                 update_upload_progress(extraction_id, "Processing", 90, "Compiling results")
                
#                 # Aggregate results
#                 all_questions = []
#                 successful_files = []
#                 failed_files = []
                
#                 for result in results:
#                     if result["status"] == "success":
#                         successful_files.append(result["filename"])
#                         all_questions.extend(result["questions"])
#                     else:
#                         failed_files.append({
#                             "filename": result["filename"],
#                             "error": result.get("error", "Unknown error")
#                         })
                
#                 # Store results in progress for retrieval
#                 final_results = {
#                     "extraction_id": extraction_id,
#                     "org_id": org_id,
#                     "total_files": total_files,
#                     "successful_files": len(successful_files),
#                     "failed_files": len(failed_files),
#                     "total_questions": len(all_questions),
#                     "files_processed": results,
#                     "all_questions": all_questions,
#                     "file_details": {
#                         "successful": successful_files,
#                         "failed": failed_files
#                     },
#                     "processing_method": "enhanced_multimodal",
#                     "supported_types": all_supported,
#                     "timestamp": datetime.now().isoformat()
#                 }
                
#                 # Update progress with final results
#                 update_upload_progress(extraction_id, "Completed", 100, "Enhanced question extraction completed", "")
                
#                 # Store results in progress data for retrieval
#                 with progress_lock:
#                     if extraction_id in upload_progress:
#                         upload_progress[extraction_id]["results"] = final_results
                
#                 print(f"✅ Enhanced question extraction completed. Total questions: {len(all_questions)}")
                
#             except Exception as e:
#                 error_msg = str(e)
#                 print(f"❌ Async processing error: {error_msg}")
#                 traceback.print_exc()
#                 update_upload_progress(extraction_id, "error", 0, f"Error: {error_msg}", "")
        
#         # Start processing in background thread
#         processing_thread = Thread(target=process_files_async)
#         processing_thread.daemon = False
#         processing_thread.start()
        
#         return jsonify({
#             "message": "Enhanced question extraction started. Check status with the extraction-status endpoint.",
#             "extraction_id": extraction_id,
#             "org_id": org_id,
#             "files_count": len(file_data_list),
#             "processing_method": "enhanced_multimodal",
#             "supported_types": ["documents", "images", "spreadsheets", "presentations", "data_files"]
#         }), 202

#     except Exception as e:
#         error_msg = str(e)
#         print(f"❌ Question Extraction Error: {error_msg}")
#         traceback.print_exc()
        
#         return jsonify({"error": "Internal server error", "details": error_msg}), 500

# @app.route("/extraction-status", methods=["GET", "OPTIONS"])
# def get_extraction_status():
#     """Get question extraction status and results"""
#     if request.method == "OPTIONS":
#         return "", 200
        
#     extraction_id = request.args.get("extractionId")
    
#     if not extraction_id:
#         return jsonify({"error": "extractionId is required"}), 400
    
#     with progress_lock:
#         exists = extraction_id in upload_progress
#         status_data = upload_progress.get(extraction_id, None)
    
#     if not exists:
#         return jsonify({"exists": False}), 200
    
#     response_data = {
#         "exists": True,
#         "status": status_data
#     }
    
#     # Include results if extraction is completed
#     if status_data and status_data.get("status") == "Completed" and "results" in status_data:
#         response_data["results"] = status_data["results"]
    
#     return jsonify(response_data), 200

# @app.route("/download-questions", methods=["GET", "OPTIONS"])
# def download_questions():
#     """Download extracted questions as JSON or text file"""
#     if request.method == "OPTIONS":
#         return "", 200
        
#     try:
#         extraction_id = request.args.get("extractionId")
#         format_type = request.args.get("format", "json")  # json or txt
        
#         if not extraction_id:
#             return jsonify({"error": "extractionId is required"}), 400
        
#         with progress_lock:
#             status_data = upload_progress.get(extraction_id, None)
        
#         if not status_data or "results" not in status_data:
#             return jsonify({"error": "No results found for this extraction ID"}), 404
        
#         results = status_data["results"]
        
#         if format_type == "txt":
#             # Generate text format
#             content = f"Enhanced Question Extraction Results\n"
#             content += f"=" * 50 + "\n"
#             content += f"Extraction ID: {extraction_id}\n"
#             content += f"Organization: {results['org_id']}\n"
#             content += f"Processing Method: {results.get('processing_method', 'enhanced_multimodal')}\n"
#             content += f"Total Files: {results['total_files']}\n"
#             content += f"Successful Files: {results['successful_files']}\n"
#             content += f"Total Questions: {results['total_questions']}\n"
#             content += f"Timestamp: {results['timestamp']}\n\n"
            
#             content += "=" * 50 + "\n"
#             content += "ALL EXTRACTED QUESTIONS\n"
#             content += "=" * 50 + "\n\n"
            
#             for i, question_obj in enumerate(results['all_questions'], 1):
#                 content += f"{i}. {question_obj['question']}\n"
#                 content += f"   Description: {question_obj['description']}\n\n"
            
#             content += "\n" + "=" * 50 + "\n"
#             content += "FILE PROCESSING DETAILS\n"
#             content += "=" * 50 + "\n\n"
            
#             for file_result in results['files_processed']:
#                 content += f"File: {file_result['filename']}\n"
#                 content += f"Type: {file_result.get('file_type', 'unknown')}\n"
#                 content += f"Status: {file_result['status']}\n"
#                 if file_result['status'] == 'success':
#                     content += f"Questions extracted: {file_result['question_count']}\n"
#                     for q in file_result['questions']:
#                         content += f"  - {q['question']}\n"
#                         content += f"    {q['description']}\n"
#                 else:
#                     content += f"Error: {file_result.get('error', 'Unknown error')}\n"
#                 content += "\n"
            
#             return Response(
#                 content,
#                 mimetype="text/plain",
#                 headers={
#                     "Content-Disposition": f"attachment; filename=questions_{extraction_id}.txt"
#                 }
#             )
#         else:
#             # Return JSON format
#             return Response(
#                 json.dumps(results, indent=2),
#                 mimetype="application/json",
#                 headers={
#                     "Content-Disposition": f"attachment; filename=questions_{extraction_id}.json"
#                 }
#             )
            
#     except Exception as e:
#         error_msg = str(e)
#         print(f"❌ Download Error: {error_msg}")
#         return jsonify({"error": "Internal server error", "details": error_msg}), 500

@app.route("/upload", methods=["POST", "OPTIONS"])
def upload_file():
    """Enhanced upload endpoint with Redis progress tracking"""
    if request.method == "OPTIONS":
        return "", 200
    
    # Enhanced: Use Redis for progress tracking when available
    if REDIS_AVAILABLE:
        print("🚀 Using Redis-enhanced threading approach")
    else:
        print("⚠️ Using memory-based progress tracking")
    if not FIREBASE_AVAILABLE:
        return jsonify({"error": "Firebase is not available"}), 503
        
    if not OPENAI_AVAILABLE:
        return jsonify({"error": "OpenAI embedding service is not available"}), 503
        
    save_path = None
    try:
        org_id = request.args.get("orgId")
        file_id = request.args.get("fileId")
        upload_id = request.args.get("uploadId", str(uuid.uuid4()))
        
        print(f"⏳ Starting enhanced upload with fileId: {file_id}, uploadId: {upload_id} for org: {org_id}")
        
        if not org_id:
            return jsonify({"error": "orgId is required"}), 400
            
        if not file_id:
            return jsonify({"error": "fileId is required"}), 400

        if "file" not in request.files:
            return jsonify({"error": "No file provided."}), 400

        file = request.files["file"]
        filename = file.filename
        file_ext = filename.split(".")[-1].lower()

        # Get currently supported types
        supported, missing = check_file_processing_dependencies()
        all_supported = supported["always"] + supported["conditional"] + ["png", "jpg", "jpeg", "gif", "webp"]
        
        if file_ext not in all_supported:
            return jsonify({
                "error": f"Unsupported file type: {file_ext}. Supported: {', '.join(sorted(all_supported))}",
                "missing_dependencies": missing,
                "install_command": f"pip install {' '.join(missing)}" if missing else None
            }), 400

        # Initialize progress
        update_upload_progress(upload_id, "Processing", 0, "Starting enhanced upload", filename)
        
        # Save file temporarily
        save_path = os.path.abspath(os.path.join(UPLOAD_FOLDER, f"{upload_id}_{filename}"))
        file.save(save_path)
        print(f"✅ File saved to {save_path}")
        
        # Process file in a separate thread
        def process_file_async():
            try:
                update_upload_progress(upload_id, "Processing", 25, "File saved", filename)
                
                # Choose processing method based on USE_LLAMAINDEX flag
                if USE_LLAMAINDEX:
                    print(f"🚀 Processing with LlamaIndex (skipping traditional chunking)...")
                    update_upload_progress(upload_id, "Processing", 50, "Processing with LlamaIndex", filename)
                    
                    try:
                        from llamaindex_integration import process_file_with_llamaindex
                        
                        result = process_file_with_llamaindex(
                            file_path=save_path,
                            file_id=file_id,
                            org_id=org_id,
                            user_id=user_id,
                            filename=filename
                        )
                        
                        if result.get("success"):
                            update_upload_progress(upload_id, "completed", 100, "File processed successfully with LlamaIndex", filename)
                            print(f"✅ LlamaIndex processing complete: {result}")
                            
                            # Update backend to mark as having embeddings
                            update_backend_embedding_status(file_id, org_id, True)
                            
                            # LlamaIndex handles everything including storage, so we can return early
                            print(f"✅ LlamaIndex handled storage to NeonDB directly")
                            return  # Exit successfully - no need for traditional storage
                        else:
                            print(f"❌ LlamaIndex processing failed, falling back to traditional processing")
                            raise Exception("LlamaIndex processing failed")
                            
                    except Exception as e:
                        print(f"❌ LlamaIndex error: {str(e)}")
                        print(f"🔄 Falling back to traditional processing...")
                        # Fall through to traditional processing
                        
                        # Extract text and create chunks (for embedding)
                        update_upload_progress(upload_id, "Processing", 50, "Extracting text with enhanced processing", filename)
                        chunks = parse_and_chunk(save_path, file_ext, chunk_size=50, max_chunks=500)
                        
                        if not chunks:
                            print(f"❌ No content extracted from {filename}")
                            update_upload_progress(upload_id, "error", 0, "No content extracted", filename)
                            # Update backend to mark as NOT having embeddings
                            update_backend_embedding_status(file_id, org_id, False)
                            return
                            
                        # Generate embeddings
                        update_upload_progress(upload_id, "Processing", 75, "Generating embeddings", filename)
                else:
                    # Traditional processing
                    # Extract text and create chunks (for embedding)
                    update_upload_progress(upload_id, "Processing", 50, "Extracting text with enhanced processing", filename)
                    chunks = parse_and_chunk(save_path, file_ext, chunk_size=50, max_chunks=500)
                    
                    if not chunks:
                        print(f"❌ No content extracted from {filename}")
                        update_upload_progress(upload_id, "error", 0, "No content extracted", filename)
                        # Update backend to mark as NOT having embeddings
                        update_backend_embedding_status(file_id, org_id, False)
                        return
                        
                    # Generate embeddings
                    update_upload_progress(upload_id, "Processing", 75, "Generating embeddings", filename)
                embeddings = embed_chunks(chunks, upload_id=upload_id, org_id=org_id, filename=filename)
                
                # Store in Firestore using fileId as document ID
                update_upload_progress(upload_id, "Processing", 90, "Storing embeddings", filename)
                
                file_doc_ref = db.collection("document_embeddings").document(f"org-{org_id}").collection("files").document(file_id)
                
                # Store file metadata
                file_doc_ref.set({
                    "filename": filename,
                    "file_id": file_id,
                    "upload_id": upload_id,
                    "file_type": file_ext,
                    "document_type": "global",
                    "created_at": firestore.SERVER_TIMESTAMP,
                    "chunk_count": len(chunks),
                    "processing_version": "3.0.0"
                })
                
                # Store chunks in batches - OPTIMIZED for 50x faster writes
                batch_size = 500  # Firestore allows up to 500 writes per batch
                total_chunks = len(chunks)
                
                print(f"🚀 Storing {total_chunks} chunks in batches of {batch_size}")
                
                for i in range(0, total_chunks, batch_size):
                    batch = db.batch()
                    end_idx = min(i + batch_size, total_chunks)
                    
                    for j in range(i, end_idx):
                        chunk_ref = file_doc_ref.collection("chunks").document(str(j))
                        batch.set(chunk_ref, {
                            "content": chunks[j],
                            "embedding": embeddings[j],
                            "index": j
                        })
                    
                    batch.commit()
                    print(f"✓ Stored batch {i//batch_size + 1}: {end_idx}/{total_chunks} chunks")
                    
                    # Only clean up memory for large batches
                    if end_idx - i >= 100:
                        del batch
                        import gc
                        gc.collect()
                    
                    progress = 90 + ((end_idx / total_chunks) * 10)
                    update_upload_progress(upload_id, "Processing", progress, 
                                        f"Storing embeddings ({end_idx}/{total_chunks})", filename)
                
                print(f"✅ Successfully processed file {filename} with fileId: {file_id}")
                update_upload_progress(upload_id, "Completed", 100, "Enhanced file processing completed", filename)
                
                # Update backend to mark file as having embeddings
                update_backend_embedding_status(file_id, org_id, True)
                
            except Exception as e:
                error_msg = str(e)
                print(f"❌ Async processing error: {error_msg}")
                traceback.print_exc()
                update_upload_progress(upload_id, "error", 0, f"Error: {error_msg}", filename)
                # Update backend to mark as NOT having embeddings on error
                update_backend_embedding_status(file_id, org_id, False)
                
            finally:
                if save_path and os.path.exists(save_path):
                    try:
                        os.remove(save_path)
                        print(f"🧹 File {save_path} deleted.")
                    except Exception as e:
                        print(f"Error deleting file: {e}")
        
        # Start processing in background thread
        processing_thread = Thread(target=process_file_async)
        processing_thread.daemon = False
        processing_thread.start()
        
        return jsonify({
            "message": "Enhanced file upload started. Check status with the upload-status endpoint.",
            "file_id": file_id,
            "upload_id": upload_id,
            "file_type": file_ext,
            "processing_version": "3.0.0"
        }), 202

    except Exception as e:
        error_msg = str(e)
        print(f"❌ Upload Error: {error_msg}")
        traceback.print_exc()
        
        if 'upload_id' in locals() and upload_id:
            update_upload_progress(upload_id, "error", 0, f"Error: {error_msg}", filename if 'filename' in locals() else "")
            
        if save_path and os.path.exists(save_path):
            try:
                os.remove(save_path)
                print(f"🧹 File {save_path} deleted after error.")
            except Exception as cleanup_error:
                print(f"Error cleaning up file: {cleanup_error}")
                
        return jsonify({"error": "Internal server error", "details": error_msg}), 500

@app.route("/upload-status", methods=["GET", "OPTIONS"])
def get_upload_status():
    """Enhanced upload status - checks both queue and traditional progress"""
    if request.method == "OPTIONS":
        return "", 200
        
    upload_id = request.args.get("uploadId")
    
    if not upload_id:
        return jsonify({"error": "uploadId is required"}), 400
    
    # First check Redis/queue system
    if QUEUE_AVAILABLE and redis_manager.is_connected():
        progress_data = redis_manager.get_progress(upload_id)
        if progress_data:
            return jsonify({
                "exists": True,
                "queue_mode": True,
                **progress_data
            })
    
    # Fall back to traditional progress tracking
    with progress_lock:
        exists = upload_id in upload_progress
        status_data = upload_progress.get(upload_id, None)
    
    if not exists and status_data is None:
        time.sleep(0.1)
        with progress_lock:
            exists = upload_id in upload_progress
            status_data = upload_progress.get(upload_id, None)
    
    if not exists:
        return jsonify({"exists": False}), 200
    
    return jsonify({
        "exists": True,
        "status": status_data
    }), 200


# @app.route("/update-embedding-status", methods=["POST", "OPTIONS"])
# def update_embedding_status():
#     """Update the isFromEmbedding flag for a file in the main database"""
#     if request.method == "OPTIONS":
#         return "", 200
        
#     try:
#         data = request.get_json(silent=True) or {}
#         file_id = data.get("fileId")
#         org_id = data.get("orgId")
#         is_from_embedding = data.get("isFromEmbedding")
        
#         if not file_id:
#             return jsonify({"error": "fileId is required"}), 400
            
#         if not org_id:
#             return jsonify({"error": "orgId is required"}), 400
            
#         if not isinstance(is_from_embedding, bool):
#             return jsonify({"error": "isFromEmbedding must be a boolean"}), 400
        
#         print(f"🔄 Manual status update for fileId: {file_id}, isFromEmbedding: {is_from_embedding}")
        
#         # Update backend status
#         success = update_backend_embedding_status(file_id, org_id, is_from_embedding)
        
#         if success:
#             return jsonify({
#                 "message": "Embedding status updated successfully",
#                 "fileId": file_id,
#                 "isFromEmbedding": is_from_embedding
#             }), 200
#         else:
#             return jsonify({
#                 "error": "Failed to update backend",
#                 "fileId": file_id
#             }), 500
            
#     except Exception as e:
#         error_msg = str(e)
#         print(f"❌ Update Status Error: {error_msg}")
#         traceback.print_exc()
#         return jsonify({"error": "Internal server error", "details": error_msg}), 500

# @app.route("/reprocess-file", methods=["POST", "OPTIONS"])
# def reprocess_file():
#     """Reprocess a file to create embeddings"""
#     if request.method == "OPTIONS":
#         return "", 200
        
#     if not FIREBASE_AVAILABLE:
#         return jsonify({"error": "Firebase is not available"}), 503
        
#     if not OPENAI_AVAILABLE:
#         return jsonify({"error": "OpenAI embedding service is not available"}), 503
        
#     try:
#         data = request.get_json(silent=True) or {}
#         file_id = data.get("fileId")
#         org_id = data.get("orgId")
#         file_url = data.get("fileUrl")  # URL to download the file
#         filename = data.get("filename")
        
#         if not all([file_id, org_id, file_url, filename]):
#             return jsonify({"error": "fileId, orgId, fileUrl, and filename are required"}), 400
        
#         print(f"🔄 Reprocessing file with enhanced processing: {filename} (fileId: {file_id})")
        
#         # Generate a new upload ID for tracking
#         upload_id = str(uuid.uuid4())
        
#         def reprocess_async():
#             save_path = None
#             try:
#                 update_upload_progress(upload_id, "Processing", 10, "Downloading file", filename)
                
#                 # Download file from URL
#                 file_response = requests.get(file_url, timeout=30)
#                 if file_response.status_code != 200:
#                     raise Exception(f"Failed to download file: {file_response.status_code}")
                
#                 # Save file temporarily
#                 file_ext = filename.split(".")[-1].lower()
#                 save_path = os.path.join(UPLOAD_FOLDER, f"{upload_id}_{filename}")
                
#                 with open(save_path, 'wb') as f:
#                     f.write(file_response.content)
                
#                 update_upload_progress(upload_id, "Processing", 30, "File downloaded", filename)
                
#                 # Extract text and create chunks
#                 update_upload_progress(upload_id, "Processing", 50, "Enhanced text extraction", filename)
#                 chunks = parse_and_chunk(save_path, file_ext, chunk_size=50, max_chunks=500)
                
#                 if not chunks:
#                     print(f"❌ No content extracted from {filename}")
#                     update_upload_progress(upload_id, "error", 0, "No content extracted", filename)
#                     update_backend_embedding_status(file_id, org_id, False)
#                     return
                
#                 # Generate embeddings
#                 update_upload_progress(upload_id, "Processing", 70, "Generating embeddings", filename)
#                 embeddings = embed_chunks(chunks, upload_id=upload_id, org_id=org_id, filename=filename)
                
#                 # Store in Firestore
#                 update_upload_progress(upload_id, "Processing", 90, "Storing embeddings", filename)
                
#                 file_doc_ref = db.collection("document_embeddings").document(f"org-{org_id}").collection("files").document(file_id)
                
#                 # Store file metadata
#                 file_doc_ref.set({
#                     "filename": filename,
#                     "file_id": file_id,
#                     "upload_id": upload_id,
#                     "file_type": file_ext,
#                     "created_at": firestore.SERVER_TIMESTAMP,
#                     "chunk_count": len(chunks),
#                     "reprocessed": True,
#                     "processing_version": "3.0.0"
#                 })
                
#                 # Store chunks in batches
#                 batch_size = 10
#                 total_chunks = len(chunks)
                
#                 for i in range(0, total_chunks, batch_size):
#                     batch = db.batch()
#                     end_idx = min(i + batch_size, total_chunks)
                    
#                     for j in range(i, end_idx):
#                         chunk_ref = file_doc_ref.collection("chunks").document(str(j))
#                         batch.set(chunk_ref, {
#                             "content": chunks[j],
#                             "embedding": embeddings[j],
#                             "index": j
#                         })
                    
#                     batch.commit()
#                     del batch
#                     import gc
#                     gc.collect()
                
#                 print(f"✅ Successfully reprocessed file {filename} with fileId: {file_id}")
#                 update_upload_progress(upload_id, "Completed", 100, "Enhanced reprocessing completed", filename)
                
#                 # Update backend status
#                 update_backend_embedding_status(file_id, org_id, True)
                
#             except Exception as e:
#                 error_msg = str(e)
#                 print(f"❌ Reprocessing error: {error_msg}")
#                 update_upload_progress(upload_id, "error", 0, f"Error: {error_msg}", filename)
#                 update_backend_embedding_status(file_id, org_id, False)
                
#             finally:
#                 if save_path and os.path.exists(save_path):
#                     try:
#                         os.remove(save_path)
#                     except Exception as e:
#                         print(f"Error deleting file: {e}")
        
#         # Start reprocessing in background
#         Thread(target=reprocess_async, daemon=False).start()
        
#         return jsonify({
#             "message": "Enhanced file reprocessing started",
#             "fileId": file_id,
#             "upload_id": upload_id,
#             "processing_version": "3.0.0"
#         }), 202
        
#     except Exception as e:
#         error_msg = str(e)
#         print(f"❌ Reprocess Error: {error_msg}")
#         return jsonify({"error": "Internal server error", "details": error_msg}), 500

@app.route("/delete", methods=["DELETE", "OPTIONS"])
def delete_file_by_file_id():
    """Delete a file and its embeddings using fileId (primary method)"""
    if request.method == "OPTIONS":
        return "", 200
        
    if not FIREBASE_AVAILABLE:
        return jsonify({"error": "Firebase is not available"}), 503
        
    try:
        org_id = request.args.get("orgId")
        file_id = request.args.get("fileId")
        filename = request.args.get("filename")
        upload_id = request.args.get("uploadId")
        rfp_id = request.args.get("rfpId")
        if not org_id:
            return jsonify({"error": "orgId is required"}), 400
            
        if not file_id and not filename and not upload_id:
            return jsonify({"error": "fileId, uploadId, or filename is required"}), 400

        print(f"🗑️ Attempting to delete embeddings for org: {org_id}, fileId: {file_id}")

        files_ref = db.collection("document_embeddings").document(f"org-{org_id}").collection("files")
        
        deleted_count = 0
        deletion_details = []

        # Method 1: Direct lookup by fileId (most reliable)
        if file_id:
            try:
                file_doc_ref = files_ref.document(file_id)
                file_doc = file_doc_ref.get()
                
                if file_doc.exists:
                    file_data = file_doc.to_dict()
                    print(f"✅ Found file document by fileId: {file_id}")
                    
                    # Delete all chunks in the file
                    chunks_deleted = delete_collection(file_doc_ref.collection("chunks"), 100)
                    
                    # Delete the file document
                    file_doc_ref.delete()
                    deleted_count += 1
                    
                    deletion_details.append({
                        "method": "fileId_direct",
                        "fileId": file_id,
                        "filename": file_data.get("filename", "unknown"),
                        "file_type": file_data.get("file_type", "unknown"),
                        "chunks_deleted": chunks_deleted,
                        "processing_version": file_data.get("processing_version", "legacy"),
                        "success": True
                    })
                    
                    print(f"✅ Successfully deleted file and {chunks_deleted} chunks by fileId: {file_id}")
                    
                    # Update backend status
                    update_backend_embedding_status(file_id, org_id, False)
                    
                else:
                    print(f"⚠️ No file found with fileId: {file_id}")
                    
            except Exception as e:
                print(f"❌ Error deleting by fileId {file_id}: {str(e)}")
                deletion_details.append({
                    "method": "fileId_direct",
                    "fileId": file_id,
                    "success": False,
                    "error": str(e)
                })

        # Response with detailed information
        if deleted_count > 0:
            return jsonify({
                "message": f"Successfully deleted {deleted_count} file(s) and their embeddings",
                "deleted_count": deleted_count,
                "deletion_details": deletion_details,
                "org_id": org_id
            }), 200
        else:
            return jsonify({
                "message": "No matching files found to delete",
                "deleted_count": 0,
                "deletion_details": deletion_details,
                "org_id": org_id,
                "warning": "File may have already been deleted or never existed"
            }), 404

    except Exception as e:
        error_msg = str(e)
        print(f"❌ Deletion Error: {error_msg}")
        traceback.print_exc()
        return jsonify({
            "error": "Internal server error", 
            "details": error_msg
        }), 500

@app.route("/api/v2/bulk-delete", methods=["DELETE", "OPTIONS"])
def bulk_delete_files_v2():
    """Bulk delete files and their embeddings/chunks from NeonDB - V2 API"""
    if request.method == "OPTIONS":
        return "", 200
        
    try:
        data = request.get_json(silent=True) or {}
        org_id = data.get("orgId")
        file_ids = data.get("fileIds", [])
        
        if not org_id:
            return jsonify({"error": "orgId is required"}), 400
            
        if not file_ids or not isinstance(file_ids, list):
            return jsonify({"error": "fileIds array is required"}), 400
            
        if len(file_ids) > 50:
            return jsonify({"error": "Cannot delete more than 50 files at once"}), 400

        print(f"🗑️ Bulk deleting embeddings for {len(file_ids)} files in org: {org_id}")

        # Use asyncio to run the deletion function
        result = asyncio.run(bulk_delete_from_neondb(org_id, file_ids))
        
        return jsonify(result["response"]), result["status_code"]

    except Exception as e:
        error_msg = str(e)
        print(f"❌ Bulk deletion error: {error_msg}")
        traceback.print_exc()
        return jsonify({
            "error": "Internal server error", 
            "details": error_msg
        }), 500


async def bulk_delete_from_neondb(org_id: str, file_ids: list) -> dict:
    """Delete files and their embeddings from NeonDB tables"""
    import asyncpg
    
    db_url = os.getenv("NEON_DATABASE_URL")
    if not db_url:
        return {
            "response": {"error": "NeonDB connection not configured"},
            "status_code": 503
        }
    
    conn = await asyncpg.connect(db_url)
    
    try:
        deletion_results = []
        total_chunks_deleted = 0
        successful_deletions = 0

        # Process each file
        for file_id in file_ids:
            try:
                # Count chunks in llamaindex_embeddings table before deletion
                llamaindex_count = await conn.fetchval("""
                    SELECT COUNT(*) FROM llamaindex_embeddings 
                    WHERE file_id = $1 AND org_id = $2
                """, file_id, org_id)
                
                if llamaindex_count > 0:
                    # Get filename for reporting
                    filename = await conn.fetchval("""
                        SELECT metadata->>'filename' FROM llamaindex_embeddings 
                        WHERE file_id = $1 AND org_id = $2 AND metadata->>'filename' IS NOT NULL
                        LIMIT 1
                    """, file_id, org_id)
                    
                    # Delete from llamaindex_embeddings table
                    await conn.execute("""
                        DELETE FROM llamaindex_embeddings 
                        WHERE file_id = $1 AND org_id = $2
                    """, file_id, org_id)
                    
                    total_chunks_deleted += llamaindex_count
                    successful_deletions += 1
                    
                    deletion_results.append({
                        "fileId": file_id,
                        "filename": filename or "unknown",
                        "chunks_deleted": llamaindex_count,
                        "success": True
                    })
                    
                    print(f"✅ Deleted file {file_id}: {llamaindex_count} chunks from llamaindex_embeddings")
                    
                else:
                    deletion_results.append({
                        "fileId": file_id,
                        "success": False,
                        "error": "File not found in embedding database"
                    })
                    print(f"⚠️ File {file_id} not found in embeddings")
                    
            except Exception as e:
                deletion_results.append({
                    "fileId": file_id,
                    "success": False,
                    "error": str(e)
                })
                print(f"❌ Error deleting {file_id}: {str(e)}")

        return {
            "response": {
                "success": True,
                "message": f"Bulk deletion completed: {successful_deletions}/{len(file_ids)} files deleted",
                "data": {
                    "total_files_processed": len(file_ids),
                    "successful_deletions": successful_deletions,
                    "total_chunks_deleted": total_chunks_deleted,
                    "deletion_results": deletion_results
                }
            },
            "status_code": 200
        }

    except Exception as e:
        print(f"❌ NeonDB bulk deletion error: {str(e)}")
        return {
            "response": {
                "error": "Database error during bulk deletion",
                "details": str(e)
            },
            "status_code": 500
        }
    finally:
        await conn.close()

# @app.route("/cleanup-orphaned", methods=["POST", "OPTIONS"])
# def cleanup_orphaned_embeddings():
#     """Clean up orphaned embeddings using fileIds from database"""
#     if request.method == "OPTIONS":
#         return "", 200
        
#     if not FIREBASE_AVAILABLE:
#         return jsonify({"error": "Firebase is not available"}), 503
        
#     try:
#         data = request.get_json(silent=True) or {}
#         org_id = data.get("orgId")
#         active_file_ids = data.get("activeFileIds", [])
        
#         if not org_id:
#             return jsonify({"error": "orgId is required"}), 400
            
#         print(f"🧹 Starting enhanced cleanup for org: {org_id}")
#         print(f"📋 Active fileIds to preserve: {len(active_file_ids)}")
        
#         files_ref = db.collection("document_embeddings").document(f"org-{org_id}").collection("files")
#         all_firestore_files = files_ref.stream()
        
#         orphaned_files = []
#         preserved_files = []
        
#         for file_doc in all_firestore_files:
#             doc_id = file_doc.id
#             file_data = file_doc.to_dict()
            
#             stored_file_id = file_data.get("file_id")
            
#             is_active = (
#                 doc_id in active_file_ids or 
#                 stored_file_id in active_file_ids
#             )
            
#             if not is_active:
#                 print(f"🗑️ Found orphaned file: {doc_id} ({file_data.get('filename', 'unknown')})")
                
#                 chunks_deleted = delete_collection(file_doc.reference.collection("chunks"), 100)
#                 file_doc.reference.delete()
                
#                 orphaned_files.append({
#                     "document_id": doc_id,
#                     "file_id": stored_file_id,
#                     "filename": file_data.get("filename", "unknown"),
#                     "file_type": file_data.get("file_type", "unknown"),
#                     "processing_version": file_data.get("processing_version", "legacy"),
#                     "chunks_deleted": chunks_deleted
#                 })
                
#                 # Update backend status
#                 if stored_file_id:
#                     update_backend_embedding_status(stored_file_id, org_id, False)
                    
#             else:
#                 preserved_files.append({
#                     "document_id": doc_id,
#                     "file_id": stored_file_id,
#                     "filename": file_data.get("filename", "unknown"),
#                     "processing_version": file_data.get("processing_version", "legacy")
#                 })
        
#         print(f"✅ Enhanced cleanup complete. Deleted {len(orphaned_files)} orphaned files, preserved {len(preserved_files)} active files")
        
#         return jsonify({
#             "message": f"Enhanced cleanup complete for org {org_id}",
#             "orphaned_files_deleted": len(orphaned_files),
#             "active_files_preserved": len(preserved_files),
#             "deletion_details": orphaned_files,
#             "cleanup_version": "3.0.0"
#         }), 200
        
#     except Exception as e:
#         error_msg = str(e)
#         print(f"❌ Cleanup Error: {error_msg}")
#         traceback.print_exc()
#         return jsonify({"error": "Internal server error", "details": error_msg}), 500

# @app.route("/files", methods=["GET", "OPTIONS"])
# def list_files():
#     """List all files for an organization with enhanced fileId information"""
#     if request.method == "OPTIONS":
#         return "", 200
        
#     if not FIREBASE_AVAILABLE:
#         return jsonify({"error": "Firebase is not available"}), 503
        
#     try:
#         org_id = request.args.get("orgId")
        
#         if not org_id:
#             return jsonify({"error": "orgId is required"}), 400
        
#         files_ref = db.collection("document_embeddings").document(f"org-{org_id}").collection("files")
#         files_docs = files_ref.stream()
        
#         files = []
#         for doc in files_docs:
#             file_data = doc.to_dict()
#             files.append({
#                 "document_id": doc.id,
#                 "file_id": file_data.get("file_id"),
#                 "upload_id": file_data.get("upload_id"),
#                 "filename": file_data.get("filename"),
#                 "file_type": file_data.get("file_type"),
#                 "created_at": file_data.get("created_at"),
#                 "chunk_count": file_data.get("chunk_count"),
#                 "reprocessed": file_data.get("reprocessed", False),
#                 "processing_version": file_data.get("processing_version", "legacy")
#             })
        
#         return jsonify({
#             "org_id": org_id,
#             "total_files": len(files),
#             "files": files,
#             "api_version": "3.0.0"
#         }), 200
        
#     except Exception as e:
#         error_msg = str(e)
#         print(f"❌ List Files Error: {error_msg}")
#         traceback.print_exc()
#         return jsonify({"error": "Internal server error", "details": error_msg}), 500

# @app.route("/migrate-to-fileid", methods=["POST", "OPTIONS"])
# def migrate_to_file_id():
#     """Migrate existing uploadId-based documents to fileId-based system"""
#     if request.method == "OPTIONS":
#         return "", 200
        
#     if not FIREBASE_AVAILABLE:
#         return jsonify({"error": "Firebase is not available"}), 503
        
#     try:
#         data = request.get_json(silent=True) or {}
#         org_id = data.get("orgId")
#         file_mappings = data.get("fileMappings", [])
        
#         if not org_id:
#             return jsonify({"error": "orgId is required"}), 400
            
#         print(f"🔄 Starting enhanced migration for org: {org_id}")
        
#         files_ref = db.collection("document_embeddings").document(f"org-{org_id}").collection("files")
#         migration_results = []
        
#         for mapping in file_mappings:
#             upload_id = mapping.get("uploadId")
#             file_id = mapping.get("fileId")
            
#             if not upload_id or not file_id:
#                 continue
                
#             try:
#                 old_doc_ref = files_ref.document(upload_id)
#                 old_doc = old_doc_ref.get()
                
#                 if old_doc.exists:
#                     old_data = old_doc.to_dict()
                    
#                     new_doc_ref = files_ref.document(file_id)
                    
#                     new_data = old_data.copy()
#                     new_data["file_id"] = file_id
#                     new_data["processing_version"] = "3.0.0"
#                     new_data["migrated"] = True
                    
#                     new_doc_ref.set(new_data)
                    
#                     old_chunks = old_doc_ref.collection("chunks").stream()
#                     batch = db.batch()
#                     chunks_copied = 0
                    
#                     for chunk_doc in old_chunks:
#                         new_chunk_ref = new_doc_ref.collection("chunks").document(chunk_doc.id)
#                         batch.set(new_chunk_ref, chunk_doc.to_dict())
#                         chunks_copied += 1
                    
#                     batch.commit()
                    
#                     delete_collection(old_doc_ref.collection("chunks"), 100)
#                     old_doc_ref.delete()
                    
#                     # Update backend status
#                     update_backend_embedding_status(file_id, org_id, True)
                    
#                     migration_results.append({
#                         "upload_id": upload_id,
#                         "file_id": file_id,
#                         "filename": old_data.get("filename"),
#                         "chunks_migrated": chunks_copied,
#                         "processing_version": "3.0.0",
#                         "success": True
#                     })
                    
#                     print(f"✅ Migrated {upload_id} -> {file_id} ({chunks_copied} chunks)")
                    
#             except Exception as e:
#                 migration_results.append({
#                     "upload_id": upload_id,
#                     "file_id": file_id,
#                     "success": False,
#                     "error": str(e)
#                 })
#                 print(f"❌ Failed to migrate {upload_id} -> {file_id}: {str(e)}")
        
#         successful_migrations = len([r for r in migration_results if r["success"]])
        
#         return jsonify({
#             "message": f"Enhanced migration complete for org {org_id}",
#             "total_attempted": len(file_mappings),
#             "successful_migrations": successful_migrations,
#             "migration_details": migration_results,
#             "migration_version": "3.0.0"
#         }), 200
        
#     except Exception as e:
#         error_msg = str(e)
#         print(f"❌ Migration Error: {error_msg}")
#         traceback.print_exc()
#         return jsonify({"error": "Internal server error", "details": error_msg}), 500

@app.route("/chat", methods=["POST", "OPTIONS"])
def chat_with_doc():
    """Chat with documents using embeddings and AI - GLOBAL SEARCH (organization knowledge base only)"""
    if request.method == "OPTIONS":
        return "", 200
        
    if not FIREBASE_AVAILABLE:
        return jsonify({"error": "Firebase is not available"}), 503
        
    if not OPENAI_AVAILABLE:
        return jsonify({"error": "OpenAI embedding service is not available"}), 503
        
    if not VERTEX_AVAILABLE:
        return jsonify({"error": "AI generation service is not available"}), 503
        
    try:
        data = request.get_json(silent=True)
        if not data:
            return jsonify({"error": "No valid JSON input found"}), 400

        query = data.get("query")
        org_id = data.get("orgId")
        file_ids = data.get("fileIds", [])  # Optional: search only specific files
        conversation_history = data.get("conversationHistory", "")
        rerank = data.get("rerank", True)  # Enable reranking by default
        enable_hybrid_search = data.get("hybridSearch", True)  # Enable hybrid search by default
        dense_weight = data.get("denseWeight", 0.7)  # Weight for semantic search
        sparse_weight = data.get("sparseWeight", 0.3)  # Weight for BM25 search
        enable_query_expansion = data.get("queryExpansion", True)  # Enable query expansion by default
        max_query_variations = data.get("maxQueryVariations", 2)  # Number of query variations
        context_type = data.get("contextType", "rfp")  # Context for query expansion
        
        if not query or not org_id:
            return jsonify({"error": "Query and orgId are required."}), 400

        print(f"🔍 Processing GLOBAL chat query: '{query}' for org: {org_id}")
        print(f"🌐 Searching organization knowledge base documents")

        # Get query embedding
        query_embedding = np.array(embed_query(query))

        retrieved_docs = []
        search_sources = []

        # PART 1: Search organization-level documents (global knowledge base)
        print("📚 Searching organization-level documents...")
        org_files_ref = db.collection("document_embeddings").document(f"org-{org_id}").collection("files")
        
        if file_ids:
            # Search only specific files if provided
            org_files = []
            for file_id in file_ids:
                doc = org_files_ref.document(file_id).get()
                if doc.exists:
                    org_files.append(doc)
        else:
            # Search all organization files
            org_files = org_files_ref.stream()
        
        org_docs_found = 0
        # Process each organization file
        for file_doc in org_files:
            file_data = file_doc.to_dict()
            
            # Get chunks for this file
            chunks_ref = file_doc.reference.collection("chunks")
            chunks = chunks_ref.stream()
            
            # Process each chunk
            for chunk_doc in chunks:
                chunk_data = chunk_doc.to_dict()
                
                # Convert to numpy array
                chunk_embedding = np.array(chunk_data["embedding"])
                
                # Calculate cosine similarity
                score = np.dot(query_embedding, chunk_embedding) / (
                    np.linalg.norm(query_embedding) * np.linalg.norm(chunk_embedding)
                )
                
                if score >= 0.2:  # Similarity threshold
                    retrieved_docs.append({
                        "content": chunk_data.get("content", chunk_data.get("text", "No content available")), 
                        "score": float(score),
                        "filename": file_data.get("filename", "Unknown"),
                        "file_id": file_data.get("file_id", file_doc.id),
                        "file_type": file_data.get("file_type", "unknown"),
                        "processing_version": file_data.get("processing_version", "legacy"),
                        "document_source": "organization",
                        "source_type": "global"
                    })
                    org_docs_found += 1

        search_sources.append({
            "source": "organization",
            "type": "global", 
            "documents_found": org_docs_found,
            "collection": f"document_embeddings/org-{org_id}/files"
        })

        # Apply enhanced search (query expansion + hybrid search + reranking)
        if retrieved_docs:
            index_key = f"chat_{org_id}"
            
            if enable_query_expansion or enable_hybrid_search:
                print(f"🚀 Applying enhanced search to {len(retrieved_docs)} documents...")
                
                # Use enhanced search with query expansion and hybrid search
                enhanced_candidates = enhanced_search(
                    query=query,
                    documents=retrieved_docs,
                    index_key=index_key,
                    context_type=context_type,
                    enable_query_expansion=enable_query_expansion,
                    max_query_variations=max_query_variations,
                    search_method="hybrid" if enable_hybrid_search else "semantic",
                    top_k=20,
                    dense_weight=dense_weight,
                    sparse_weight=sparse_weight
                )
                
                if rerank and enhanced_candidates:
                    print(f"🔄 Applying reranking to {len(enhanced_candidates)} enhanced candidates...")
                    top_chunks = rerank_documents(query, enhanced_candidates, top_k=10)
                else:
                    top_chunks = enhanced_candidates[:10]
                    
            elif rerank:
                print(f"🔄 Applying reranking to {len(retrieved_docs)} documents...")
                # Get more candidates for reranking (top 20), then rerank to get top 10
                similarity_candidates = sorted(retrieved_docs, key=lambda x: x["score"], reverse=True)[:20]
                top_chunks = rerank_documents(query, similarity_candidates, top_k=10)
            else:
                # Get top chunks by similarity (from all sources combined)
                top_chunks = sorted(retrieved_docs, key=lambda x: x["score"], reverse=True)[:10]
        else:
            top_chunks = []

        if not top_chunks:
            # No relevant chunks found in organization documents
            search_summary = {
                "organization_docs": org_docs_found,
                "total_sources": len(search_sources)
            }
            
            no_results_message = f"I couldn't find any relevant information in your organization's knowledge base to answer your question."
            no_results_message += f" I searched {org_docs_found} organization documents."
            
            if org_docs_found == 0:
                no_results_message += " No documents were found in the system. Please upload relevant documents first."
            else:
                no_results_message += " The documents exist but don't contain information relevant to your specific query. Try rephrasing your question or check if the right documents are uploaded."
            
            return jsonify({
                "query": query, 
                "retrieved_chunks": [], 
                "answer": no_results_message,
                "source_files": [],
                "relevance_scores": [],
                "chat_type": "global",
                "search_sources": search_sources,
                "search_summary": search_summary,
                "api_version": "4.2.0"
            }), 200

        # Generate answer using combined context
        context_chunks = [doc["content"] for doc in top_chunks]
        answer = generate_answer_with_gcp(query, context_chunks, conversation_history)
        
        # Get unique source files from organization documents
        source_files = []
        seen_files = set()
        
        for doc in top_chunks:
            file_key = f"{doc['filename']}_{doc['file_type']}_{doc['document_source']}"
            if file_key not in seen_files:
                file_info = {
                    "filename": doc["filename"],
                    "file_type": doc["file_type"],
                    "file_id": doc["file_id"],
                    "processing_version": doc["processing_version"],
                    "document_source": doc["document_source"],
                    "source_type": doc["source_type"]
                }
                
                source_files.append(file_info)
                seen_files.add(file_key)

        # Create search summary
        search_summary = {
            "total_chunks_found": len(retrieved_docs),
            "top_chunks_used": len(top_chunks),
            "organization_files": len(source_files),
            "sources_searched": len(search_sources),
            "search_scope": "organization_knowledge_base"
        }

        print(f"✅ Global chat search completed:")
        print(f"  Organization documents: {org_docs_found} chunks from {len(source_files)} files")
        print(f"  Total relevant chunks: {len(retrieved_docs)}")
        print(f"  Used for answer: {len(top_chunks)}")

        return jsonify({
            "query": query, 
            "retrieved_chunks": context_chunks, 
            "answer": answer,
            "source_files": source_files,
            "relevance_scores": [doc["score"] for doc in top_chunks],
            "chat_type": "global",
            "search_sources": search_sources,
            "search_summary": search_summary,
            "api_version": "4.2.0"
        }), 200

    except Exception as e:
        error_msg = str(e)
        print(f"❌ Global Chat Error: {error_msg}")
        traceback.print_exc()
        return jsonify({"error": "Internal server error", "details": error_msg}), 500

@app.route("/chat-ragie", methods=["POST", "OPTIONS"])
def chat_with_ragie():
    """Chat endpoint using Ragie API for retrieval and local GPT for response generation"""
    if request.method == "OPTIONS":
        return "", 200
        
    if not VERTEX_AVAILABLE:
        return jsonify({"error": "AI generation service is not available"}), 503
        
    try:
        data = request.get_json(silent=True)
        if not data:
            return jsonify({"error": "No valid JSON input found"}), 400

        query = data.get("query")
        org_id = data.get("orgId")
        top_k = data.get("top_k", 10)
        rerank = data.get("rerank", False)
        
        if not query or not org_id:
            return jsonify({"error": "Query and orgId are required."}), 400

        print(f"🔍 Processing Ragie chat query: '{query}' for org: {org_id}")

        # Step 1: Call Ragie API for retrieval
        ragie_url = "https://api.ragie.ai/retrievals"
        ragie_payload = {
            "rerank": rerank,
            "query": query,
            "top_k": top_k
        }
        ragie_headers = {
            "accept": "application/json",
            "content-type": "application/json",
            "authorization": "Bearer tnt_7kuuKbBhbh4_RxOTauUTuIhEADMwd5w6BitGJbK41uSucAARlFqgeJi"
        }

        print("📡 Calling Ragie API for document retrieval...")
        ragie_response = requests.post(ragie_url, json=ragie_payload, headers=ragie_headers, timeout=30)
        
        if ragie_response.status_code != 200:
            print(f"❌ Ragie API error: {ragie_response.status_code} - {ragie_response.text}")
            return jsonify({"error": f"Ragie API error: {ragie_response.status_code}"}), 500

        ragie_data = ragie_response.json()
        scored_chunks = ragie_data.get("scored_chunks", [])
        
        print(f"✅ Ragie API returned {len(scored_chunks)} chunks")

        # Step 2: Prepare context for GPT generation using only Ragie results
        context_chunks = []
        
        # Add Ragie chunks to context
        for chunk in scored_chunks:
            context_chunks.append(f"[{chunk.get('document_name', 'Unknown')} (Score: {chunk.get('score', 0):.3f})]\n{chunk.get('text', '')}")

        if not context_chunks:
            return jsonify({
                "response": "I couldn't find relevant information to answer your query.",
                "query": query,
                "sources": {
                    "ragie_chunks": 0
                },
                "api_version": "1.0.0"
            }), 200

        print(f"🎯 Generating response with {len(context_chunks)} context chunks from Ragie")

        # Step 3: Generate response using local GPT/Vertex AI
        ai_response = generate_answer_with_gcp(query, context_chunks, "")

        # Step 4: Prepare response
        response_data = {
            "response": ai_response,
            "query": query,
            "sources": {
                "ragie_chunks": len(scored_chunks),
                "total_context_chunks": len(context_chunks)
            },
            "ragie_documents": [
                {
                    "text": chunk.get("text", "")[:200] + "...",
                    "score": chunk.get("score", 0),
                    "document_name": chunk.get("document_name", "Unknown"),
                    "document_id": chunk.get("document_id", "")
                } for chunk in scored_chunks[:5]  # Show top 5
            ],
            "api_version": "1.0.0"
        }

        print(f"✅ Chat response generated successfully")
        return jsonify(response_data), 200

    except requests.RequestException as req_error:
        print(f"❌ Request error: {str(req_error)}")
        return jsonify({"error": "External API request failed", "details": str(req_error)}), 500
    except Exception as e:
        error_msg = str(e)
        print(f"❌ Chat Ragie Error: {error_msg}")
        traceback.print_exc()
        return jsonify({"error": "Internal server error", "details": error_msg}), 500
    
# ================================
# RFP SUPPORT DOCUMENT ENDPOINTS
# ================================

@app.route("/upload-support-document", methods=["POST", "OPTIONS"])
def upload_support_document():
    """Upload and process a support document for a specific RFP (organized by orgId)"""
    if request.method == "OPTIONS":
        return "", 200
        
    if not FIREBASE_AVAILABLE:
        return jsonify({"error": "Firebase is not available"}), 503
        
    if not OPENAI_AVAILABLE:
        return jsonify({"error": "OpenAI embedding service is not available"}), 503
        
    save_path = None
    try:
        org_id = request.args.get("orgId")
        rfp_id = request.args.get("rfpId")
        file_id = request.args.get("fileId")
        upload_id = request.args.get("uploadId", str(uuid.uuid4()))
        
        print(f"📄 Starting RFP support document upload - Org: {org_id}, RFP: {rfp_id}, fileId: {file_id}, uploadId: {upload_id}")
        
        if not org_id:
            return jsonify({"error": "orgId is required"}), 400
            
        if not rfp_id:
            return jsonify({"error": "rfpId is required"}), 400
            
        if not file_id:
            return jsonify({"error": "fileId is required"}), 400

        if "file" not in request.files:
            return jsonify({"error": "No file provided."}), 400

        file = request.files["file"]
        filename = file.filename
        file_ext = filename.split(".")[-1].lower()

        # Get currently supported types
        supported, missing = check_file_processing_dependencies()
        all_supported = supported["always"] + supported["conditional"] + ["png", "jpg", "jpeg", "gif", "webp"]
        
        if file_ext not in all_supported:
            return jsonify({
                "error": f"Unsupported file type: {file_ext}. Supported: {', '.join(sorted(all_supported))}",
                "missing_dependencies": missing,
                "install_command": f"pip install {' '.join(missing)}" if missing else None
            }), 400

        # Initialize progress
        update_upload_progress(upload_id, "Processing", 0, "Starting RFP support document upload", filename)
        
        # Save file temporarily
        save_path = os.path.abspath(os.path.join(UPLOAD_FOLDER, f"{upload_id}_{filename}"))
        file.save(save_path)
        print(f"✅ File saved to {save_path}")
        
        # Process file in a separate thread
        def process_support_file_async():
            try:
                update_upload_progress(upload_id, "Processing", 25, "File saved, extracting content", filename)
                
                # Extract text and create chunks
                chunks = parse_and_chunk(save_path, file_ext, chunk_size=50, max_chunks=500)
                
                if not chunks:
                    print(f"❌ No content extracted from {filename}")
                    update_upload_progress(upload_id, "error", 0, "No content extracted", filename)
                    update_backend_embedding_status(file_id, org_id, False)
                    return
                    
                # Generate embeddings
                update_upload_progress(upload_id, "Processing", 50, "Generating embeddings", filename)
                embeddings = embed_chunks(chunks, upload_id=upload_id, org_id=org_id, filename=filename)
                
                # Store in NEW RFP-specific Firestore collection structure
                update_upload_progress(upload_id, "Processing", 75, "Creating document structure", filename)
                
                # STEP 1: Create/update the ORG parent document
                print(f"🏗️ Creating org parent document: org-{org_id}")
                org_doc_ref = db.collection("org_rfp_support_embeddings").document(f"org-{org_id}")
                org_doc_ref.set({
                    "org_id": org_id,
                    "created_at": firestore.SERVER_TIMESTAMP,
                    "last_updated": firestore.SERVER_TIMESTAMP,
                    "rfp_count": firestore.Increment(1)  # Increment RFP count
                }, merge=True)  # Use merge=True to not overwrite existing data
                
                print(f"✅ Created/updated org parent document: org-{org_id}")
                
                # STEP 2: Create/update the RFP parent document  
                print(f"🏗️ Creating RFP parent document: rfp-{rfp_id}")
                rfp_parent_ref = org_doc_ref.collection("rfps").document(f"rfp-{rfp_id}")
                rfp_parent_ref.set({
                    "rfp_id": rfp_id,
                    "org_id": org_id,
                    "created_at": firestore.SERVER_TIMESTAMP,
                    "last_updated": firestore.SERVER_TIMESTAMP,
                    "file_count": firestore.Increment(1)  # Increment file count
                }, merge=True)  # Use merge=True to not overwrite existing data
                
                print(f"✅ Created/updated RFP parent document: rfp-{rfp_id}")
                
                # STEP 3: Create the file document
                update_upload_progress(upload_id, "Processing", 80, "Storing RFP support document embeddings", filename)
                
                rfp_doc_ref = rfp_parent_ref.collection("files").document(file_id)
                
                print(f"🗂️ Storing document at: org_rfp_support_embeddings/org-{org_id}/rfps/rfp-{rfp_id}/files/{file_id}")
                
                # Store file metadata
                rfp_doc_ref.set({
                    "filename": filename,
                    "file_id": file_id,
                    "upload_id": upload_id,
                    "rfp_id": rfp_id,
                    "org_id": org_id,
                    "file_type": file_ext,
                    "document_type": "support",
                    "created_at": firestore.SERVER_TIMESTAMP,
                    "chunk_count": len(chunks),
                    "processing_version": "4.2.0"  # Updated version for fixed structure
                })
                
                print(f"✅ Created file document: {file_id}")
                
                # STEP 4: Store chunks in batches - OPTIMIZED for 50x faster writes
                update_upload_progress(upload_id, "Processing", 85, "Storing embeddings", filename)
                
                batch_size = 500  # Firestore allows up to 500 writes per batch
                total_chunks = len(chunks)
                
                print(f"🚀 Storing {total_chunks} RFP chunks in batches of {batch_size}")
                
                for i in range(0, total_chunks, batch_size):
                    batch = db.batch()
                    end_idx = min(i + batch_size, total_chunks)
                    
                    for j in range(i, end_idx):
                        chunk_ref = rfp_doc_ref.collection("chunks").document(str(j))
                        batch.set(chunk_ref, {
                            "content": chunks[j],
                            "embedding": embeddings[j],
                            "index": j
                        })
                    
                    batch.commit()
                    print(f"✓ Stored RFP batch {i//batch_size + 1}: {end_idx}/{total_chunks} chunks")
                    
                    # Only clean up memory for large batches
                    if end_idx - i >= 100:
                        del batch
                        import gc
                        gc.collect()
                    
                    progress = 85 + ((end_idx / total_chunks) * 15)
                    update_upload_progress(upload_id, "Processing", progress, 
                                        f"Storing chunks ({end_idx}/{total_chunks})", filename)
                
                print(f"✅ Successfully processed RFP support document {filename} for Org: {org_id}, RFP: {rfp_id}")
                print(f"📊 Created complete structure: org-{org_id}/rfps/rfp-{rfp_id}/files/{file_id} with {len(chunks)} chunks")
                
                # STEP 5: Verify the structure was created correctly
                try:
                    # Verify we can now list RFPs
                    verification_rfps = list(org_doc_ref.collection("rfps").stream())
                    print(f"✅ Verification: Can list {len(verification_rfps)} RFPs for org")
                    
                    # Verify RFP document exists
                    rfp_exists = rfp_parent_ref.get().exists
                    print(f"✅ Verification: RFP document exists: {rfp_exists}")
                    
                    # Verify file document exists
                    file_exists = rfp_doc_ref.get().exists
                    print(f"✅ Verification: File document exists: {file_exists}")
                    
                except Exception as verification_error:
                    print(f"⚠️ Verification failed (but upload succeeded): {str(verification_error)}")
                
                update_upload_progress(upload_id, "Completed", 100, "RFP support document processing completed", filename)
                
                # Update backend to mark file as having embeddings
                update_backend_embedding_status(file_id, org_id, True)
                
            except Exception as e:
                error_msg = str(e)
                print(f"❌ Async processing error: {error_msg}")
                traceback.print_exc()
                update_upload_progress(upload_id, "error", 0, f"Error: {error_msg}", filename)
                update_backend_embedding_status(file_id, org_id, False)
                
            finally:
                if save_path and os.path.exists(save_path):
                    try:
                        os.remove(save_path)
                        print(f"🧹 File {save_path} deleted.")
                    except Exception as e:
                        print(f"Error deleting file: {e}")
        
        # Start processing in background thread
        processing_thread = Thread(target=process_support_file_async)
        processing_thread.daemon = False
        processing_thread.start()
        
        return jsonify({
            "message": "RFP support document upload started. Check status with the upload-status endpoint.",
            "file_id": file_id,
            "upload_id": upload_id,
            "rfp_id": rfp_id,
            "org_id": org_id,
            "file_type": file_ext,
            "document_type": "support",
            "storage_structure": f"org_rfp_support_embeddings/org-{org_id}/rfps/rfp-{rfp_id}/files/{file_id}",
            "processing_version": "4.2.0",
            "structure_fix": "Creates parent documents for proper subcollection listing"
        }), 202

    except Exception as e:
        error_msg = str(e)
        print(f"❌ Support Document Upload Error: {error_msg}")
        traceback.print_exc()
        
        if 'upload_id' in locals() and upload_id:
            update_upload_progress(upload_id, "error", 0, f"Error: {error_msg}", filename if 'filename' in locals() else "")
            
        if save_path and os.path.exists(save_path):
            try:
                os.remove(save_path)
                print(f"🧹 File {save_path} deleted after error.")
            except Exception as cleanup_error:
                print(f"Error cleaning up file: {cleanup_error}")
                
        return jsonify({"error": "Internal server error", "details": error_msg}), 500
    
@app.route("/check-rfp-support-documents", methods=["GET", "OPTIONS"])
def check_rfp_support_documents():
    """Check if an RFP has support documents (NEW ORG-BASED STRUCTURE)"""
    if request.method == "OPTIONS":
        return "", 200
        
    if not FIREBASE_AVAILABLE:
        return jsonify({"error": "Firebase is not available"}), 503
        
    try:
        org_id = request.args.get("orgId")  # NEW: Required parameter
        rfp_id = request.args.get("rfpId")
        
        
        if not rfp_id:
            return jsonify({"error": "rfpId is required"}), 400
            
        if not org_id:
            return jsonify({"error": "orgId is required"}), 400
        
        print(f"🔍 Checking RFP support documents for Org: {org_id}, RFP: {rfp_id}")
        
        # NEW STRUCTURE: Check if RFP support documents exist
        rfp_files_ref = (db.collection("org_rfp_support_embeddings")
                        .document(f"org-{org_id}")
                        .collection("rfps")
                        .document(f"rfp-{rfp_id}")
                        .collection("files"))
        
        files = list(rfp_files_ref.limit(1).stream())
        has_support_documents = len(files) > 0
        
        if has_support_documents:
            # Get all files for this RFP
            all_files = rfp_files_ref.stream()
            file_list = []
            
            for file_doc in all_files:
                file_data = file_doc.to_dict()
                file_list.append({
                    "file_id": file_data.get("file_id"),
                    "filename": file_data.get("filename"),
                    "file_type": file_data.get("file_type"),
                    "created_at": file_data.get("created_at"),
                    "chunk_count": file_data.get("chunk_count", 0),
                    "processing_version": file_data.get("processing_version", "legacy")
                })
        else:
            file_list = []
        
        return jsonify({
            "rfp_id": rfp_id,
            "org_id": org_id,
            "has_support_documents": has_support_documents,
            "document_count": len(file_list),
            "files": file_list,
            "storage_structure": f"org_rfp_support_embeddings/org-{org_id}/rfps/rfp-{rfp_id}/files"
        }), 200
        
    except Exception as e:
        error_msg = str(e)
        print(f"❌ Check RFP Support Documents Error: {error_msg}")
        return jsonify({"error": "Internal server error", "details": error_msg}), 500
    
@app.route("/delete-rfp-support-document", methods=["DELETE", "OPTIONS"])
def delete_rfp_support_document():
    """Delete a specific RFP support document (NEW ORG-BASED STRUCTURE)"""
    if request.method == "OPTIONS":
        return "", 200
        
    if not FIREBASE_AVAILABLE:
        return jsonify({"error": "Firebase is not available"}), 503
        
    try:
        rfp_id = request.args.get("rfpId")
        file_id = request.args.get("fileId")
        org_id = request.args.get("orgId")  # NEW: Required parameter

        if not rfp_id or not file_id or not org_id:
            return jsonify({"error": "rfpId, fileId, and orgId are required"}), 400

        print(f"🗑️ Deleting RFP support document - Org: {org_id}, RFP: {rfp_id}, fileId: {file_id}")

        # NEW STRUCTURE: Access document using org-based path
        rfp_file_ref = (db.collection("org_rfp_support_embeddings")
                       .document(f"org-{org_id}")
                       .collection("rfps")
                       .document(f"rfp-{rfp_id}")
                       .collection("files")
                       .document(file_id))
        
        file_doc = rfp_file_ref.get()
        
        if file_doc.exists:
            file_data = file_doc.to_dict()
            
            # Delete all chunks
            chunks_deleted = delete_collection(rfp_file_ref.collection("chunks"), 100)
            
            # Delete the file document
            rfp_file_ref.delete()
            
            print(f"✅ Successfully deleted RFP support document and {chunks_deleted} chunks")
            
            # Update backend status
            update_backend_embedding_status(file_id, org_id, False)
            
            return jsonify({
                "message": "Successfully deleted RFP support document",
                "rfp_id": rfp_id,
                "org_id": org_id,
                "file_id": file_id,
                "filename": file_data.get("filename", "unknown"),
                "chunks_deleted": chunks_deleted,
                "storage_structure": f"org_rfp_support_embeddings/org-{org_id}/rfps/rfp-{rfp_id}/files/{file_id}"
            }), 200
        else:
            return jsonify({
                "message": "RFP support document not found",
                "rfp_id": rfp_id,
                "org_id": org_id,
                "file_id": file_id,
                "storage_structure": f"org_rfp_support_embeddings/org-{org_id}/rfps/rfp-{rfp_id}/files/{file_id}"
            }), 404

    except Exception as e:
        error_msg = str(e)
        print(f"❌ Delete RFP Support Document Error: {error_msg}")
        return jsonify({"error": "Internal server error", "details": error_msg}), 500

@app.route("/delete-project-data", methods=["DELETE", "OPTIONS"])
def delete_project_data():
    """Delete all project support documents and embeddings from Firestore"""
    if request.method == "OPTIONS":
        return "", 200
        
    if not FIREBASE_AVAILABLE:
        return jsonify({"error": "Firebase is not available"}), 503
        
    try:
        project_id = request.args.get("projectId")
        org_id = request.args.get("orgId")

        if not project_id or not org_id:
            return jsonify({"error": "projectId and orgId are required"}), 400

        print(f"🗑️ Deleting project data - Org: {org_id}, Project: {project_id}")

        # Access the project document in the org-based structure
        project_ref = (db.collection("org_project_support_embeddings")
                      .document(f"org-{org_id}")
                      .collection("projects")
                      .document(f"project-{project_id}"))
        
        project_doc = project_ref.get()
        total_files_deleted = 0
        total_chunks_deleted = 0
        
        if project_doc.exists:
            project_data = project_doc.to_dict()
            
            # Delete all files and their chunks
            files_collection = project_ref.collection("files")
            files = files_collection.stream()
            
            for file_doc in files:
                file_data = file_doc.to_dict()
                file_id = file_doc.id
                
                print(f"🗑️ Deleting project file: {file_data.get('filename', 'unknown')} (ID: {file_id})")
                
                # Delete all chunks for this file
                chunks_deleted = delete_collection(file_doc.reference.collection("chunks"), 100)
                total_chunks_deleted += chunks_deleted
                
                # Delete the file document
                file_doc.reference.delete()
                total_files_deleted += 1
            
            # Delete the project document itself
            project_ref.delete()
            
            print(f"✅ Successfully deleted project data: {total_files_deleted} files, {total_chunks_deleted} chunks")
            
            return jsonify({
                "message": "Successfully deleted project data",
                "project_id": project_id,
                "org_id": org_id,
                "files_deleted": total_files_deleted,
                "chunks_deleted": total_chunks_deleted,
                "storage_structure": f"org_project_support_embeddings/org-{org_id}/projects/project-{project_id}"
            }), 200
        else:
            return jsonify({
                "message": "Project data not found",
                "project_id": project_id,
                "org_id": org_id,
                "files_deleted": 0,
                "chunks_deleted": 0,
                "storage_structure": f"org_project_support_embeddings/org-{org_id}/projects/project-{project_id}"
            }), 404

    except Exception as e:
        error_msg = str(e)
        print(f"❌ Delete Project Data Error: {error_msg}")
        return jsonify({"error": "Internal server error", "details": error_msg}), 500

@app.route("/delete-project-support-document", methods=["DELETE", "OPTIONS"])
def delete_project_support_document():
    """Delete a specific project support document from Firestore"""
    if request.method == "OPTIONS":
        return "", 200
        
    if not FIREBASE_AVAILABLE:
        return jsonify({"error": "Firebase is not available"}), 503
        
    try:
        project_id = request.args.get("projectId")
        file_id = request.args.get("fileId")
        org_id = request.args.get("orgId")

        if not project_id or not file_id or not org_id:
            return jsonify({"error": "projectId, fileId, and orgId are required"}), 400

        print(f"🗑️ Deleting project support document - Org: {org_id}, Project: {project_id}, File: {file_id}")

        # Access the file document in the project-based structure
        project_file_ref = (db.collection("org_project_support_embeddings")
                           .document(f"org-{org_id}")
                           .collection("projects")
                           .document(f"project-{project_id}")
                           .collection("files")
                           .document(file_id))
        
        file_doc = project_file_ref.get()
        chunks_deleted = 0
        
        if file_doc.exists:
            file_data = file_doc.to_dict()
            
            print(f"🗑️ Deleting project support file: {file_data.get('filename', 'unknown')} (ID: {file_id})")
            
            # Delete all chunks for this file
            chunks_deleted = delete_collection(file_doc.reference.collection("chunks"), 100)
            
            # Delete the file document
            file_doc.reference.delete()
            
            print(f"✅ Successfully deleted project support document and {chunks_deleted} chunks")
            
            return jsonify({
                "message": "Successfully deleted project support document",
                "project_id": project_id,
                "org_id": org_id,
                "file_id": file_id,
                "filename": file_data.get("filename", "unknown"),
                "chunks_deleted": chunks_deleted,
                "storage_structure": f"org_project_support_embeddings/org-{org_id}/projects/project-{project_id}/files/{file_id}"
            }), 200
        else:
            return jsonify({
                "message": "Project support document not found",
                "project_id": project_id,
                "org_id": org_id,
                "file_id": file_id,
                "chunks_deleted": 0,
                "storage_structure": f"org_project_support_embeddings/org-{org_id}/projects/project-{project_id}/files/{file_id}"
            }), 404

    except Exception as e:
        error_msg = str(e)
        print(f"❌ Delete Project Support Document Error: {error_msg}")
        return jsonify({"error": "Internal server error", "details": error_msg}), 500

@app.route("/run-agent-v2", methods=["POST", "OPTIONS"])
def run_agent_v2():
    """V2 Agent runner with WebSocket support and database storage"""
    if request.method == "OPTIONS":
        return "", 200
        
    if not VERTEX_AVAILABLE:
        return jsonify({"error": "AI generation service is not available"}), 503
        
    try:
        data = request.get_json(silent=True) or {}
        agent_run_id = data.get("agentRunId")
        agent_id = data.get("agentId")
        agent_type = data.get("agentType")
        org_id = data.get("orgId")
        user_id = data.get("userId")
        project_id = data.get("projectId")
        project_file = data.get("projectFile", {})
        
        print(f"🤖 Starting V2 Agent: {agent_id} (Type: {agent_type}), Run ID: {agent_run_id}")
        
        if not agent_run_id or not agent_id or not org_id or not user_id:
            return jsonify({"error": "agentRunId, agentId, orgId, and userId are required"}), 400
            
        if not project_file or not project_file.get("gcpUrl"):
            return jsonify({"error": "projectFile with gcpUrl is required"}), 400

        # Send initial status update
        notify_agent_status(agent_run_id, user_id, "running", False)
        
        def process_agent_v2_async():
            start_time = time.time()
            try:
                # Download file from GCP
                notify_agent_status(agent_run_id, user_id, "downloading", False)
                
                file_content = download_file_from_gcp(project_file["gcpUrl"])
                if not file_content:
                    raise Exception("Failed to download file from GCP")
                
                filename = project_file["filename"]
                file_type = project_file["fileType"]
                
                # Save file temporarily
                temp_file_path = os.path.join(UPLOAD_FOLDER, f"{agent_run_id}_{filename}")
                with open(temp_file_path, 'wb') as f:
                    f.write(file_content)
                
                notify_agent_status(agent_run_id, user_id, "processing", False)
                
                # Process based on agent type
                if agent_type == "question_generator":
                    result = process_question_generator_v2(temp_file_path, filename, agent_run_id, user_id)
                else:
                    raise Exception(f"Unsupported agent type: {agent_type}")
                
                # Calculate duration
                end_time = time.time()
                duration = int(end_time - start_time)
                
                # Send results to backend
                notify_agent_completion(agent_run_id, user_id, result, duration)
                
                print(f"✅ V2 Agent completed successfully: {agent_run_id}")
                
            except Exception as e:
                error_msg = str(e)
                print(f"❌ V2 Agent Error: {error_msg}")
                traceback.print_exc()
                
                # Calculate duration even for failures
                end_time = time.time()
                duration = int(end_time - start_time)
                
                # Send error status
                notify_agent_error(agent_run_id, user_id, error_msg, duration)
                
            finally:
                # Clean up temp file
                if 'temp_file_path' in locals() and os.path.exists(temp_file_path):
                    try:
                        os.remove(temp_file_path)
                    except:
                        pass
        
        # Start processing in background
        agent_thread = Thread(target=process_agent_v2_async)
        agent_thread.daemon = False
        agent_thread.start()
        
        return jsonify({
            "message": "V2 Agent started successfully",
            "agent_run_id": agent_run_id,
            "agent_id": agent_id,
            "agent_type": agent_type,
            "status": "running"
        }), 202

    except Exception as e:
        error_msg = str(e)
        print(f"❌ V2 Agent startup error: {error_msg}")
        return jsonify({"error": "Internal server error", "details": error_msg}), 500

def download_file_from_gcp(gcp_url):
    """Download file from GCP URL"""
    try:
        # Extract bucket and path from gs:// URL
        if not gcp_url.startswith("gs://"):
            raise Exception("Invalid GCP URL format")
            
        # Remove gs:// prefix
        path = gcp_url[5:]
        bucket_name = path.split('/')[0]
        file_path = '/'.join(path.split('/')[1:])
        
        print(f"📥 Downloading from GCP: {bucket_name}/{file_path}")
        
        # Download using Google Cloud Storage client with credentials
        from google.cloud import storage
        
        # For GCS buckets, use environment credentials (care-proposals-451406 project)
        cred_path = os.environ.get("GOOGLE_APPLICATION_CREDENTIALS", "fire.json")
        if cred_path.startswith("{"):
            # JSON string from environment - this has bucket access
            import json
            from google.oauth2 import service_account
            cred_dict = json.loads(cred_path)
            bucket_credentials = service_account.Credentials.from_service_account_info(cred_dict)
            client = storage.Client(credentials=bucket_credentials, project=os.environ.get("GOOGLE_CLOUD_PROJECT"))
            print("✅ Using environment credentials for GCS bucket access")
        else:
            # Fallback to file path or default credentials
            client = storage.Client()
            print("✅ Using default/file credentials for GCS bucket access")
        bucket = client.bucket(bucket_name)
        blob = bucket.blob(file_path)
        
        return blob.download_as_bytes()
        
    except Exception as e:
        print(f"❌ GCP download error: {str(e)}")
        return None

def process_question_generator_v2(file_path, filename, agent_run_id, user_id):
    """Process question generation for V2 agent"""
    try:
        notify_agent_status(agent_run_id, user_id, "extracting", False)
        
        # Extract questions using existing AI function
        questions = extract_questions_with_ai_direct(file_path, filename)
        
        notify_agent_status(agent_run_id, user_id, "generating", False)
        
        if questions:
            # Handle both old and new sectioned formats
            if isinstance(questions, dict) and questions.get("format") == "sectioned":
                # New sectioned format
                total_questions = questions.get("total_questions", 0)
                sections_count = questions.get("total_sections", 0)
                
                result = {
                    "agent_type": "question_generator",
                    "format": "sectioned",
                    "sections": questions.get("sections", []),
                    "total_sections": sections_count,
                    "total_questions": total_questions,
                    "file_processed": {
                        "filename": filename,
                        "questions_extracted": total_questions,
                        "sections_created": sections_count
                    },
                    "processing_version": "v2.1.0",
                    "timestamp": datetime.now().isoformat()
                }
                
                print(f"✅ Generated {sections_count} sections with {total_questions} total questions from {filename}")
                return result
            elif isinstance(questions, list) and len(questions) > 0:
                # Legacy format - convert to sectioned format
                sectioned_result = {
                    "format": "sectioned",
                    "sections": [{
                        "section_id": "legacy_questions",
                        "section_title": "Document Questions",
                        "section_description": "Questions extracted from the document",
                        "questions": questions[:15]
                    }],
                    "total_sections": 1,
                    "total_questions": len(questions)
                }
                
                result = {
                    "agent_type": "question_generator",
                    "format": "sectioned",
                    "sections": sectioned_result["sections"],
                    "total_sections": 1,
                    "total_questions": len(questions),
                    "file_processed": {
                        "filename": filename,
                        "questions_extracted": len(questions),
                        "sections_created": 1
                    },
                    "processing_version": "v2.1.0",
                    "timestamp": datetime.now().isoformat()
                }
                
                print(f"✅ Generated 1 section with {len(questions)} questions from {filename} (legacy format converted)")
                return result
            else:
                raise Exception("Invalid questions format returned from extraction")
        else:
            raise Exception("No questions could be extracted from the file")
            
    except Exception as e:
        print(f"❌ Question generation error: {str(e)}")
        raise e

def notify_agent_status(agent_run_id, user_id, status, is_complete, error=None):
    """Send agent status update to backend via webhook"""
    try:
        backend_api_url = os.environ.get("BACKEND_API_URL", "http://localhost:8080")
        webhook_url = f"{backend_api_url}/api/v2/agents/webhook/status"
        
        payload = {
            "agentRunId": agent_run_id,
            "userId": user_id,
            "status": status,
            "isComplete": is_complete,
            "timestamp": datetime.now().isoformat(),
            "source": "flask_ai_v2"
        }
        
        if error:
            payload["error"] = error
            
        response = requests.post(
            webhook_url,
            json=payload,
            headers={"Content-Type": "application/json"},
            timeout=10
        )
        
        if response.status_code == 200:
            print(f"✅ Agent status webhook sent: {agent_run_id} → {status}")
        else:
            print(f"⚠️ Agent webhook failed: {response.status_code}")
            
    except Exception as e:
        print(f"❌ Agent webhook error: {str(e)}")

def notify_agent_completion(agent_run_id, user_id, result, duration):
    """Send agent completion notification with results"""
    try:
        backend_api_url = os.environ.get("BACKEND_API_URL", "http://localhost:8080")
        webhook_url = f"{backend_api_url}/api/v2/agents/webhook/completion"
        
        payload = {
            "agentRunId": agent_run_id,
            "userId": user_id,
            "status": "completed",
            "isComplete": True,
            "duration": duration,
            "result": result,
            "timestamp": datetime.now().isoformat(),
            "source": "flask_ai_v2"
        }
            
        response = requests.post(
            webhook_url,
            json=payload,
            headers={"Content-Type": "application/json"},
            timeout=30
        )
        
        if response.status_code == 200:
            print(f"✅ Agent completion webhook sent: {agent_run_id}")
        else:
            print(f"⚠️ Agent completion webhook failed: {response.status_code}")
            
    except Exception as e:
        print(f"❌ Agent completion webhook error: {str(e)}")

def notify_agent_error(agent_run_id, user_id, error_message, duration):
    """Send agent error notification"""
    try:
        backend_api_url = os.environ.get("BACKEND_API_URL", "http://localhost:8080")
        webhook_url = f"{backend_api_url}/api/v2/agents/webhook/error"
        
        payload = {
            "agentRunId": agent_run_id,
            "userId": user_id,
            "status": "failed",
            "isComplete": True,
            "duration": duration,
            "error": error_message,
            "timestamp": datetime.now().isoformat(),
            "source": "flask_ai_v2"
        }
            
        response = requests.post(
            webhook_url,
            json=payload,
            headers={"Content-Type": "application/json"},
            timeout=10
        )
        
        if response.status_code == 200:
            print(f"✅ Agent error webhook sent: {agent_run_id}")
        else:
            print(f"⚠️ Agent error webhook failed: {response.status_code}")
            
    except Exception as e:
        print(f"❌ Agent error webhook error: {str(e)}")

# Add this new endpoint to your Python Flask application

@app.route("/run-question-agent", methods=["POST", "OPTIONS"])
def run_question_agent():
    """Run Question Agent - Fetch project files and generate questions"""
    if request.method == "OPTIONS":
        return "", 200
        
    if not VERTEX_AVAILABLE:
        return jsonify({"error": "AI generation service is not available"}), 503
        
    try:
        data = request.get_json(silent=True) or {}
        rfp_id = data.get("rfpId")
        org_id = data.get("orgId")
        agent_run_id = data.get("agentRunId", str(uuid.uuid4()))
        auth_token = data.get("authToken") 
        
        print(f"🤖 Starting Question Agent for RFP: {rfp_id}, Agent Run ID: {agent_run_id}")
        
        if not rfp_id:
            return jsonify({"error": "rfpId is required"}), 400
            
        if not org_id:
            return jsonify({"error": "orgId is required"}), 400

        # Initialize progress tracking for agent
        update_upload_progress(agent_run_id, "Processing", 0, "Starting Question Agent", "")
        
        def run_question_agent_async():
            try:
                # Step 1: Fetch project files from backend
                update_upload_progress(agent_run_id, "Processing", 10, "Fetching project files", "")
                
                backend_api_url = os.environ.get("BACKEND_API_URL", "http://localhost:8080")
                project_files_url = f"{backend_api_url}/api/rfps/project-files/{rfp_id}"
                headers = {"Content-Type": "application/json"}
                if auth_token:
                    headers["Authorization"] = f"Bearer {auth_token}"
                
                try:
                    print(f"🔍 Fetching project files from: {project_files_url}")
                    files_response = requests.get(project_files_url, headers=headers, timeout=30)
                    
                    if files_response.status_code != 200:
                        error_text = files_response.text[:300] if files_response.text else "No response"
                        print(f"❌ Failed to fetch project files: {files_response.status_code}")
                        print(f"Response: {error_text}")
                        raise Exception(f"Failed to fetch project files: {files_response.status_code} - {error_text}")
                    
                    files_data = files_response.json()
                    project_files = files_data.get("files", [])
                    
                    if not project_files:
                        print("⚠️ No project files found")
                        update_upload_progress(agent_run_id, "Completed", 100, "No project files found to process", "")
                        
                        # Store empty results
                        with progress_lock:
                            if agent_run_id in upload_progress:
                                upload_progress[agent_run_id]["agent_results"] = {
                                    "agent_run_id": agent_run_id,
                                    "agent_type": "question_generator",
                                    "rfp_id": rfp_id,
                                    "org_id": org_id,
                                    "total_files": 0,
                                    "successful_files": 0,
                                    "failed_files": 0,
                                    "total_questions": 0,
                                    "questions": [],
                                    "processed_files": [],
                                    "failed_files": [],
                                    "timestamp": datetime.now().isoformat(),
                                    "status": "completed"
                                }
                        return
                    
                    print(f"📁 Found {len(project_files)} project files to process")
                    update_upload_progress(agent_run_id, "Processing", 15, f"Found {len(project_files)} project files", "")
                    
                except Exception as fetch_error:
                    raise Exception(f"Error fetching project files: {str(fetch_error)}")
                
                # Step 2: Download and process each file using ZIP method
                all_questions = []
                processed_files = []
                failed_files = []
                
                total_files = len(project_files)
                
                for index, file_info in enumerate(project_files):
                    file_id = file_info.get("id")
                    filename = file_info.get("name") or file_info.get("fileName", f"file_{file_id}")
                    file_url = file_info.get("url")
                    file_type = file_info.get("fileType", "")
                    file_ext = filename.split(".")[-1].lower() if "." in filename else ""
                    
                    try:
                        # Update progress
                        progress = 15 + ((index / total_files) * 65)  # 15% to 80%
                        update_upload_progress(agent_run_id, "Processing", progress, 
                                            f"Processing file {index + 1}/{total_files}: {filename}")
                        
                        print(f"\n📄 Processing file {index + 1}/{total_files}:")
                        print(f"  ID: {file_id}")
                        print(f"  Name: {filename}")
                        print(f"  Type: {file_type}")
                        print(f"  Extension: {file_ext}")
                        
                        # Validate file info
                        if not file_id:
                            raise Exception("File ID is missing")
                        
                        if not filename or filename.strip() == "":
                            filename = f"file_{file_id}"
                            print(f"⚠️ Using generated filename: {filename}")
                        
                        # Check if file type is supported
                        supported, missing = check_file_processing_dependencies()
                        all_supported = supported["always"] + supported["conditional"] + ["png", "jpg", "jpeg", "gif", "webp"]
                        
                        if file_ext and file_ext not in all_supported:
                            error_msg = f"Unsupported file type: {file_ext}. Supported: {', '.join(sorted(all_supported))}"
                            print(f"⚠️ {error_msg}")
                            failed_files.append({
                                "filename": filename,
                                "file_id": file_id,
                                "error": error_msg
                            })
                            continue
                        
                        # Download file using ZIP method
                        update_upload_progress(agent_run_id, "Processing", progress + 1, 
                                            f"Downloading {filename} as ZIP")
                        
                        try:
                            file_content = download_single_file_as_zip(
                                backend_api_url, file_id, filename, headers
                            )
                            
                            if not file_content or len(file_content) == 0:
                                raise Exception("Downloaded file content is empty")
                                
                            print(f"✅ Successfully downloaded {filename} ({len(file_content)} bytes)")
                            
                        except Exception as download_error:
                            error_msg = f"Download failed: {str(download_error)}"
                            print(f"❌ {error_msg}")
                            failed_files.append({
                                "filename": filename,
                                "file_id": file_id,
                                "error": error_msg
                            })
                            continue
                        
                        # Save file temporarily
                        update_upload_progress(agent_run_id, "Processing", progress + 2, 
                                            f"Saving {filename} temporarily")
                        
                        temp_file_path = os.path.join(UPLOAD_FOLDER, f"{agent_run_id}_{index}_{filename}")
                        
                        try:
                            with open(temp_file_path, 'wb') as temp_file:
                                temp_file.write(file_content)
                            
                            # Verify saved file
                            saved_size = os.path.getsize(temp_file_path)
                            if saved_size != len(file_content):
                                raise Exception(f"File save verification failed: expected {len(file_content)}, got {saved_size}")
                                
                            print(f"💾 File saved temporarily: {temp_file_path} (verified: {saved_size} bytes)")
                            
                        except Exception as save_error:
                            error_msg = f"Failed to save file: {str(save_error)}"
                            print(f"❌ {error_msg}")
                            failed_files.append({
                                "filename": filename,
                                "file_id": file_id,
                                "error": error_msg
                            })
                            continue
                        
                        # Extract questions using AI
                        update_upload_progress(agent_run_id, "Processing", progress + 3, 
                                            f"AI analyzing {filename}")
                        
                        try:
                            questions = extract_questions_with_ai_direct(temp_file_path, filename)
                            
                            if questions and len(questions) > 0:
                                processed_files.append({
                                    "filename": filename,
                                    "file_id": file_id,
                                    "file_type": file_ext,
                                    "questions": questions,
                                    "question_count": len(questions)
                                })
                                all_questions.extend(questions)
                                print(f"✅ Extracted {len(questions)} questions from {filename}")
                            else:
                                error_msg = "AI could not extract questions from this file"
                                print(f"⚠️ {error_msg}")
                                failed_files.append({
                                    "filename": filename,
                                    "file_id": file_id,
                                    "error": error_msg
                                })
                                
                        except Exception as ai_error:
                            error_msg = f"AI processing failed: {str(ai_error)}"
                            print(f"❌ {error_msg}")
                            failed_files.append({
                                "filename": filename,
                                "file_id": file_id,
                                "error": error_msg
                            })
                        
                        # Clean up temporary file
                        try:
                            os.remove(temp_file_path)
                            print(f"🧹 Cleaned up temp file: {temp_file_path}")
                        except Exception as cleanup_error:
                            print(f"⚠️ Failed to remove temp file {temp_file_path}: {cleanup_error}")
                            
                    except Exception as file_error:
                        error_msg = str(file_error)
                        print(f"❌ Error processing file {filename}: {error_msg}")
                        traceback.print_exc()
                        failed_files.append({
                            "filename": filename,
                            "file_id": file_id,
                            "error": error_msg
                        })
                        
                        # Ensure temp file cleanup even on error
                        temp_file_path = os.path.join(UPLOAD_FOLDER, f"{agent_run_id}_{index}_{filename}")
                        if os.path.exists(temp_file_path):
                            try:
                                os.remove(temp_file_path)
                                print(f"🧹 Cleaned up temp file after error: {temp_file_path}")
                            except:
                                pass
                
                # Step 3: Compile and store results
                update_upload_progress(agent_run_id, "Processing", 85, "Compiling agent results")
                
                print(f"\n📊 Question Agent Results Summary:")
                print(f"  Total files: {total_files}")
                print(f"  Successfully processed: {len(processed_files)}")
                print(f"  Failed: {len(failed_files)}")
                print(f"  Total questions generated: {len(all_questions)}")
                
                # Store agent results
                agent_results = {
                    "agent_run_id": agent_run_id,
                    "agent_type": "question_generator",
                    "rfp_id": rfp_id,
                    "org_id": org_id,
                    "total_files": total_files,
                    "successful_files": len(processed_files),
                    "failed_files": len(failed_files),
                    "total_questions": len(all_questions),
                    "questions": all_questions,
                    "processed_files": processed_files,
                    "failed_files": failed_files,
                    "processing_method": "zip_extraction",
                    "timestamp": datetime.now().isoformat(),
                    "status": "completed"
                }
                
                # Update progress with completion
                update_upload_progress(agent_run_id, "Completed", 100, "Question Agent completed successfully")
                
                # Store results in progress for retrieval
                with progress_lock:
                    if agent_run_id in upload_progress:
                        upload_progress[agent_run_id]["agent_results"] = agent_results
                if len(all_questions) > 0:
                    postAgentResultsAndQuestions(rfp_id, agent_results, auth_token)
                
                print(f"✅ Question Agent completed successfully!")
                print(f"📈 Final Stats: {len(all_questions)} questions from {len(processed_files)}/{total_files} files")
                
            except Exception as e:
                error_msg = str(e)
                print(f"❌ Question Agent error: {error_msg}")
                traceback.print_exc()
                update_upload_progress(agent_run_id, "error", 0, f"Agent Error: {error_msg}")
                
                # Store error results
                with progress_lock:
                    if agent_run_id in upload_progress:
                        upload_progress[agent_run_id]["agent_results"] = {
                            "agent_run_id": agent_run_id,
                            "agent_type": "question_generator",
                            "rfp_id": rfp_id,
                            "org_id": org_id,
                            "error": error_msg,
                            "timestamp": datetime.now().isoformat(),
                            "status": "failed"
                        }
                
                # Store results in progress for retrieval
                update_upload_progress(agent_run_id, "Completed", 100, "Question Agent completed successfully")
                
                with progress_lock:
                    if agent_run_id in upload_progress:
                        upload_progress[agent_run_id]["agent_results"] = agent_results
                
                print(f"✅ Question Agent completed. Generated {len(all_questions)} questions from {len(processed_files)} files")
                
            except Exception as e:
                error_msg = str(e)
                print(f"❌ Question Agent error: {error_msg}")
                traceback.print_exc()
                update_upload_progress(agent_run_id, "error", 0, f"Agent Error: {error_msg}")
        
        # Start agent processing in background thread
        agent_thread = Thread(target=run_question_agent_async)
        agent_thread.daemon = False
        agent_thread.start()
        
        return jsonify({
            "message": "Question Agent started successfully",
            "agent_run_id": agent_run_id,
            "agent_type": "question_generator",
            "rfp_id": rfp_id,
            "status": "running"
        }), 202

    except Exception as e:
        error_msg = str(e)
        print(f"❌ Question Agent Error: {error_msg}")
        traceback.print_exc()
        return jsonify({"error": "Internal server error", "details": error_msg}), 500

@app.route("/agent-status", methods=["GET", "OPTIONS"])
def get_agent_status():
    """Get agent run status and results"""
    if request.method == "OPTIONS":
        return "", 200
        
    agent_run_id = request.args.get("agentRunId")
    
    if not agent_run_id:
        return jsonify({"error": "agentRunId is required"}), 400
    
    with progress_lock:
        exists = agent_run_id in upload_progress
        status_data = upload_progress.get(agent_run_id, None)
    
    if not exists:
        return jsonify({"exists": False}), 200
    
    response_data = {
        "exists": True,
        "status": status_data
    }
    
    # Include agent results if completed
    if status_data and status_data.get("status") == "Completed" and "agent_results" in status_data:
        response_data["agent_results"] = status_data["agent_results"]
    
    return jsonify(response_data), 200

# Add this new endpoint after the existing /run-question-agent endpoint

@app.route("/run-question-agent-with-upload", methods=["POST", "OPTIONS"])
def run_question_agent_with_upload():
    """Run Question Agent with direct file upload - Upload files and generate questions"""
    if request.method == "OPTIONS":
        return "", 200
        
    if not VERTEX_AVAILABLE:
        return jsonify({"error": "AI generation service is not available"}), 503
        
    try:
        rfp_id = request.form.get("rfpId")
        org_id = request.form.get("orgId") 
        agent_run_id = request.form.get("agentRunId", str(uuid.uuid4()))
        auth_token = request.form.get("authToken")
        agent_id = request.form.get("agentId")
        
        print(f"🤖 Starting Question Agent with Upload for RFP: {rfp_id}, Agent Run ID: {agent_run_id}")
        
        if not rfp_id:
            return jsonify({"error": "rfpId is required"}), 400
            
        if not org_id:
            return jsonify({"error": "orgId is required"}), 400
        if not agent_id:  # ADD THIS VALIDATION
                    return jsonify({"error": "agentId is required"}), 400
        # Check if files were uploaded
        if not request.files:
            return jsonify({"error": "No files provided"}), 400

        uploaded_files = request.files.getlist('files')
        if not uploaded_files:
            return jsonify({"error": "No files found in request"}), 400

        # Validate files
        valid_files = []
        for file in uploaded_files:
            if file.filename and file.filename.strip():
                valid_files.append(file)
        
        if not valid_files:
            return jsonify({"error": "No valid files found"}), 400

        print(f"📁 Processing {len(valid_files)} files for question extraction")

        # Read all files into memory
        file_data_list = []
        for file in valid_files:
            try:
                file.seek(0)
                file_content = file.read()
                
                if len(file_content) == 0:
                    print(f"⚠️ Skipping empty file: {file.filename}")
                    continue
                    
                file_data_list.append({
                    'filename': file.filename,
                    'content': file_content,
                    'size': len(file_content),
                    'file_object': file  # Keep reference for backend upload
                })
                print(f"✅ Read {file.filename} ({len(file_content)} bytes)")
                
            except Exception as e:
                print(f"❌ Error reading file {file.filename}: {str(e)}")
                return jsonify({"error": f"Failed to read file {file.filename}: {str(e)}"}), 400

        if not file_data_list:
            return jsonify({"error": "No files could be read successfully"}), 400

        # Initialize progress tracking
        update_upload_progress(agent_run_id, "Processing", 0, "Starting Question Agent with file upload", "")
        
        def process_files_with_upload_async():
            try:
                # Step 1: Upload files to backend uploadProjectFiles API
                update_upload_progress(agent_run_id, "Processing", 5, "Uploading files to backend", "")
                
                backend_api_url = os.environ.get("BACKEND_API_URL", "http://localhost:8080")
                upload_url = f"{backend_api_url}/api/rfps/{rfp_id}/project-files"
                
                # Prepare files for backend upload
                files_for_upload = []
                for file_data in file_data_list:
                    # Reset file pointer and create new file object for upload
                    file_content = file_data['content']
                    files_for_upload.append(
                        ('files', (file_data['filename'], io.BytesIO(file_content), 'application/octet-stream'))
                    )
                
                headers = {}
                if auth_token:
                    headers["Authorization"] = f"Bearer {auth_token}"
                
                try:
                    print(f"📤 Uploading {len(files_for_upload)} files to backend: {upload_url}")
                    backend_response = requests.post(
                        upload_url,
                        files=files_for_upload,
                        headers=headers,
                        timeout=60
                    )
                    
                    if backend_response.status_code in [200, 201]:
                        print(f"✅ Successfully uploaded files to backend")
                        update_upload_progress(agent_run_id, "Processing", 15, "Files uploaded to backend successfully", "")
                    else:
                        print(f"⚠️ Backend upload failed: {backend_response.status_code} - {backend_response.text[:200]}")
                        update_upload_progress(agent_run_id, "Processing", 15, "Backend upload failed, continuing with AI processing", "")
                        
                except Exception as upload_error:
                    print(f"⚠️ Backend upload error: {str(upload_error)}")
                    update_upload_progress(agent_run_id, "Processing", 15, "Backend upload failed, continuing with AI processing", "")

                # Step 2: Process files for question extraction
                update_upload_progress(agent_run_id, "Processing", 20, "Starting AI question extraction", "")
                
                all_questions = []
                processed_files = []
                failed_files = []
                total_files = len(file_data_list)
                
                # Get supported file types
                supported, missing = check_file_processing_dependencies()
                all_supported = supported["always"] + supported["conditional"] + ["png", "jpg", "jpeg", "gif", "webp"]
                
                for index, file_data in enumerate(file_data_list):
                    filename = file_data['filename']
                    file_content = file_data['content']
                    file_ext = filename.split(".")[-1].lower()
                    
                    # Validate file type
                    if file_ext not in all_supported:
                        failed_files.append({
                            "filename": filename,
                            "error": f"Unsupported file type: {file_ext}. Supported: {', '.join(sorted(all_supported))}"
                        })
                        continue
                    
                    try:
                        # Update progress
                        progress = 20 + ((index / total_files) * 60)  # 20% to 80%
                        update_upload_progress(agent_run_id, "Processing", progress, 
                                            f"AI processing file {index + 1}/{total_files}: {filename}")
                        
                        # Save file content to temporary file
                        save_path = os.path.join(UPLOAD_FOLDER, f"{agent_run_id}_{index}_{filename}")
                        with open(save_path, 'wb') as f:
                            f.write(file_content)
                        
                        # Extract questions using AI
                        questions = extract_questions_with_ai_direct(save_path, filename)
                        
                        if not questions:
                            failed_files.append({
                                "filename": filename,
                                "error": "AI could not extract questions from this file"
                            })
                        else:
                            processed_files.append({
                                "filename": filename,
                                "file_type": file_ext,
                                "questions": questions,
                                "question_count": len(questions)
                            })
                            all_questions.extend(questions)
                            print(f"✅ Extracted {len(questions)} questions from {filename}")
                        
                        # Clean up temporary file
                        try:
                            os.remove(save_path)
                        except Exception as e:
                            print(f"⚠️ Failed to remove temp file {save_path}: {e}")
                            
                    except Exception as file_error:
                        print(f"❌ Error processing file {filename}: {str(file_error)}")
                        failed_files.append({
                            "filename": filename,
                            "error": str(file_error)
                        })
                
                # Step 3: Compile and store results
                update_upload_progress(agent_run_id, "Processing", 85, "Compiling agent results")
                
                print(f"\n📊 Question Agent Results Summary:")
                print(f"  Total files: {total_files}")
                print(f"  Successfully processed: {len(processed_files)}")
                print(f"  Failed: {len(failed_files)}")
                print(f"  Total questions generated: {len(all_questions)}")
                
                # Store agent results (same format as before)
                agent_results = {
                    "agent_run_id": agent_run_id,
                    "agent_type": "question_generator",
                    "agent_id": agent_id,  # ADD THIS LINE
                    "rfp_id": rfp_id,
                    "org_id": org_id,
                    "total_files": total_files,
                    "successful_files": len(processed_files),
                    "failed_files": len(failed_files),
                    "total_questions": len(all_questions),
                    "questions": all_questions,
                    "processed_files": processed_files,
                    "failed_files": failed_files,
                    "processing_method": "direct_upload",
                    "timestamp": datetime.now().isoformat(),
                    "status": "completed"
                }
                
                # Update progress with completion
                update_upload_progress(agent_run_id, "Completed", 100, "Question Agent completed successfully")
                
                # Store results in progress for retrieval
                with progress_lock:
                    if agent_run_id in upload_progress:
                        upload_progress[agent_run_id]["agent_results"] = agent_results
                if len(all_questions) > 0:
                    postAgentResultsAndQuestions(rfp_id, agent_results, auth_token,agent_id)
                
                print(f"✅ Question Agent with upload completed successfully!")
                print(f"📈 Final Stats: {len(all_questions)} questions from {len(processed_files)}/{total_files} files")
                
            except Exception as e:
                error_msg = str(e)
                print(f"❌ Question Agent error: {error_msg}")
                traceback.print_exc()
                update_upload_progress(agent_run_id, "error", 0, f"Agent Error: {error_msg}")
        
        # Start processing in background thread
        processing_thread = Thread(target=process_files_with_upload_async)
        processing_thread.daemon = False
        processing_thread.start()
        
        return jsonify({
            "message": "Question Agent with file upload started successfully",
            "agent_run_id": agent_run_id,
            "agent_type": "question_generator",
            "rfp_id": rfp_id,
            "files_count": len(file_data_list),
            "status": "running"
        }), 202

    except Exception as e:
        error_msg = str(e)
        print(f"❌ Question Agent with Upload Error: {error_msg}")
        traceback.print_exc()
        return jsonify({"error": "Internal server error", "details": error_msg}), 500
    

@app.route("/text-operations", methods=["POST", "OPTIONS"])
def text_operations():
    """Perform AI operations on selected text (rewrite, summarize, expand, etc.)"""
    if request.method == "OPTIONS":
        return "", 200
        
    if not VERTEX_AVAILABLE:
        return jsonify({"error": "AI generation service is not available"}), 503
        
    try:
        data = request.get_json(silent=True)
        if not data:
            return jsonify({"error": "No valid JSON input found"}), 400

        selected_text = data.get("selectedText")
        operation = data.get("operation")  # rewrite, summarize, expand, improve, etc.
        context = data.get("context", "")  # Optional surrounding context
        file_type = data.get("fileType", "")  # Optional file type context
        
        if not selected_text or not operation:
            return jsonify({"error": "selectedText and operation are required"}), 400

        print(f"🔧 Text operation: {operation} on text: {selected_text[:100]}...")

        # Define operation prompts - UPDATED WITH ALL FRONTEND OPERATIONS
        operation_prompts = {
            "rewrite": f"""
Rewrite the following text to make it clearer and more professional while maintaining the same meaning:

Text: {selected_text}

Provide only the rewritten text without any additional explanations or formatting.
""",
            "expand": f"""
Expand the following text with more detail and explanation while maintaining the same tone and style:

Text: {selected_text}

Provide only the expanded text without any additional explanations or formatting.
""",
            "simplify": f"""
Simplify the following text to make it easier to understand while keeping the essential meaning:

Text: {selected_text}

Provide only the simplified text without any additional explanations or formatting.
""",
            "formalize": f"""
Make the following text more formal and professional in tone:

Text: {selected_text}

Provide only the formalized text without any additional explanations or formatting.
""",
            "casual": f"""
Make the following text more casual and conversational while maintaining professionalism:

Text: {selected_text}

Provide only the casual version without any additional explanations or formatting.
""",
            "shorten": f"""
Shorten the following text while keeping the key information and main message:

Text: {selected_text}

Provide only the shortened text without any additional explanations or formatting.
""",
            "grammar": f"""
Fix any grammar, spelling, and punctuation errors in the following text:

Text: {selected_text}

Provide only the corrected text without any additional explanations or formatting.
""",
            "professional": f"""
Make the following text more professional and business-appropriate:

Text: {selected_text}

Provide only the professional version without any additional explanations or formatting.
""",
            # Additional operations for completeness
            "summarize": f"""
Summarize the following text concisely while retaining the key points:

Text: {selected_text}

Provide only the summary without any additional explanations or formatting.
""",
            "improve": f"""
Improve the following text for better clarity, grammar, flow, and professionalism:

Text: {selected_text}

Provide only the improved text without any additional explanations or formatting.
""",
            "bullet_points": f"""
Convert the following text into clear, well-organized bullet points:

Text: {selected_text}

Provide only the bullet points without any additional explanations or formatting.
"""
        }

        if operation not in operation_prompts:
            supported_ops = list(operation_prompts.keys())
            return jsonify({
                "error": f"Unsupported operation: {operation}. Supported operations: {', '.join(supported_ops)}"
            }), 400

        prompt = operation_prompts[operation]

        # Add context if provided
        if context and context.strip():
            prompt += f"\n\nContext (for reference): {context[:500]}..."  # Limit context length

        # Add file type context if provided
        if file_type:
            prompt += f"\n\nNote: This text is from a {file_type.upper()} file."

        from vertexai.generative_models import GenerativeModel, Part, SafetySetting
        model = GenerativeModel("gemini-2.0-flash")

        response = model.generate_content(
            [Part.from_text(prompt)],
            generation_config={
                "max_output_tokens": 2048,
                "temperature": 0.3,
                "top_p": 0.8,
            },
            safety_settings=[
                SafetySetting(
                    category=SafetySetting.HarmCategory.HARM_CATEGORY_HATE_SPEECH,
                    threshold=SafetySetting.HarmBlockThreshold.OFF
                ),
                SafetySetting(
                    category=SafetySetting.HarmCategory.HARM_CATEGORY_DANGEROUS_CONTENT,
                    threshold=SafetySetting.HarmBlockThreshold.OFF
                ),
                SafetySetting(
                    category=SafetySetting.HarmCategory.HARM_CATEGORY_SEXUALLY_EXPLICIT,
                    threshold=SafetySetting.HarmBlockThreshold.OFF
                ),
                SafetySetting(
                    category=SafetySetting.HarmCategory.HARM_CATEGORY_HARASSMENT,
                    threshold=SafetySetting.HarmBlockThreshold.OFF
                ),
            ],
            stream=False
        )

        processed_text = response.text.strip()

        # Basic validation of the response
        if not processed_text:
            return jsonify({"error": "AI returned empty response"}), 500

        print(f"✅ Text operation '{operation}' completed successfully")

        return jsonify({
            "original_text": selected_text,
            "operation": operation,
            "processed_text": processed_text,
            "success": True,
            "character_count": {
                "original": len(selected_text),
                "processed": len(processed_text)
            }
        }), 200

    except Exception as e:
        error_msg = str(e)
        print(f"❌ Text Operation Error: {error_msg}")
        traceback.print_exc()
        return jsonify({
            "error": "Internal server error", 
            "details": error_msg,
            "success": False
        }), 500

@app.route("/api/v2/text-operations", methods=["POST", "OPTIONS"])
def text_operations_v2():
    """V2 Text operations with knowledge base support and custom prompts"""
    if request.method == "OPTIONS":
        return "", 200
        
    if not VERTEX_AVAILABLE:
        return jsonify({"error": "AI generation service is not available"}), 503
        
    if not FIREBASE_AVAILABLE:
        return jsonify({"error": "Firebase is not available"}), 503
        
    if not OPENAI_AVAILABLE:
        return jsonify({"error": "OpenAI embedding service is not available"}), 503
        
    try:
        data = request.get_json(silent=True)
        if not data:
            return jsonify({"error": "No valid JSON input found"}), 400

        selected_text = data.get("selectedText")
        operation = data.get("operation")  # rewrite, summarize, expand, improve, custom, etc.
        custom_prompt = data.get("customPrompt", "")  # Custom prompt for 'custom' operation
        project_id = data.get("projectId")
        org_id = data.get("orgId")
        knowledge_base_option = data.get("knowledgeBaseOption", "global")  # "global" or "specific"
        
        if not selected_text:
            return jsonify({"error": "selectedText is required"}), 400
            
        if not operation:
            return jsonify({"error": "operation is required"}), 400
            
        if operation == "custom" and not custom_prompt:
            return jsonify({"error": "customPrompt is required for custom operation"}), 400

        print(f"🔧 V2 Text operation: {operation} on text: {selected_text[:100]}... (Knowledge: {knowledge_base_option})")

        # Define operation prompts - same as v1 but enhanced
        operation_prompts = {
            "rewrite": f"""Rewrite the following text to make it clearer and more professional while maintaining the same meaning:

Text: {selected_text}

Provide only the rewritten text without any additional explanations or formatting.""",
            "expand": f"""Expand the following text with more detail and explanation while maintaining the same tone and style:

Text: {selected_text}

Provide only the expanded text without any additional explanations or formatting.""",
            "simplify": f"""Simplify the following text to make it easier to understand while keeping the essential meaning:

Text: {selected_text}

Provide only the simplified text without any additional explanations or formatting.""",
            "formalize": f"""Make the following text more formal and professional in tone:

Text: {selected_text}

Provide only the formalized text without any additional explanations or formatting.""",
            "casual": f"""Make the following text more casual and conversational while maintaining professionalism:

Text: {selected_text}

Provide only the casual version without any additional explanations or formatting.""",
            "shorten": f"""Shorten the following text while keeping the key information and main message:

Text: {selected_text}

Provide only the shortened text without any additional explanations or formatting.""",
            "grammar": f"""Fix any grammar, spelling, and punctuation errors in the following text:

Text: {selected_text}

Provide only the corrected text without any additional explanations or formatting.""",
            "professional": f"""Make the following text more professional and business-appropriate:

Text: {selected_text}

Provide only the professional version without any additional explanations or formatting.""",
            "summarize": f"""Summarize the following text concisely while retaining the key points:

Text: {selected_text}

Provide only the summary without any additional explanations or formatting.""",
            "improve": f"""Improve the following text for better clarity, grammar, flow, and professionalism:

Text: {selected_text}

Provide only the improved text without any additional explanations or formatting.""",
            "bullet_points": f"""Convert the following text into clear, well-organized bullet points:

Text: {selected_text}

Provide only the bullet points without any additional explanations or formatting.""",
            "custom": custom_prompt + f"""

Text to process: {selected_text}

Provide only the processed text without any additional explanations or formatting."""
        }

        if operation not in operation_prompts:
            supported_ops = list(operation_prompts.keys())
            return jsonify({
                "error": f"Unsupported operation: {operation}. Supported operations: {', '.join(supported_ops)}"
            }), 400

        base_prompt = operation_prompts[operation]
        
        # Add knowledge base context if project info is provided
        context_chunks = []
        if project_id and org_id:
            query_embedding = np.array(embed_query(selected_text))
            
            if knowledge_base_option == "global":
                # Search organization-level documents (global knowledge base)
                print("🌐 Using global knowledge base for text operations")
                org_files_ref = db.collection("document_embeddings").document(f"org-{org_id}").collection("files")
                org_files = org_files_ref.stream()
                
                retrieved_docs = []
                for file_doc in org_files:
                    file_data = file_doc.to_dict()
                    chunks_ref = file_doc.reference.collection("chunks")
                    chunks = chunks_ref.stream()
                    
                    for chunk_doc in chunks:
                        chunk_data = chunk_doc.to_dict()
                        chunk_embedding = np.array(chunk_data["embedding"])
                        score = np.dot(query_embedding, chunk_embedding) / (
                            np.linalg.norm(query_embedding) * np.linalg.norm(chunk_embedding)
                        )
                        
                        if score >= 0.15:  # Lower threshold for context
                            # Handle different chunk storage formats: "content" (v3.0.0) or "text" (v2)
                            chunk_content = chunk_data.get("content") or chunk_data.get("text", "")
                            retrieved_docs.append({
                                "content": chunk_content, 
                                "score": float(score)
                            })
                
                context_chunks = [doc["content"] for doc in sorted(retrieved_docs, key=lambda x: x["score"], reverse=True)[:3]]
                
            elif knowledge_base_option == "specific":
                # Use project-specific support documents
                print("📁 Using specific project support documents for text operations")
                project_files_ref = (db.collection("org_project_support_embeddings")
                                   .document(f"org-{org_id}")
                                   .collection("projects")
                                   .document(f"project-{project_id}")
                                   .collection("files"))
                
                files = project_files_ref.stream()
                retrieved_docs = []
                
                for file_doc in files:
                    file_data = file_doc.to_dict()
                    chunks_ref = file_doc.reference.collection("chunks")
                    chunks = chunks_ref.stream()
                    
                    for chunk_doc in chunks:
                        chunk_data = chunk_doc.to_dict()
                        chunk_embedding = np.array(chunk_data["embedding"])
                        score = np.dot(query_embedding, chunk_embedding) / (
                            np.linalg.norm(query_embedding) * np.linalg.norm(chunk_embedding)
                        )
                        
                        if score >= 0.15:  # Lower threshold for context
                            # Handle different chunk storage formats: "content" (v3.0.0) or "text" (v2)
                            chunk_content = chunk_data.get("content") or chunk_data.get("text", "")
                            retrieved_docs.append({
                                "content": chunk_content, 
                                "score": float(score)
                            })
                
                context_chunks = [doc["content"] for doc in sorted(retrieved_docs, key=lambda x: x["score"], reverse=True)[:3]]

        # Enhance prompt with knowledge base context
        final_prompt = base_prompt
        if context_chunks:
            context_text = "\n\n".join(context_chunks)
            final_prompt += f"\n\nRelevant context from your knowledge base:\n{context_text[:1000]}...\n\nUse this context to inform your response but focus on processing the text as requested."

        from vertexai.generative_models import GenerativeModel, Part, SafetySetting
        model = GenerativeModel("gemini-2.0-flash")

        response = model.generate_content(
            [Part.from_text(final_prompt)],
            generation_config={
                "max_output_tokens": 2048,
                "temperature": 0.3,
                "top_p": 0.8,
            },
            safety_settings=[
                SafetySetting(
                    category=SafetySetting.HarmCategory.HARM_CATEGORY_HATE_SPEECH,
                    threshold=SafetySetting.HarmBlockThreshold.OFF
                ),
                SafetySetting(
                    category=SafetySetting.HarmCategory.HARM_CATEGORY_DANGEROUS_CONTENT,
                    threshold=SafetySetting.HarmBlockThreshold.OFF
                ),
                SafetySetting(
                    category=SafetySetting.HarmCategory.HARM_CATEGORY_SEXUALLY_EXPLICIT,
                    threshold=SafetySetting.HarmBlockThreshold.OFF
                ),
                SafetySetting(
                    category=SafetySetting.HarmCategory.HARM_CATEGORY_HARASSMENT,
                    threshold=SafetySetting.HarmBlockThreshold.OFF
                ),
            ],
            stream=False
        )

        processed_text = response.text.strip()

        if not processed_text:
            return jsonify({"error": "AI returned empty response"}), 500

        print(f"✅ V2 Text operation '{operation}' completed successfully with {knowledge_base_option} knowledge")

        return jsonify({
            "original_text": selected_text,
            "operation": operation,
            "processed_text": processed_text,
            "success": True,
            "knowledge_type": knowledge_base_option,
            "context_used": len(context_chunks) > 0,
            "character_count": {
                "original": len(selected_text),
                "processed": len(processed_text)
            }
        }), 200

    except Exception as e:
        error_msg = str(e)
        print(f"❌ V2 Text Operation Error: {error_msg}")
        traceback.print_exc()
        return jsonify({
            "error": "Internal server error", 
            "details": error_msg,
            "success": False
        }), 500

# Add this endpoint to your existing Python Flask application after the /run-question-agent-with-upload endpoint

# Add this endpoint to your existing Python Flask application after the /run-question-agent-with-upload endpoint

@app.route("/run-proposal-narrative-agent", methods=["POST", "OPTIONS"])
def run_proposal_narrative_agent():
    """Run Proposal Narrative Writer Agent - Process RFP and generate proposal sections"""
    if request.method == "OPTIONS":
        return "", 200
        
    if not VERTEX_AVAILABLE:
        return jsonify({"error": "AI generation service is not available"}), 503
        
    try:
        rfp_id = request.form.get("rfpId")
        org_id = request.form.get("orgId") 
        agent_run_id = request.form.get("agentRunId", str(uuid.uuid4()))
        auth_token = request.form.get("authToken")
        
        print(f"🎯 Starting Proposal Narrative Agent for RFP: {rfp_id}, Agent Run ID: {agent_run_id}")
        
        if not rfp_id:
            return jsonify({"error": "rfpId is required"}), 400
            
        if not org_id:
            return jsonify({"error": "orgId is required"}), 400

        # Get project file (RFP document) - required
        project_file = request.files.get('projectFile')
        if not project_file or not project_file.filename:
            return jsonify({"error": "Project file (RFP document) is required"}), 400

        # Get support files - optional
        support_files = request.files.getlist('supportFiles')
        
        print(f"📄 Processing RFP document: {project_file.filename}")
        if support_files:
            print(f"📚 Processing {len(support_files)} support files for knowledge base")

        # Read all files into memory before starting async processing
        try:
            # Read project file content
            project_file.seek(0)
            project_content = project_file.read()
            project_file_data = {
                'filename': project_file.filename,
                'content': project_content,
                'size': len(project_content)
            }
            
            if len(project_file_data['content']) == 0:
                return jsonify({"error": "Project file is empty"}), 400
            
            # Read support files content
            support_files_data = []
            for support_file in support_files:
                if support_file.filename and support_file.filename.strip():
                    support_file.seek(0)
                    support_content = support_file.read()
                    
                    if len(support_content) > 0:
                        support_files_data.append({
                            'filename': support_file.filename,
                            'content': support_content,
                            'size': len(support_content)
                        })
                        print(f"✅ Read support file: {support_file.filename} ({len(support_content)} bytes)")
                    else:
                        print(f"⚠️ Skipping empty support file: {support_file.filename}")
            
            print(f"✅ Read project file: {project_file_data['filename']} ({len(project_file_data['content'])} bytes)")
            print(f"✅ Read {len(support_files_data)} support files")
            
        except Exception as e:
            print(f"❌ Error reading file contents: {str(e)}")
            return jsonify({"error": f"Failed to read file contents: {str(e)}"}), 400

        # Initialize progress tracking
        update_upload_progress(agent_run_id, "Processing", 0, "Starting Proposal Narrative Agent", "")
        
        def process_proposal_agent_async():
            try:
                # Step 1: Process support files for knowledge base (if any)
                support_results = []
                if support_files_data and len(support_files_data) > 0:
                    update_upload_progress(agent_run_id, "Processing", 10, "Processing support files for knowledge base", "")
                    
                    for idx, support_file_data in enumerate(support_files_data):
                        filename = support_file_data['filename']
                        file_content = support_file_data['content']
                        
                        try:
                            # Save support file temporarily
                            support_path = os.path.join(UPLOAD_FOLDER, f"{agent_run_id}_support_{idx}_{filename}")
                            with open(support_path, 'wb') as f:
                                f.write(file_content)
                            
                            # Upload to RFP-specific knowledge base
                            support_upload_id = f"support_{agent_run_id}_{idx}"
                            
                            # Process for embedding using RFP support document endpoint logic
                            file_ext = filename.split(".")[-1].lower()
                            chunks = parse_and_chunk(support_path, file_ext, chunk_size=50, max_chunks=500)
                            
                            if chunks:
                                embeddings = embed_chunks(chunks, upload_id=support_upload_id, org_id=org_id, filename=filename)
                                
                                # Store in RFP-specific Firestore collection
                                rfp_doc_ref = db.collection("org_rfp_support_embeddings").document(f"rfp-{rfp_id}").collection("files").document(support_upload_id)
                                
                                # Store file metadata
                                rfp_doc_ref.set({
                                    "filename": filename,
                                    "file_id": support_upload_id,
                                    "upload_id": support_upload_id,
                                    "rfp_id": rfp_id,
                                    "org_id": org_id,
                                    "file_type": file_ext,
                                    "document_type": "support",
                                    "agent_run_id": agent_run_id,
                                    "created_at": firestore.SERVER_TIMESTAMP,
                                    "chunk_count": len(chunks),
                                    "processing_version": "4.0.0"
                                })
                                
                                # Store chunks in batches - OPTIMIZED for 50x faster writes
                                batch_size = 500  # Firestore allows up to 500 writes per batch
                                total_chunks = len(chunks)
                                
                                print(f"🚀 Storing {total_chunks} agent chunks in batches of {batch_size}")
                                
                                for i in range(0, total_chunks, batch_size):
                                    batch = db.batch()
                                    end_idx = min(i + batch_size, total_chunks)
                                    
                                    for j in range(i, end_idx):
                                        chunk_ref = rfp_doc_ref.collection("chunks").document(str(j))
                                        batch.set(chunk_ref, {
                                            "content": chunks[j],
                                            "embedding": embeddings[j],
                                            "index": j
                                        })
                                    
                                    batch.commit()
                                    print(f"✓ Stored agent batch {i//batch_size + 1}: {end_idx}/{total_chunks} chunks")
                                    
                                    # Only clean up memory for large batches
                                    if end_idx - i >= 100:
                                        del batch
                                        import gc
                                        gc.collect()
                                
                                support_results.append({
                                    "filename": filename,
                                    "status": "success",
                                    "chunks": len(chunks)
                                })
                                
                                print(f"✅ Support file processed for knowledge base: {filename}")
                            else:
                                support_results.append({
                                    "filename": filename,
                                    "status": "failed",
                                    "error": "No content extracted"
                                })
                            
                            # Clean up temp file
                            try:
                                os.remove(support_path)
                            except Exception as e:
                                print(f"⚠️ Failed to remove temp file {support_path}: {e}")
                                
                        except Exception as support_error:
                            print(f"❌ Error processing support file {filename}: {str(support_error)}")
                            support_results.append({
                                "filename": filename,
                                "status": "failed",
                                "error": str(support_error)
                            })
                
                # Step 2: Process main RFP document and extract proposal sections
                update_upload_progress(agent_run_id, "Processing", 40, "Analyzing RFP document and extracting sections", "")
                
                # Save project file temporarily using the content we read earlier
                project_filename = project_file_data['filename']
                project_content = project_file_data['content']
                project_path = os.path.join(UPLOAD_FOLDER, f"{agent_run_id}_project_{project_filename}")
                
                try:
                    # Write project file content to temporary file
                    with open(project_path, 'wb') as f:
                        f.write(project_content)
                    
                    print(f"✅ Project file saved temporarily: {project_path}")
                    
                    # Extract proposal sections and requirements from RFP using AI
                    proposal_sections = extract_proposal_sections_with_ai(project_path, project_filename)
                    
                    if not proposal_sections:
                        raise Exception("Could not extract proposal sections from RFP document")
                    
                    print(f"✅ Extracted {len(proposal_sections)} proposal sections from RFP")
                    
                    # Clean up project file
                    try:
                        os.remove(project_path)
                        print(f"🧹 Cleaned up project file: {project_path}")
                    except Exception as e:
                        print(f"⚠️ Failed to remove temp file {project_path}: {e}")
                        
                except Exception as rfp_error:
                    raise Exception(f"Error processing RFP document: {str(rfp_error)}")
                
                # Step 3: Compile results
                update_upload_progress(agent_run_id, "Processing", 80, "Compiling proposal narrative results", "")
                
                # Create agent results
                agent_results = {
                    "agent_run_id": agent_run_id,
                    "agent_type": "proposal_narrative_writer",
                    "rfp_id": rfp_id,
                    "org_id": org_id,
                    "project_file": {
                        "filename": project_filename,
                        "status": "success"
                    },
                    "support_files_processed": len(support_results),
                    "support_files_success": len([r for r in support_results if r["status"] == "success"]),
                    "support_files": support_results,
                    "sections": proposal_sections,
                    "total_sections": len(proposal_sections),
                    "knowledge_base_enhanced": len([r for r in support_results if r["status"] == "success"]) > 0,
                    "processing_method": "proposal_narrative_generation",
                    "timestamp": datetime.now().isoformat(),
                    "status": "completed"
                }
                
                # Update progress with completion
                update_upload_progress(agent_run_id, "Completed", 100, "Proposal Narrative Agent completed successfully")
                
                # Store results in progress for retrieval
                with progress_lock:
                    if agent_run_id in upload_progress:
                        upload_progress[agent_run_id]["agent_results"] = agent_results
                
                # Post results to backend
                post_proposal_agent_results_to_backend(rfp_id, agent_results, auth_token)
                
                print(f"✅ Proposal Narrative Agent completed successfully!")
                print(f"📈 Final Stats: {len(proposal_sections)} sections extracted, {len([r for r in support_results if r['status'] == 'success'])} support files processed")
                
            except Exception as e:
                error_msg = str(e)
                print(f"❌ Proposal Narrative Agent error: {error_msg}")
                traceback.print_exc()
                update_upload_progress(agent_run_id, "error", 0, f"Agent Error: {error_msg}")
                
                # Store error results
                with progress_lock:
                    if agent_run_id in upload_progress:
                        upload_progress[agent_run_id]["agent_results"] = {
                            "agent_run_id": agent_run_id,
                            "agent_type": "proposal_narrative_writer",
                            "rfp_id": rfp_id,
                            "org_id": org_id,
                            "error": error_msg,
                            "timestamp": datetime.now().isoformat(),
                            "status": "failed"
                        }
        
        # Start processing in background thread
        processing_thread = Thread(target=process_proposal_agent_async)
        processing_thread.daemon = False
        processing_thread.start()
        
        return jsonify({
            "message": "Proposal Narrative Agent started successfully",
            "agent_run_id": agent_run_id,
            "agent_type": "proposal_narrative_writer",
            "rfp_id": rfp_id,
            "project_file": project_file.filename,
            "support_files_count": len(support_files) if support_files else 0,
            "status": "running"
        }), 202

    except Exception as e:
        error_msg = str(e)
        print(f"❌ Proposal Narrative Agent Error: {error_msg}")
        traceback.print_exc()
        return jsonify({"error": "Internal server error", "details": error_msg}), 500

def extract_proposal_sections_with_ai(file_path, filename):
    """Extract proposal sections and requirements from RFP document using Vertex AI"""
    if not VERTEX_AVAILABLE:
        print("❌ Vertex AI not available")
        return []
        
    try:
        # Verify file exists and is readable
        if not os.path.exists(file_path):
            print(f"❌ File does not exist: {file_path}")
            return []
        
        file_size = os.path.getsize(file_path)
        if file_size == 0:
            print(f"❌ File is empty: {file_path}")
            return []
            
        print(f"📄 Processing RFP file: {filename} ({file_size} bytes)")
        
        # Get file extension to determine how to handle it
        file_ext = filename.split(".")[-1].lower()
        
        # Customize prompt for proposal section extraction
        prompt = f"""
Analyze the uploaded RFP document and extract all sections that require a proposal response. 

For each section that needs a response in the proposal, provide:
1. The section title/name
2. The type of section (e.g., "Technical Approach", "Methodology", "Qualifications", "Experience", "Timeline", "Budget", "Compliance", "Company Profile")
3. The specific requirements, questions, or criteria that need to be addressed

Focus on sections like:
- Technical approach and methodology
- Company qualifications and experience
- Project team and personnel
- Case studies and past performance
- Compliance and certifications
- Timeline and project management
- Budget and pricing (if mentioned)
- Company profile and capabilities
- Risk management
- Quality assurance
- Any specific questions posed to bidders

Format your response as a valid JSON array with objects containing "title", "type", and "content" fields:

[
  {{
    "title": "Technical Approach",
    "type": "Technical Approach",
    "content": "Describe your technical approach to implementing the solution, including methodology, tools, and processes."
  }},
  {{
    "title": "Company Experience", 
    "type": "Qualifications",
    "content": "Provide examples of similar projects completed in the last 5 years, including client references and outcomes achieved."
  }}
]

IMPORTANT: Return ONLY the JSON array, no additional text or formatting.
"""

        model = GenerativeModel("gemini-2.0-flash")
        
        # Prepare the content based on file type
        parts = [Part.from_text(prompt)]
        
        # Define MIME type mapping
        mime_mapping = {
            "pdf": "application/pdf",
            "txt": "text/plain", 
            "csv": "text/csv",
            "json": "application/json",
            "png": "image/png",
            "jpg": "image/jpeg",
            "jpeg": "image/jpeg", 
            "gif": "image/gif",
            "webp": "image/webp",
        }
        
        # Files that Gemini can handle directly
        gemini_supported_files = [
            # Documents
            "pdf", "txt", "csv", "md", "markdown", "html", "htm", "xml", "yaml", "yml",
            
            # Images  
            "png", "jpg", "jpeg", "gif", "webp", "bmp", "tiff", "tif", "svg", "heic", "heif",
            
            # Code files (text-based)
            "py", "js", "ts", "jsx", "tsx", "java", "cpp", "c", "cs", "php", "rb", "go", "rs", 
            "swift", "kt", "scala", "sql", "sh", "ps1",
            
            # Configuration files
            "log", "cfg", "conf", "ini", "properties", "env", "lock", "gitignore", "dockerfile",
            
            # Audio (Gemini 2.0 supports audio)
            "mp3", "wav", "m4a", "aac", "ogg", "flac",
            
            # Video (Gemini 2.0 supports video)
            "mp4", "mov", "avi", "wmv", "webm", "mkv", "m4v"
        ]
        
        if file_ext in gemini_supported_files:
            try:
                # Read file as bytes for direct upload
                with open(file_path, "rb") as f:
                    file_data = f.read()
                
                mime_type = mime_mapping.get(file_ext, "application/octet-stream")
                
                # Validate file size for Gemini (max ~10MB)
                if len(file_data) > 10 * 1024 * 1024:  # 10MB limit
                    print(f"⚠️ File too large for direct upload ({len(file_data)} bytes), falling back to text extraction")
                    content = extract_text_from_file(file_path, file_ext)
                    if content and not content.startswith("ERROR:"):
                        if len(content) > 30000:
                            content = content[:30000] + "..."
                        parts.append(Part.from_text(f"\nRFP Document Content:\n{content}"))
                    else:
                        return []
                else:
                    # Add file part directly
                    file_part = Part.from_data(data=file_data, mime_type=mime_type)
                    parts.append(file_part)
                    
                    print(f"✅ Using direct multimodal upload for {filename} ({mime_type}, {len(file_data)} bytes)")
                
            except Exception as e:
                print(f"⚠️ Direct upload failed for {filename}, falling back to text extraction: {e}")
                content = extract_text_from_file(file_path, file_ext)
                if content and not content.startswith("ERROR:"):
                    if len(content) > 30000:
                        content = content[:30000] + "..."
                    parts.append(Part.from_text(f"\nRFP Document Content:\n{content}"))
                else:
                    return []
        else:
            # For Excel, Word, and other files - extract text first
            print(f"📄 Extracting text from {filename} (file type: {file_ext})")
            content = extract_text_from_file(file_path, file_ext)
            if not content or content.startswith("ERROR:"):
                print(f"❌ No content extracted from {filename}: {content}")
                return []
            
            # Limit content length
            if len(content) > 30000:
                content = content[:30000] + "..."
            
            parts.append(Part.from_text(f"\nRFP Document Content:\n{content}"))
        
        # Generate response
        print(f"🤖 Sending request to Gemini for proposal section extraction: {filename}")
        response = model.generate_content(
            parts,
            generation_config={
                "max_output_tokens": 4096,
                "temperature": 0.3,
                "top_p": 0.8,
            },
            safety_settings=[
                SafetySetting(
                    category=SafetySetting.HarmCategory.HARM_CATEGORY_HATE_SPEECH,
                    threshold=SafetySetting.HarmBlockThreshold.OFF
                ),
                SafetySetting(
                    category=SafetySetting.HarmCategory.HARM_CATEGORY_DANGEROUS_CONTENT,
                    threshold=SafetySetting.HarmBlockThreshold.OFF
                ),
                SafetySetting(
                    category=SafetySetting.HarmCategory.HARM_CATEGORY_SEXUALLY_EXPLICIT,
                    threshold=SafetySetting.HarmBlockThreshold.OFF
                ),
                SafetySetting(
                    category=SafetySetting.HarmCategory.HARM_CATEGORY_HARASSMENT,
                    threshold=SafetySetting.HarmBlockThreshold.OFF
                ),
            ],
            stream=False
        )
        
        print(f"✅ Received response from Gemini for {filename}")
        
        # Parse the JSON response
        response_text = response.text.strip()
        print(f"📝 Raw response length: {len(response_text)} characters")
        
        # Clean up the response - remove any markdown formatting
        if response_text.startswith("```json"):
            response_text = response_text.replace("```json", "").replace("```", "").strip()
        elif response_text.startswith("```"):
            response_text = response_text.replace("```", "").strip()
        
        # Try to extract JSON from the response
        try:
            sections = json.loads(response_text)
            
            # Validate the structure
            if isinstance(sections, list):
                valid_sections = []
                for section in sections:
                    if isinstance(section, dict) and "title" in section and "type" in section and "content" in section:
                        valid_sections.append({
                            "title": str(section["title"]).strip(),
                            "type": str(section["type"]).strip(),
                            "content": str(section["content"]).strip()
                        })
                
                print(f"✅ Successfully parsed {len(valid_sections)} proposal sections from {filename}")
                return valid_sections[:20]  # Limit to 20 sections
            else:
                print(f"⚠️ Invalid JSON structure from AI for {filename} (not a list)")
                return []
                
        except json.JSONDecodeError as e:
            print(f"⚠️ JSON parsing error for {filename}: {e}")
            print(f"Raw response preview: {response_text[:500]}...")
            
            # Fallback: try to extract sections manually using regex
            sections = []
            
            # Try to find JSON objects in the text
            json_pattern = r'\{\s*"title":\s*"([^"]+)"\s*,\s*"type":\s*"([^"]+)"\s*,\s*"content":\s*"([^"]+)"\s*\}'
            matches = re.findall(json_pattern, response_text)
            
            for title, section_type, content in matches:
                sections.append({
                    "title": title.strip(),
                    "type": section_type.strip(),
                    "content": content.strip()
                })
            
            if sections:
                print(f"✅ Extracted {len(sections)} sections using fallback regex for {filename}")
            else:
                print(f"❌ Could not extract any sections from {filename}")
                
            return sections[:20]
            
    except Exception as e:
        print(f"❌ Error extracting proposal sections from {filename}: {str(e)}")
        traceback.print_exc()
        return []

def post_proposal_agent_results_to_backend(rfp_id, agent_results, auth_token):
    """Post proposal agent results to backend"""
    if not rfp_id or not agent_results:
        return

    try:
        backend_api_url = os.environ.get("BACKEND_API_URL", "http://localhost:8080")
        agent_response_url = f"{backend_api_url}/api/rfps/{rfp_id}/agent-response"
        
        headers = {"Content-Type": "application/json"}
        if auth_token:
            headers["Authorization"] = f"Bearer {auth_token}"

        # Parse results into string format
        response_string = json.dumps({
            "summary": {
                "total_sections": agent_results.get("total_sections", 0),
                "support_files_processed": agent_results.get("support_files_processed", 0),
                "support_files_success": agent_results.get("support_files_success", 0),
                "knowledge_base_enhanced": agent_results.get("knowledge_base_enhanced", False),
            },
            "sections": agent_results.get("sections", []),
            "project_file": agent_results.get("project_file", {}),
            "support_files": agent_results.get("support_files", []),
            "timestamp": agent_results.get("timestamp"),
            "processing_method": agent_results.get("processing_method"),
        })

        response = requests.post(
            agent_response_url,
            json={
                "agent": "Proposal Narrative Writer",
                "response": response_string,
            },
            headers=headers,
            timeout=30
        )

        if response.status_code in [200, 201]:
            print("✅ Proposal agent results posted to backend successfully")
        else:
            print(f"⚠️ Failed to post proposal agent results: {response.status_code} - {response.text[:200]}")

    except Exception as e:
        print(f"❌ Error posting proposal agent results: {str(e)}")

# ================================
# V2 API ROUTES FOR MVP
# ================================

# ================================
# LLAMAINDEX INTEGRATION FLAGS
# ================================
USE_LLAMAINDEX = os.getenv("USE_LLAMAINDEX", "true").lower() == "true"

@app.route("/api/v2/upload", methods=["POST", "OPTIONS"])
def upload_file_v2():
    """V2 Upload endpoint enhanced with Redis progress tracking"""
    if request.method == "OPTIONS":
        return "", 200
    
    # Enhanced: Use Redis for progress tracking when available
    if REDIS_AVAILABLE:
        print("🚀 Using Redis-enhanced threading approach")
    else:
        print("⚠️ Using memory-based progress tracking")
    if not FIREBASE_AVAILABLE:
        return jsonify({"error": "Firebase is not available"}), 503
        
    if not OPENAI_AVAILABLE:
        return jsonify({"error": "OpenAI embedding service is not available"}), 503
        
    save_path = None
    try:
        org_id = request.args.get("orgId")
        file_id = request.args.get("fileId")
        user_id = request.args.get("userId")
        
        print(f"📥 V2 Upload: fileId={file_id}, orgId={org_id}, userId={user_id}")
        
        if not org_id or not file_id or not user_id:
            return jsonify({"error": "orgId, fileId, and userId are required"}), 400

        if "file" not in request.files:
            return jsonify({"error": "No file provided"}), 400

        file = request.files["file"]
        filename = file.filename
        
        # Fix double extension issue (e.g., "file.pdf.pdf" -> "file.pdf")
        if filename.count('.') > 1:
            parts = filename.split('.')
            # Check if the last two parts are the same extension
            if len(parts) >= 2 and parts[-1].lower() == parts[-2].lower():
                filename = '.'.join(parts[:-1])  # Remove the duplicate extension
                print(f"🔧 Fixed double extension: {file.filename} -> {filename}")
        
        file_ext = filename.split(".")[-1].lower()

        # Accept all file types - AI will handle extraction
        print(f"📋 V2 Processing file type: {file_ext} (all types accepted with AI extraction)")

        # Notify Node.js backend: processing started
        notify_backend_status(file_id, user_id, 'processing', False)
        
        # Save file temporarily - use absolute path to prevent working directory issues
        save_path = os.path.abspath(os.path.join(UPLOAD_FOLDER, f"{file_id}_{filename}"))
        file.save(save_path)
        print(f"✅ V2 File saved: {save_path}")
        
        # Immediate verification that file was saved
        if os.path.exists(save_path):
            file_size = os.path.getsize(save_path)
            print(f"✅ File confirmed saved: {save_path} ({file_size} bytes)")
            
            # Create a backup copy with different name to prevent Flask cleanup
            backup_path = os.path.abspath(os.path.join(UPLOAD_FOLDER, f"backup_{file_id}_{filename}"))
            import shutil
            shutil.copy2(save_path, backup_path)
            print(f"✅ Backup created: {backup_path}")
            
        else:
            print(f"❌ WARNING: File not found immediately after save: {save_path}")
            return jsonify({"error": "File save failed"}), 500
        
        # Process file in background thread
        def process_file_v2_async():
            try:
                print(f"🔄 V2 Processing file: {filename}")
                
                # Use backup path for processing to avoid Flask interference
                processing_path = backup_path
                
                # Verify file still exists before processing
                if not os.path.exists(processing_path):
                    print(f"❌ Backup file missing: {processing_path}")
                    notify_backend_status(file_id, user_id, 'failed', False, "File not found")
                    return
                
                print(f"✅ Using backup file for processing: {processing_path}")
                
                # Extract text and create chunks
                chunks = parse_and_chunk(processing_path, file_ext, chunk_size=50, max_chunks=500)
                
                if not chunks:
                    print(f"❌ V2 No content extracted from {filename}")
                    notify_backend_status(file_id, user_id, 'failed', False, "No content extracted")
                    return
                
                # Notify: embedding phase started
                notify_backend_status(file_id, user_id, 'embedding', False)
                    
                # Choose processing method based on USE_LLAMAINDEX flag
                print(f"🔧 USE_LLAMAINDEX flag: {USE_LLAMAINDEX}")
                if USE_LLAMAINDEX:
                    print(f"🚀 V2 Processing with LlamaIndex...")
                    print(f"📄 File: {filename}, Path: {save_path}")
                    print(f"🆔 IDs - File: {file_id}, Org: {org_id}, User: {user_id}")
                    try:
                        print(f"📥 Importing LlamaIndex integration...")
                        from llamaindex_integration import process_file_with_llamaindex
                        
                        print(f"🔄 Calling LlamaIndex processor...")
                        result = process_file_with_llamaindex(
                            file_path=save_path,
                            file_id=file_id,
                            org_id=org_id,
                            user_id=user_id,
                            filename=filename
                        )
                        print(f"✅ LlamaIndex result: {result}")
                        
                        print(f"✅ V2 LlamaIndex processing complete: {result}")
                        
                        if result.get("success"):
                            # Notify Node.js backend: completed
                            notify_backend_status(file_id, user_id, 'completed', True)
                            
                            # Clean up temporary files
                            if 'save_path' in locals() and save_path and os.path.exists(save_path):
                                os.remove(save_path)
                                print(f"🗑️ V2 Cleaned up temporary file: {save_path}")
                            
                            return  # Exit successfully
                        else:
                            # Notify failure and fallback to traditional processing
                            print(f"❌ V2 LlamaIndex failed, falling back to traditional processing")
                            
                    except Exception as llamaindex_error:
                        print(f"❌ V2 LlamaIndex error: {str(llamaindex_error)}")
                        print(f"🔄 V2 Falling back to traditional processing...")
                
                # Traditional processing (fallback or when USE_LLAMAINDEX=false)
                print(f"🔄 V2 Using traditional processing...")
                
                # Generate embeddings
                print(f"🧠 V2 Generating embeddings for {len(chunks)} chunks")
                embeddings = embed_chunks(chunks, upload_id=file_id, org_id=org_id, filename=filename)
                
                if not embeddings:
                    print(f"❌ V2 Embedding generation failed for {filename}")
                    notify_backend_status(file_id, user_id, 'failed', False, "Embedding generation failed")
                    return
                
                print(f"🔄 V2 Starting Firestore storage for {len(embeddings)} embeddings...")
                
                # Store in Firestore using nested structure
                file_doc_ref = db.collection("document_embeddings").document(f"org-{org_id}").collection("files").document(file_id)
                
                print(f"📝 V2 Storing file metadata...")
                
                try:
                    file_doc_ref.set({
                        "filename": filename,
                        "file_id": file_id,
                        "file_type": file_ext,
                        "chunk_count": len(chunks),
                        "upload_timestamp": time.time(),
                        "created_at": datetime.now().isoformat(),
                        "org_id": org_id,
                        "user_id": user_id,
                        "status": "completed"
                    })
                    print(f"✅ V2 File metadata stored successfully")
                except Exception as metadata_error:
                    print(f"❌ V2 Failed to store file metadata: {str(metadata_error)}")
                    raise
                
                # Store embeddings in subcollection
                print(f"📦 V2 Starting to store {len(embeddings)} embedding chunks...")
                chunks_collection = file_doc_ref.collection("chunks")
                
                try:
                    # Use batched writes for better performance and reliability
                    batch_size = 5  # Smaller batches for better reliability
                    total_chunks = len(embeddings)
                    
                    for i in range(0, total_chunks, batch_size):
                        batch = db.batch()
                        end_idx = min(i + batch_size, total_chunks)
                        
                        print(f"🔹 V2 Storing batch {i//batch_size + 1} (chunks {i+1}-{end_idx})...")
                        
                        for j in range(i, end_idx):
                            chunk_ref = chunks_collection.document(f"chunk_{j}")
                            batch.set(chunk_ref, {
                                "text": chunks[j],
                                "embedding": embeddings[j],
                                "chunk_index": j,
                                "created_at": datetime.now().isoformat()
                            })
                        
                        # Commit the batch
                        batch.commit()
                        print(f"✓ V2 Batch committed - stored chunks {i+1}-{end_idx}")
                        
                        # Clean up batch object
                        del batch
                        import gc
                        gc.collect()
                    
                    print(f"✅ V2 All embedding chunks stored successfully using batched writes")
                except Exception as chunk_error:
                    print(f"❌ V2 Failed to store chunk {i}: {str(chunk_error)}")
                    raise
                
                print(f"✅ V2 Embeddings stored successfully for {filename}")
                
                # Notify Node.js backend: completed
                notify_backend_status(file_id, user_id, 'completed', True)
                
                # Clean up temporary files only after successful completion
                if 'save_path' in locals() and save_path and os.path.exists(save_path):
                    os.remove(save_path)
                    print(f"🗑️ V2 Cleaned up temporary file: {save_path}")
                
                return  # Exit successfully
                
            except Exception as e:
                print(f"❌ V2 Processing error for {filename}: {str(e)}")
                traceback.print_exc()
                notify_backend_status(file_id, user_id, 'failed', False, str(e))
                
                # Clean up temporary files on error too
                for cleanup_path in [save_path, backup_path]:
                    if cleanup_path and os.path.exists(cleanup_path):
                        os.remove(cleanup_path)
                        print(f"🗑️ V2 Cleaned up after error: {cleanup_path}")
        
        # Start processing in background
        processing_thread = Thread(target=process_file_v2_async, daemon=True)
        processing_thread.start()
        
        return jsonify({
            "success": True,
            "message": "File received and processing started",
            "fileId": file_id,
            "filename": filename,
            "status": "processing"
        }), 202

    except Exception as e:
        error_msg = str(e)
        print(f"❌ V2 Upload error: {error_msg}")
        traceback.print_exc()
        
        # Notify backend of failure
        if 'file_id' in locals() and 'user_id' in locals():
            notify_backend_status(file_id, user_id, 'failed', False, error_msg)
        
        return jsonify({"error": "Internal server error", "details": error_msg}), 500
    finally:
        # Clean up on immediate error
        if save_path and os.path.exists(save_path):
            try:
                os.remove(save_path)
            except:
                pass

@app.route("/api/v3/upload", methods=["POST", "OPTIONS"])
def upload_file_v3():
    """V3 Upload endpoint - Direct LlamaIndex processing with custom HF embeddings"""
    if request.method == "OPTIONS":
        return "", 200
    
    if not USE_LLAMAINDEX:
        return jsonify({"error": "V3 upload requires LlamaIndex (USE_LLAMAINDEX=true)"}), 503
        
    save_path = None
    try:
        org_id = request.args.get("orgId")
        file_id = request.args.get("fileId")
        user_id = request.args.get("userId")
        
        print(f"📥 V3 Upload: fileId={file_id}, orgId={org_id}, userId={user_id}")
        
        if not org_id or not file_id or not user_id:
            return jsonify({"error": "orgId, fileId, and userId are required"}), 400

        if "file" not in request.files:
            return jsonify({"error": "No file provided"}), 400

        file = request.files["file"]
        filename = file.filename
        
        # Fix double extension issue
        if filename.count('.') > 1:
            parts = filename.split('.')
            if len(parts) >= 2 and parts[-1].lower() == parts[-2].lower():
                filename = '.'.join(parts[:-1])
                print(f"🔧 Fixed double extension: {file.filename} -> {filename}")
        
        file_ext = filename.split(".")[-1].lower()
        print(f"📋 V3 Processing file type: {file_ext} (direct LlamaIndex processing)")

        # Notify backend: processing started
        notify_backend_status(file_id, user_id, 'processing', False, source="flask_ai_v3")
        
        # Save file temporarily
        timestamp = int(time.time())
        safe_filename = f"{file_id}_{filename}"
        save_path = os.path.join(UPLOAD_FOLDER, safe_filename)
        
        file.save(save_path)
        print(f"✅ V3 File saved: {save_path}")
        
        if not os.path.exists(save_path):
            return jsonify({"error": "Failed to save file"}), 500
            
        file_size = os.path.getsize(save_path)
        print(f"📏 File size: {file_size} bytes")
        
        # Notify backend: embedding phase started
        notify_backend_status(file_id, user_id, 'embedding', False, source="flask_ai_v3")
        
        # Start LlamaIndex processing in background
        print(f"🚀 V3 Starting LlamaIndex processing in background...")
        print(f"📄 File: {filename}, Path: {save_path}")
        print(f"🆔 IDs - File: {file_id}, Org: {org_id}, User: {user_id}")
        
        def process_file_v3_async():
            try:
                print(f"🚀 V3 Starting direct NeonDB processing...")
                from direct_neondb_storage import process_file_direct_storage
                
                result = process_file_direct_storage(
                    file_path=save_path,
                    file_id=file_id,
                    org_id=org_id,
                    user_id=user_id,
                    filename=filename
                )
                
                print(f"✅ V3 Direct storage result: {result}")
                
                if result.get("success"):
                    # Success - notify completion
                    notify_backend_status(file_id, user_id, 'completed', True, source="flask_ai_v3")
                    print(f"✅ V3 Processing complete: {result.get('chunks_stored', 0)} chunks stored")
                else:
                    # Direct storage failed
                    error_msg = result.get("error", "Unknown storage error")
                    print(f"❌ V3 Direct storage failed: {error_msg}")
                    notify_backend_status(file_id, user_id, 'failed', False, f"Direct storage failed: {error_msg}")
                    
            except Exception as e:
                print(f"❌ V3 Background processing error: {str(e)}")
                notify_backend_status(file_id, user_id, 'failed', False, str(e))
            finally:
                # Cleanup files
                if save_path and os.path.exists(save_path):
                    os.unlink(save_path)
                    print(f"🗑️ V3 Cleaned up: {save_path}")
        
        # Start background processing
        Thread(target=process_file_v3_async, daemon=True).start()
        
        # Return immediate response
        return jsonify({
            "message": "File upload successful, processing started with direct NeonDB storage",
            "file_id": file_id,
            "filename": filename,
            "status": "processing",
            "processing_method": "direct_neondb_v3",
            "embedding_model": "custom_hf_endpoint",
            "vector_dimension": os.getenv("VECTOR_DIM", "2560"),
            "note": "Using LlamaIndex parsing with direct database storage, bypassing vector store"
        }), 202
            
    except Exception as e:
        print(f"❌ V3 Upload error: {str(e)}")
        
        # Cleanup on error
        if save_path and os.path.exists(save_path):
            os.unlink(save_path)
            print(f"🗑️ V3 Cleaned up after error: {save_path}")
        
        # Notify backend of failure
        if 'user_id' in locals() and 'file_id' in locals():
            notify_backend_status(file_id, user_id, 'failed', False, str(e))
        
        return jsonify({"error": str(e)}), 500

@app.route("/api/v2/search", methods=["POST", "OPTIONS"])
def search_v2():
    """V2 Search endpoint with LlamaIndex integration"""
    if request.method == "OPTIONS":
        return "", 200
    
    try:
        data = request.get_json()
        query = data.get("query")
        org_id = data.get("orgId")
        top_k = data.get("top_k", 10)
        file_ids = data.get("fileIds")  # Optional
        
        if not query or not org_id:
            return jsonify({"error": "Missing required fields: query, orgId"}), 400
        
        if USE_LLAMAINDEX:
            print(f"🔍 V2 LlamaIndex Search: '{query}' in org: {org_id}")
            try:
                from llamaindex_integration import search_with_llamaindex
                
                results = search_with_llamaindex(
                    query=query,
                    org_id=org_id,
                    top_k=top_k,
                    file_ids=file_ids
                )
                
                return jsonify({
                    "query": query,
                    "results": results,
                    "total_found": len(results),
                    "search_type": "llamaindex_vector_search",
                    "storage": "neondb_postgresql"
                }), 200
                
            except Exception as llamaindex_error:
                print(f"❌ V2 LlamaIndex search error: {str(llamaindex_error)}")
                # Could fallback to traditional search here if implemented
                return jsonify({"error": f"LlamaIndex search failed: {str(llamaindex_error)}"}), 500
        else:
            # Traditional search would go here
            return jsonify({"error": "Traditional search not implemented. Enable LlamaIndex."}), 501
        
    except Exception as e:
        print(f"❌ V2 Search error: {str(e)}")
        return jsonify({"error": str(e)}), 500

@app.route("/api/v2/answer", methods=["POST", "OPTIONS"])
def generate_answer_v2():
    """V2 Answer generation endpoint with LlamaIndex + Qwen2"""
    if request.method == "OPTIONS":
        return "", 200
    
    try:
        data = request.get_json()
        query = data.get("query")
        org_id = data.get("orgId")
        file_ids = data.get("fileIds")  # Optional
        
        if not query or not org_id:
            return jsonify({"error": "Missing required fields: query, orgId"}), 400
        
        if USE_LLAMAINDEX:
            print(f"🤖 V2 LlamaIndex Answer Generation: '{query}' in org: {org_id}")
            try:
                from llamaindex_integration import get_llamaindex_processor
                
                processor = get_llamaindex_processor()
                result = processor.generate_answer(
                    query=query,
                    org_id=org_id,
                    file_ids=file_ids
                )
                
                return jsonify({
                    "query": query,
                    "answer": result["answer"],
                    "sources": result["sources"],
                    "confidence": result["confidence"],
                    "total_sources_found": result["total_sources_found"],
                    "model": "qwen2.5-7b-instruct",
                    "method": "llamaindex_rag",
                    "storage": "neondb_postgresql"
                }), 200
                
            except Exception as llamaindex_error:
                print(f"❌ V2 LlamaIndex answer error: {str(llamaindex_error)}")
                return jsonify({"error": f"LlamaIndex answer generation failed: {str(llamaindex_error)}"}), 500
        else:
            return jsonify({"error": "Answer generation requires LlamaIndex. Enable USE_LLAMAINDEX."}), 501
        
    except Exception as e:
        print(f"❌ V2 Answer error: {str(e)}")
        return jsonify({"error": str(e)}), 500

def notify_backend_status(file_id, user_id, status, embedding_complete, error=None, source="flask_ai"):
    """Notify Node.js backend of file processing status via webhook"""
    try:
        backend_api_url = os.environ.get("BACKEND_API_URL", "http://localhost:8080")
        webhook_url = f"{backend_api_url}/api/v2/files/webhook/ai-status"
        
        payload = {
            "fileId": file_id,
            "userId": user_id,
            "status": status,
            "embeddingComplete": embedding_complete,
            "timestamp": datetime.now().isoformat(),
            "source": source
        }
        
        if error:
            payload["error"] = error
            
        response = requests.post(
            webhook_url,
            json=payload,
            headers={"Content-Type": "application/json"},
            timeout=10
        )
        
        if response.status_code == 200:
            print(f"✅ V2 Status webhook sent: {file_id} → {status}")
        else:
            print(f"⚠️ V2 Webhook failed: {response.status_code}")
            
    except Exception as e:
        print(f"❌ V2 Webhook error: {str(e)}")


@app.route("/api/v2/delete-org-collections", methods=["DELETE", "OPTIONS"])
def delete_org_collections():
    """Delete all collections and documents for a specific organization"""
    if request.method == "OPTIONS":
        return "", 200
        
    if not FIREBASE_AVAILABLE:
        return jsonify({"error": "Firebase is not available"}), 503
        
    try:
        org_id = request.args.get("orgId")
        
        if not org_id:
            return jsonify({"error": "orgId is required"}), 400
            
        print(f"🗑️ Starting deletion of all collections for org: {org_id}")
        
        deleted_counts = {}
        total_deleted = 0
        
        # Define all collection patterns for this org
        collections_to_delete = [
            f"document_embeddings/org-{org_id}",
            f"org_rfp_support_embeddings/org-{org_id}",
            f"org_project_support_embeddings/org-{org_id}"
        ]
        
        for collection_path in collections_to_delete:
            try:
                print(f"🔍 Checking collection: {collection_path}")
                
                # Handle nested collection paths
                if "/" in collection_path:
                    parts = collection_path.split("/")
                    if len(parts) == 2:
                        # Format: collection/document
                        parent_collection = parts[0]
                        document_id = parts[1]
                        
                        # Get the document reference
                        doc_ref = db.collection(parent_collection).document(document_id)
                        
                        # Check if document exists
                        if doc_ref.get().exists:
                            print(f"📄 Found org document: {document_id}")
                            
                            # Delete all subcollections with nested cleanup
                            subcollection_names = ["files", "rfps", "projects"]  # Known subcollections
                            subcollection_count = 0
                            
                            for subcol_name in subcollection_names:
                                try:
                                    subcol_ref = doc_ref.collection(subcol_name)
                                    
                                    # Special handling for files collection - delete chunks first
                                    if subcol_name == "files":
                                        print(f"  🗂️ Processing files collection...")
                                        files_docs = list(subcol_ref.stream())
                                        
                                        for file_doc in files_docs:
                                            file_id = file_doc.id
                                            print(f"    📄 Deleting file: {file_id}")
                                            
                                            # Delete chunks subcollection first
                                            chunks_ref = file_doc.reference.collection("chunks")
                                            chunks_deleted = delete_collection(chunks_ref, batch_size=10)
                                            if chunks_deleted > 0:
                                                print(f"      ✅ Deleted {chunks_deleted} chunks")
                                                subcollection_count += chunks_deleted
                                            
                                            # Delete the file document
                                            file_doc.reference.delete()
                                            subcollection_count += 1
                                            print(f"      ✅ Deleted file document: {file_id}")
                                    
                                    # Special handling for rfps collection - delete files and chunks
                                    elif subcol_name == "rfps":
                                        print(f"  🗂️ Processing RFPs collection...")
                                        rfp_docs = list(subcol_ref.stream())
                                        
                                        for rfp_doc in rfp_docs:
                                            rfp_id = rfp_doc.id
                                            print(f"    📁 Deleting RFP: {rfp_id}")
                                            
                                            # Delete files subcollection in RFP
                                            rfp_files_ref = rfp_doc.reference.collection("files")
                                            rfp_files = list(rfp_files_ref.stream())
                                            
                                            for rfp_file in rfp_files:
                                                # Delete chunks in RFP file
                                                rfp_chunks_ref = rfp_file.reference.collection("chunks")
                                                rfp_chunks_deleted = delete_collection(rfp_chunks_ref, batch_size=10)
                                                if rfp_chunks_deleted > 0:
                                                    print(f"        ✅ Deleted {rfp_chunks_deleted} RFP file chunks")
                                                    subcollection_count += rfp_chunks_deleted
                                                
                                                # Delete RFP file
                                                rfp_file.reference.delete()
                                                subcollection_count += 1
                                            
                                            # Delete the RFP document
                                            rfp_doc.reference.delete()
                                            subcollection_count += 1
                                            print(f"      ✅ Deleted RFP document: {rfp_id}")
                                    
                                    else:
                                        # Standard deletion for other collections
                                        count = delete_collection(subcol_ref, batch_size=10)
                                        if count > 0:
                                            print(f"  ✅ Deleted {count} documents from {subcol_name}")
                                            subcollection_count += count
                                            
                                except Exception as subcol_error:
                                    print(f"  ⚠️ Error deleting subcollection {subcol_name}: {subcol_error}")
                            
                            # Delete the parent document
                            doc_ref.delete()
                            subcollection_count += 1
                            print(f"  ✅ Deleted parent document: {document_id}")
                            
                            deleted_counts[collection_path] = subcollection_count
                            total_deleted += subcollection_count
                        else:
                            print(f"  ℹ️ Document not found: {document_id}")
                            deleted_counts[collection_path] = 0
                else:
                    # Direct collection reference
                    collection_ref = db.collection(collection_path)
                    count = delete_collection(collection_ref, batch_size=10)
                    deleted_counts[collection_path] = count
                    total_deleted += count
                    
                    if count > 0:
                        print(f"✅ Deleted {count} documents from {collection_path}")
                    else:
                        print(f"ℹ️ No documents found in {collection_path}")
                        
            except Exception as col_error:
                print(f"❌ Error processing collection {collection_path}: {str(col_error)}")
                deleted_counts[collection_path] = f"Error: {str(col_error)}"
        
        print(f"🎉 Deletion completed for org {org_id}. Total documents deleted: {total_deleted}")
        
        return jsonify({
            "success": True,
            "message": f"Successfully deleted collections for organization {org_id}",
            "org_id": org_id,
            "total_deleted": total_deleted,
            "deleted_counts": deleted_counts,
            "collections_processed": collections_to_delete
        }), 200
        
    except Exception as e:
        error_msg = str(e)
        print(f"❌ Error deleting org collections: {error_msg}")
        return jsonify({
            "error": "Failed to delete org collections", 
            "details": error_msg
        }), 500

@app.route("/api/v2/delete-org-files", methods=["DELETE", "OPTIONS"])
def delete_org_files_only():
    """Delete only the files collections and their chunks for a specific organization"""
    if request.method == "OPTIONS":
        return "", 200
        
    if not FIREBASE_AVAILABLE:
        return jsonify({"error": "Firebase is not available"}), 503
        
    try:
        org_id = request.args.get("orgId")
        
        if not org_id:
            return jsonify({"error": "orgId is required"}), 400
            
        print(f"🗑️ Starting deletion of files collections for org: {org_id}")
        
        deleted_counts = {}
        total_deleted = 0
        
        # Target only the files collections
        collections_to_check = [
            f"document_embeddings/org-{org_id}",
            f"org_rfp_support_embeddings/org-{org_id}",
            f"org_project_support_embeddings/org-{org_id}"
        ]
        
        for collection_path in collections_to_check:
            try:
                print(f"🔍 Checking files in: {collection_path}")
                
                parts = collection_path.split("/")
                parent_collection = parts[0]
                document_id = parts[1]
                
                # Get the document reference
                doc_ref = db.collection(parent_collection).document(document_id)
                
                # Check if document exists
                if doc_ref.get().exists:
                    print(f"📄 Found org document: {document_id}")
                    
                    # Only delete files collection
                    files_ref = doc_ref.collection("files")
                    files_docs = list(files_ref.stream())
                    files_deleted = 0
                    
                    if files_docs:
                        print(f"  🗂️ Found {len(files_docs)} files to delete...")
                        
                        for file_doc in files_docs:
                            file_id = file_doc.id
                            print(f"    📄 Deleting file: {file_id}")
                            
                            # Delete chunks subcollection first
                            chunks_ref = file_doc.reference.collection("chunks")
                            chunks_deleted = delete_collection(chunks_ref, batch_size=10)
                            if chunks_deleted > 0:
                                print(f"      ✅ Deleted {chunks_deleted} chunks")
                                files_deleted += chunks_deleted
                            
                            # Delete the file document
                            file_doc.reference.delete()
                            files_deleted += 1
                            print(f"      ✅ Deleted file document: {file_id}")
                    else:
                        print(f"  ℹ️ No files found in {collection_path}")
                    
                    deleted_counts[collection_path] = files_deleted
                    total_deleted += files_deleted
                else:
                    print(f"  ℹ️ Document not found: {document_id}")
                    deleted_counts[collection_path] = 0
                        
            except Exception as col_error:
                print(f"❌ Error processing collection {collection_path}: {str(col_error)}")
                deleted_counts[collection_path] = f"Error: {str(col_error)}"
        
        print(f"🎉 Files deletion completed for org {org_id}. Total documents deleted: {total_deleted}")
        
        return jsonify({
            "success": True,
            "message": f"Successfully deleted files for organization {org_id}",
            "org_id": org_id,
            "total_deleted": total_deleted,
            "deleted_counts": deleted_counts,
            "collections_processed": collections_to_check
        }), 200
        
    except Exception as e:
        error_msg = str(e)
        print(f"❌ Error deleting org files: {error_msg}")
        return jsonify({
            "error": "Failed to delete org files", 
            "details": error_msg
        }), 500

@app.route("/api/v2/status", methods=["GET", "OPTIONS"])
def get_file_status_v2():
    """V2 Get file processing status"""
    if request.method == "OPTIONS":
        return "", 200
        
    try:
        file_id = request.args.get("fileId")
        org_id = request.args.get("orgId")
        
        if not file_id or not org_id:
            return jsonify({"error": "fileId and orgId are required"}), 400
            
        # Check if file exists in Firestore
        file_doc_ref = db.collection("document_embeddings").document(f"org-{org_id}").collection("files").document(file_id)
        file_doc = file_doc_ref.get()
        
        if file_doc.exists:
            data = file_doc.to_dict()
            return jsonify({
                "success": True,
                "fileId": file_id,
                "status": data.get("status", "unknown"),
                "filename": data.get("filename"),
                "chunkCount": data.get("chunk_count", 0),
                "embeddingComplete": data.get("status") == "completed",
                "createdAt": data.get("created_at")
            })
        else:
            return jsonify({
                "success": False,
                "message": "File not found in embedding database",
                "fileId": file_id
            }), 404
            
    except Exception as e:
        print(f"❌ V2 Status check error: {str(e)}")
        return jsonify({"error": "Internal server error"}), 500

@app.route("/process-project-support-embedding", methods=["POST", "OPTIONS"])
def process_project_support_embedding():
    """Process embeddings for project support documents using V3 approach (LlamaParse + HF + NeonDB)"""
    if request.method == "OPTIONS":
        return "", 200
        
    if not USE_LLAMAINDEX:
        return jsonify({"error": "Project support processing requires LlamaIndex (USE_LLAMAINDEX=true)"}), 503
        
    save_path = None
    try:
        org_id = request.args.get("orgId")
        project_id = request.args.get("projectId")
        file_id = request.args.get("fileId")
        upload_id = request.args.get("uploadId", str(uuid.uuid4()))
        user_id = request.args.get("userId")  # For WebSocket notifications
        
        if not org_id or not project_id or not file_id:
            return jsonify({"error": "orgId, projectId, and fileId are required"}), 400

        if "file" not in request.files:
            return jsonify({"error": "No file provided"}), 400

        file = request.files["file"]
        filename = file.filename
        file_ext = filename.split(".")[-1].lower()

        # Save file temporarily
        save_path = os.path.abspath(os.path.join(UPLOAD_FOLDER, f"{upload_id}_{filename}"))
        file.save(save_path)
        
        print(f"🚀 Processing project support file: {filename} (Project: {project_id})")
        
        # Send initial status
        if user_id:
            notify_backend_status(file_id, user_id, "processing", False)
        
        # Process file in background thread
        def process_project_support_file_async():
            try:
                # Send processing status
                if user_id:
                    notify_backend_status(file_id, user_id, "extracting", False)
                
                # Use V3 approach: LlamaParse + HF embeddings + NeonDB storage
                print(f"🚀 V3 Project Support: Starting LlamaParse processing...")
                from direct_neondb_storage import process_file_direct_storage_project_support
                
                result = process_file_direct_storage_project_support(
                    file_path=save_path,
                    file_id=file_id,
                    org_id=org_id,
                    user_id=user_id,
                    project_id=project_id,
                    filename=filename
                )
                
                print(f"✅ V3 Project Support result: {result}")
                
                if not result.get("success"):
                    error_msg = result.get("error", "Processing failed")
                    print(f"❌ V3 Project Support failed: {error_msg}")
                    if user_id:
                        notify_backend_status(file_id, user_id, "failed", False, error_msg)
                    return
                
                # V3 storage handled by the function above
                print(f"✅ Successfully processed project support document {filename} for Project: {project_id}")
                print(f"📊 Stored {result.get('chunks_stored', 0)} chunks in project_support_embeddings table")
                
                # Send completion status
                if user_id:
                    notify_backend_status(file_id, user_id, "completed", True)
                
            except Exception as e:
                error_msg = str(e)
                print(f"❌ Async project support processing error: {error_msg}")
                traceback.print_exc()
                
                # Send error status
                if user_id:
                    notify_backend_status(file_id, user_id, "failed", False, error_msg)
                
            finally:
                if save_path and os.path.exists(save_path):
                    try:
                        os.remove(save_path)
                    except Exception as e:
                        print(f"Error deleting file: {e}")
        
        # Start processing in background thread
        processing_thread = Thread(target=process_project_support_file_async)
        processing_thread.daemon = False
        processing_thread.start()
        
        return jsonify({
            "message": "Project support document embedding started",
            "file_id": file_id,
            "upload_id": upload_id,
            "project_id": project_id,
            "org_id": org_id,
            "file_type": file_ext,
            "document_type": "project_support"
        }), 202

    except Exception as e:
        error_msg = str(e)
        print(f"❌ Project Support Document Embedding Error: {error_msg}")
        return jsonify({"error": "Internal server error", "details": error_msg}), 500

# ================================
# AGENT RESPONSES API v2
# ================================

@app.route("/api/v2/agent-responses/run/<agent_run_id>", methods=["GET", "OPTIONS"])
def get_agent_responses_for_run(agent_run_id):
    """Get agent responses for a specific agent run"""
    if request.method == "OPTIONS":
        return "", 200
        
    try:
        # Get query parameters for filtering
        status = request.args.get('status')
        assigned_user = request.args.get('assignedUser')
        priority = request.args.get('priority')
        
        print(f"🤖 Getting agent responses for run: {agent_run_id}")
        
        # Make request to backend
        backend_api_url = os.getenv('BACKEND_API_URL', 'http://localhost:8080')
        url = f"{backend_api_url}/api/v2/agent-responses/run/{agent_run_id}"
        
        params = {}
        if status:
            params['status'] = status
        if assigned_user:
            params['assignedUser'] = assigned_user
        if priority:
            params['priority'] = priority
            
        headers = {}
        if 'Authorization' in request.headers:
            headers['Authorization'] = request.headers['Authorization']
        
        response = requests.get(url, params=params, headers=headers)
        
        if response.status_code == 200:
            return jsonify(response.json()), 200
        else:
            return jsonify({"error": "Failed to fetch agent responses"}), response.status_code
            
    except Exception as e:
        print(f"❌ Get agent responses error: {str(e)}")
        return jsonify({"error": "Internal server error"}), 500

@app.route("/api/v2/agent-responses/<response_id>", methods=["GET", "OPTIONS"])
def get_agent_response_details(response_id):
    """Get specific agent response details"""
    if request.method == "OPTIONS":
        return "", 200
        
    try:
        print(f"🤖 Getting agent response details: {response_id}")
        
        # Make request to backend
        backend_api_url = os.getenv('BACKEND_API_URL', 'http://localhost:8080')
        url = f"{backend_api_url}/api/v2/agent-responses/{response_id}"
        
        headers = {}
        if 'Authorization' in request.headers:
            headers['Authorization'] = request.headers['Authorization']
        
        response = requests.get(url, headers=headers)
        
        if response.status_code == 200:
            return jsonify(response.json()), 200
        else:
            return jsonify({"error": "Failed to fetch agent response details"}), response.status_code
            
    except Exception as e:
        print(f"❌ Get agent response details error: {str(e)}")
        return jsonify({"error": "Internal server error"}), 500

@app.route("/api/v2/agent-responses/<response_id>", methods=["PUT", "OPTIONS"])
def update_agent_response(response_id):
    """Update an agent response"""
    if request.method == "OPTIONS":
        return "", 200
        
    try:
        data = request.get_json()
        if not data:
            return jsonify({"error": "No data provided"}), 400
            
        print(f"🤖 Updating agent response: {response_id}")
        
        # Make request to backend
        backend_api_url = os.getenv('BACKEND_API_URL', 'http://localhost:8080')
        url = f"{backend_api_url}/api/v2/agent-responses/{response_id}"
        
        headers = {'Content-Type': 'application/json'}
        if 'Authorization' in request.headers:
            headers['Authorization'] = request.headers['Authorization']
        
        response = requests.put(url, json=data, headers=headers)
        
        if response.status_code == 200:
            return jsonify(response.json()), 200
        else:
            return jsonify({"error": "Failed to update agent response"}), response.status_code
            
    except Exception as e:
        print(f"❌ Update agent response error: {str(e)}")
        return jsonify({"error": "Internal server error"}), 500

@app.route("/api/v2/agent-responses/<response_id>/assign", methods=["POST", "OPTIONS"])
def assign_users_to_agent_response(response_id):
    """Assign users to an agent response"""
    if request.method == "OPTIONS":
        return "", 200
        
    try:
        data = request.get_json()
        if not data:
            return jsonify({"error": "No data provided"}), 400
            
        print(f"🤖 Assigning users to agent response: {response_id}")
        
        # Make request to backend
        backend_api_url = os.getenv('BACKEND_API_URL', 'http://localhost:8080')
        url = f"{backend_api_url}/api/v2/agent-responses/{response_id}/assign"
        
        headers = {'Content-Type': 'application/json'}
        if 'Authorization' in request.headers:
            headers['Authorization'] = request.headers['Authorization']
        
        response = requests.post(url, json=data, headers=headers)
        
        if response.status_code == 200:
            return jsonify(response.json()), 200
        else:
            return jsonify({"error": "Failed to assign users"}), response.status_code
            
    except Exception as e:
        print(f"❌ Assign users error: {str(e)}")
        return jsonify({"error": "Internal server error"}), 500

# ================================
# RESPONSE GENERATION API v2
# ================================

@app.route("/api/v2/generate-response-old", methods=["POST", "OPTIONS"])
def generate_response_v2_old():
    """Generate response based on project knowledge type (global or specific)"""
    if request.method == "OPTIONS":
        return "", 200
        
    if not FIREBASE_AVAILABLE:
        return jsonify({"error": "Firebase is not available"}), 503
        
    if not OPENAI_AVAILABLE:
        return jsonify({"error": "OpenAI embedding service is not available"}), 503
        
    if not VERTEX_AVAILABLE:
        return jsonify({"error": "AI generation service is not available"}), 503
        
    try:
        data = request.get_json(silent=True)
        if not data:
            return jsonify({"error": "No valid JSON input found"}), 400

        query = data.get("query")
        project_id = data.get("projectId")
        org_id = data.get("orgId")
        knowledge_base_option = data.get("knowledgeBaseOption", "global")  # "global" or "specific"
        rerank = data.get("rerank", True)  # Enable reranking by default
        enable_hybrid_search = data.get("hybridSearch", True)  # Enable hybrid search by default
        dense_weight = data.get("denseWeight", 0.7)  # Weight for semantic search
        sparse_weight = data.get("sparseWeight", 0.3)  # Weight for BM25 search
        enable_query_expansion = data.get("queryExpansion", True)  # Enable query expansion by default
        max_query_variations = data.get("maxQueryVariations", 2)  # Number of query variations
        context_type = data.get("contextType", "general")  # Context for query expansion
        
        if not query or not project_id or not org_id:
            return jsonify({"error": "Query, projectId, and orgId are required."}), 400

        print(f"🤖 Processing v2 response generation: '{query}' for Project: {project_id}, Org: {org_id}, Knowledge: {knowledge_base_option}")

        if knowledge_base_option == "global":
            # Use global knowledge base (like /chat endpoint)
            print("🌐 Using global knowledge base")
            
            # Get query embedding
            query_embedding = np.array(embed_query(query))
            retrieved_docs = []

            # Search organization-level documents (global knowledge base)
            org_files_ref = db.collection("document_embeddings").document(f"org-{org_id}").collection("files")
            org_files = org_files_ref.stream()
            
            # Process each organization file
            for file_doc in org_files:
                file_data = file_doc.to_dict()
                
                # Get chunks for this file
                chunks_ref = file_doc.reference.collection("chunks")
                chunks = chunks_ref.stream()
                
                # Process each chunk
                for chunk_doc in chunks:
                    chunk_data = chunk_doc.to_dict()
                    
                    # Convert to numpy array
                    chunk_embedding = np.array(chunk_data["embedding"])
                    
                    # Calculate cosine similarity
                    score = np.dot(query_embedding, chunk_embedding) / (
                        np.linalg.norm(query_embedding) * np.linalg.norm(chunk_embedding)
                    )
                    
                    if score >= 0.2:  # Similarity threshold
                        # Handle different chunk storage formats: "content" (v3.0.0) or "text" (v2)
                        chunk_content = chunk_data.get("content") or chunk_data.get("text", "")
                        retrieved_docs.append({
                            "content": chunk_content, 
                            "score": float(score),
                            "filename": file_data.get("filename", "Unknown"),
                            "file_id": file_data.get("file_id", file_doc.id),
                            "file_type": file_data.get("file_type", "unknown"),
                            "document_source": "organization",
                            "source_type": "global"
                        })

            # Apply enhanced search (global knowledge)
            if retrieved_docs:
                index_key = f"generate_global_{org_id}"
                
                if enable_query_expansion or enable_hybrid_search:
                    print(f"🚀 Applying enhanced search to {len(retrieved_docs)} documents (global knowledge)...")
                    
                    enhanced_candidates = enhanced_search(
                        query=query,
                        documents=retrieved_docs,
                        index_key=index_key,
                        context_type=context_type,
                        enable_query_expansion=enable_query_expansion,
                        max_query_variations=max_query_variations,
                        search_method="hybrid" if enable_hybrid_search else "semantic",
                        top_k=15,
                        dense_weight=dense_weight,
                        sparse_weight=sparse_weight
                    )
                    
                    if rerank and enhanced_candidates:
                        print(f"🔄 Applying reranking to {len(enhanced_candidates)} enhanced candidates...")
                        top_chunks = rerank_documents(query, enhanced_candidates, top_k=5)
                    else:
                        top_chunks = enhanced_candidates[:5]
                        
                elif rerank:
                    print(f"🔄 Applying reranking to {len(retrieved_docs)} documents (global knowledge)...")
                    # Get more candidates for reranking (top 15), then rerank to get top 5
                    similarity_candidates = sorted(retrieved_docs, key=lambda x: x["score"], reverse=True)[:15]
                    top_chunks = rerank_documents(query, similarity_candidates, top_k=5)
                else:
                    # Get top chunks by similarity
                    top_chunks = sorted(retrieved_docs, key=lambda x: x["score"], reverse=True)[:5]
            else:
                top_chunks = []

            if not top_chunks:
                return jsonify({
                    "query": query,
                    "project_id": project_id,
                    "org_id": org_id,
                    "knowledge_type": "global",
                    "answer": "No relevant information found in the global knowledge base.",
                    "source_files": []
                }), 200

            # Generate answer
            context_chunks = [doc["content"] for doc in top_chunks]
            answer = generate_answer_with_gcp(query, context_chunks, "")
            
            # Get unique source files
            source_files = []
            seen_files = set()
            for doc in top_chunks:
                file_key = f"{doc['filename']}_{doc['file_type']}"
                if file_key not in seen_files:
                    source_files.append({
                        "filename": doc["filename"],
                        "file_id": doc["file_id"],
                        "file_type": doc["file_type"],
                        "source": doc["document_source"]
                    })
                    seen_files.add(file_key)

            return jsonify({
                "query": query,
                "project_id": project_id,
                "org_id": org_id,
                "knowledge_type": "global",
                "answer": answer,
                "source_files": source_files,
                "retrieved_chunks": len(top_chunks)
            }), 200

        else:
            # Use specific project support documents
            print("📁 Using specific project support documents")
            
            # Get query embedding
            query_embedding = np.array(embed_query(query))

            # Get project support documents using new structure
            project_files_ref = (db.collection("org_project_support_embeddings")
                               .document(f"org-{org_id}")
                               .collection("projects")
                               .document(f"project-{project_id}")
                               .collection("files"))
            
            files = project_files_ref.stream()
            retrieved_docs = []
            
            # Process each file
            for file_doc in files:
                file_data = file_doc.to_dict()
                
                # Get chunks for this file
                chunks_ref = file_doc.reference.collection("chunks")
                chunks = chunks_ref.stream()
                
                # Process each chunk
                for chunk_doc in chunks:
                    chunk_data = chunk_doc.to_dict()
                    
                    # Convert to numpy array
                    chunk_embedding = np.array(chunk_data["embedding"])
                    
                    # Calculate cosine similarity
                    score = np.dot(query_embedding, chunk_embedding) / (
                        np.linalg.norm(query_embedding) * np.linalg.norm(chunk_embedding)
                    )
                    
                    if score >= 0.2:  # Similarity threshold
                        # Handle different chunk storage formats: "content" (v3.0.0) or "text" (v2)
                        chunk_content = chunk_data.get("content") or chunk_data.get("text", "")
                        retrieved_docs.append({
                            "content": chunk_content, 
                            "score": float(score),
                            "filename": file_data.get("filename", "Unknown"),
                            "file_id": file_data.get("file_id", file_doc.id),
                            "file_type": file_data.get("file_type", "unknown"),
                            "document_type": "support",
                            "processing_version": file_data.get("processing_version", "legacy")
                        })

            # Apply enhanced search (specific knowledge)
            if retrieved_docs:
                index_key = f"generate_specific_{org_id}_{project_id}"
                
                if enable_query_expansion or enable_hybrid_search:
                    print(f"🚀 Applying enhanced search to {len(retrieved_docs)} documents (specific knowledge)...")
                    
                    enhanced_candidates = enhanced_search(
                        query=query,
                        documents=retrieved_docs,
                        index_key=index_key,
                        context_type=context_type,
                        enable_query_expansion=enable_query_expansion,
                        max_query_variations=max_query_variations,
                        search_method="hybrid" if enable_hybrid_search else "semantic",
                        top_k=15,
                        dense_weight=dense_weight,
                        sparse_weight=sparse_weight
                    )
                    
                    if rerank and enhanced_candidates:
                        print(f"🔄 Applying reranking to {len(enhanced_candidates)} enhanced candidates...")
                        top_chunks = rerank_documents(query, enhanced_candidates, top_k=5)
                    else:
                        top_chunks = enhanced_candidates[:5]
                        
                elif rerank:
                    print(f"🔄 Applying reranking to {len(retrieved_docs)} documents (specific knowledge)...")
                    # Get more candidates for reranking (top 15), then rerank to get top 5
                    similarity_candidates = sorted(retrieved_docs, key=lambda x: x["score"], reverse=True)[:15]
                    top_chunks = rerank_documents(query, similarity_candidates, top_k=5)
                else:
                    # Get top chunks by similarity
                    top_chunks = sorted(retrieved_docs, key=lambda x: x["score"], reverse=True)[:5]
            else:
                top_chunks = []

            if not top_chunks:
                return jsonify({
                    "query": query,
                    "project_id": project_id,
                    "org_id": org_id,
                    "knowledge_type": "specific",
                    "answer": "No relevant information found in the project support documents.",
                    "source_files": [],
                    "storage_structure": f"org_project_support_embeddings/org-{org_id}/projects/project-{project_id}/files"
                }), 200

            # Generate answer
            context_chunks = [doc["content"] for doc in top_chunks]
            answer = generate_answer_with_gcp(query, context_chunks, "")
            
            # Get unique source files
            source_files = []
            seen_files = set()
            for doc in top_chunks:
                file_key = f"{doc['filename']}_{doc['file_type']}"
                if file_key not in seen_files:
                    source_files.append({
                        "filename": doc["filename"],
                        "file_id": doc["file_id"],
                        "file_type": doc["file_type"],
                        "document_type": doc["document_type"]
                    })
                    seen_files.add(file_key)

            return jsonify({
                "query": query,
                "project_id": project_id,
                "org_id": org_id,
                "knowledge_type": "specific",
                "answer": answer,
                "source_files": source_files,
                "retrieved_chunks": len(top_chunks),
                "storage_structure": f"org_project_support_embeddings/org-{org_id}/projects/project-{project_id}/files"
            }), 200

    except Exception as e:
        error_msg = str(e)
        print(f"❌ Generate Response v2 Error: {error_msg}")
        traceback.print_exc()
        return jsonify({"error": "Internal server error", "details": error_msg}), 500

# NEW MODERNIZED V2 GENERATE RESPONSE - USING V3 ARCHITECTURE
@app.route("/api/v2/generate-response", methods=["POST", "OPTIONS"])
def generate_response_v2():
    """Generate response using modern V3 architecture with NeonDB"""
    if request.method == "OPTIONS":
        return "", 200
        
    try:
        data = request.get_json(silent=True)
        if not data:
            return jsonify({"error": "No valid JSON input found"}), 400

        query = data.get("query")
        project_id = data.get("projectId")
        org_id = data.get("orgId")
        knowledge_base_option = data.get("knowledgeBaseOption", "global")  # "global" or "specific"
        conversation_history = data.get("conversationHistory", "")
        max_results = data.get("maxResults", 5)
        
        if not query or not org_id:
            return jsonify({"error": "Query and orgId are required."}), 400

        print(f"🤖 V2 Generate Response: '{query[:100]}...' for Org: {org_id}, Knowledge: {knowledge_base_option}")
        
        # Import modern retrieval and LLM systems
        from retrieval_system import search_documents
        from llm_integration import generate_rag_answer
        
        # Determine file filtering based on knowledge base option
        file_ids = None  # Default: search all files
        if knowledge_base_option == "specific" and project_id:
            print(f"📁 Using project-specific knowledge for project: {project_id}")
            # Note: This would require project-specific file mapping in NeonDB
            # For now, we'll search all files and add project context
        else:
            print("🌐 Using global knowledge base")
        
        # Search documents using modern V3 system
        print("🎯 Searching documents with modern retrieval system...")
        search_results = search_documents(
            query=query,
            org_id=org_id,
            file_ids=file_ids,
            top_k=max_results
        )
        
        # Build context from search results
        context_result = {
            "context": "\n\n".join([result["text"][:500] for result in search_results[:3]]),  # Top 3 results, 500 chars each
            "sources": list(set([result["file_id"] for result in search_results])),
            "total_chunks": len(search_results),
            "avg_similarity": sum([result["similarity_score"] for result in search_results]) / len(search_results) if search_results else 0,
            "chunks_metadata": search_results
        }
        
        if not search_results:
            return jsonify({
                "query": query,
                "project_id": project_id,
                "org_id": org_id,
                "knowledge_type": knowledge_base_option,
                "answer": "I couldn't find relevant information in the knowledge base to answer this question.",
                "source_files": [],
                "retrieved_chunks": 0,
                "search_results": []
            }), 200
        
        # Generate LLM answer using retrieved context
        print("🤖 Generating LLM answer...")
        
        # Prepare enhanced context with conversation history
        enhanced_context = context_result["context"]
        if conversation_history:
            enhanced_context = f"Previous conversation:\n{conversation_history}\n\nRelevant context:\n{context_result['context']}"
        
        llm_result = generate_rag_answer(
            query=query,
            context=enhanced_context,
            max_length=800,
            temperature=0.7
        )
        
        if llm_result["success"]:
            # Clean the LLM response
            raw_answer = llm_result["answer"]
            cleaned_answer = raw_answer
            
            # Clean up response formatting
            if "assistant" in raw_answer and raw_answer.count("assistant") > 0:
                parts = raw_answer.split("assistant")
                if len(parts) > 1:
                    cleaned_answer = parts[-1].strip()
            
            # Remove system/user prefixes
            for prefix in ["system\n", "user\n", "Answer:\n", "Answer:", "\nassistant\n"]:
                if cleaned_answer.startswith(prefix):
                    cleaned_answer = cleaned_answer[len(prefix):].strip()
            
            # Remove trailing artifacts
            for suffix in ["<|im_end|>", "<|endoftext|>"]:
                if cleaned_answer.endswith(suffix):
                    cleaned_answer = cleaned_answer[:-len(suffix)].strip()
            
            answer = cleaned_answer
        else:
            answer = f"I found {len(search_results)} relevant results, but couldn't generate a complete answer. Please check the search results below."
        
        # Collect source files
        source_files = []
        seen_files = set()
        for result in search_results:
            file_id = result["file_id"]
            if file_id not in seen_files:
                source_files.append({
                    "file_id": file_id,
                    "filename": result.get("filename", "unknown"),
                    "similarity_score": result["similarity_score"]
                })
                seen_files.add(file_id)
        
        # Return response in V2 format for compatibility
        response = {
            "query": query,
            "project_id": project_id,
            "org_id": org_id,
            "knowledge_type": knowledge_base_option,
            "answer": answer,
            "source_files": source_files,
            "retrieved_chunks": len(search_results),
            "search_results": search_results,
            "context_metadata": {
                "sources": context_result["sources"],
                "total_chunks": context_result["total_chunks"],
                "avg_similarity": context_result["avg_similarity"]
            },
            "llm_metadata": {
                "model": llm_result.get("model", "unknown"),
                "context_used": llm_result.get("context_used", False),
                "context_length": llm_result.get("context_length", 0)
            } if llm_result["success"] else {},
            "retrieval_method": "vector_similarity_v3",
            "processing_version": "v2.1.0_modernized"
        }
        
        print(f"✅ V2 Generate Response completed: {len(search_results)} results, answer length: {len(answer)}")
        return jsonify(response), 200

    except Exception as e:
        error_msg = str(e)
        print(f"❌ Generate Response v2 Error: {error_msg}")
        traceback.print_exc()
        return jsonify({"error": "Internal server error", "details": error_msg}), 500

# ================================
# QUESTIONNAIRE RESPONSE GENERATION API
# ================================

@app.route("/api/v2/questionnaires/generate-response", methods=["POST", "OPTIONS"])
def generate_questionnaire_response():
    """Generate contextually appropriate responses for questionnaire questions based on response type"""
    
    if request.method == "OPTIONS":
        return "", 200
    
    try:
        data = request.get_json()
        if not data:
            return jsonify({'success': False, 'error': 'No JSON data provided'}), 400
        
        questionnaire_id = data.get('questionnaireId')
        question_text = data.get('questionText')
        response_type = data.get('responseType', 'text_answer')
        choices = data.get('choices', [])
        context = data.get('context', {})
        knowledge_base_option = data.get('knowledgeBaseOption', 'global')
        org_id = data.get('orgId')
        project_id = data.get('projectId')
        rerank = data.get('rerank', True)  # Enable reranking by default
        enable_hybrid_search = data.get('hybridSearch', True)  # Enable hybrid search by default
        dense_weight = data.get('denseWeight', 0.7)  # Weight for semantic search
        sparse_weight = data.get('sparseWeight', 0.3)  # Weight for BM25 search
        enable_query_expansion = data.get('queryExpansion', True)  # Enable query expansion by default
        max_query_variations = data.get('maxQueryVariations', 2)  # Number of query variations
        context_type = data.get('contextType', "compliance")  # Context for query expansion (compliance for questionnaires)
        
        print(f"🎯 Questionnaire Generate: {question_text} ({response_type})")
        
        if not all([questionnaire_id, question_text, org_id, project_id]):
            return jsonify({
                'success': False,
                'error': 'Missing required fields'
            }), 400
        
        if not VERTEX_AVAILABLE:
            return jsonify({
                'success': False,
                'error': 'Vertex AI service not available'
            }), 500
        
        start_time = time.time()
        
        try:
            knowledge_context = get_knowledge_base_context(
                query=question_text,
                org_id=org_id,
                project_id=project_id,
                knowledge_base_option=knowledge_base_option,
                rerank=rerank,
                enable_hybrid_search=enable_hybrid_search,
                dense_weight=dense_weight,
                sparse_weight=sparse_weight,
                enable_query_expansion=enable_query_expansion,
                max_query_variations=max_query_variations,
                context_type=context_type
            )
        except:
            knowledge_context = "No specific context available."
        
        # Generate response based on type with strict formatting
        base_prompt = f"""You are assisting in responding to an RFP. You will be provided one question and the relevant document text that contains the information needed to answer it.

Important Rules:
• Only use the information provided in the relevant text.
• Do not assume, invent, or include any external facts.
• Maintain a professional and proposal-ready tone.
• Follow the response format exactly as specified for each answer type.

Relevant Text:
{knowledge_context}

Question: {question_text}

"""

        if response_type == 'yes_no':
            prompt = base_prompt + "Respond with ONLY ONE WORD: either 'Yes' or 'No'. Nothing else."
        elif response_type == 'true_false':
            prompt = base_prompt + "Respond with ONLY ONE WORD: either 'True' or 'False'. Nothing else."
        elif response_type == 'single_choice' and choices:
            prompt = base_prompt + f"Available options: {', '.join(choices)}\n\nRespond with ONLY the exact option text from the list above. Nothing else."
        elif response_type == 'multi_choice' and choices:
            prompt = base_prompt + f"Available options: {', '.join(choices)}\n\nRespond with ONLY the selected option names separated by commas. Example: 'Option1, Option3'. Nothing else."
        elif response_type == 'completed_incomplete':
            prompt = base_prompt + "Respond with ONLY ONE WORD: either 'Complete' or 'Incomplete'. Nothing else."
        elif response_type == 'long_answer':
            prompt = base_prompt + "Provide a comprehensive detailed answer with examples, background information, and specific details. Use multiple paragraphs if needed."
        elif response_type == 'text_answer':
            prompt = base_prompt + "Provide a detailed answer with examples, background information, and specific details. Use multiple paragraphs if needed."
        else:  # short_answer or default
            prompt = base_prompt + "Provide a concise, professional answer in 1-2 sentences."
        
        # Use Vertex AI Gemini with very low temperature
        from vertexai.generative_models import GenerativeModel, GenerationConfig
        model = GenerativeModel("gemini-2.0-flash")
        response = model.generate_content(
            prompt,
            generation_config=GenerationConfig(
                temperature=0.1,
                max_output_tokens=500,
                top_p=0.8,
                top_k=10
            )
        )
        
        generated_text = response.text.strip()
        
        return jsonify({
            'success': True,
            'generatedResponse': generated_text,
            'responseType': response_type,
            'confidence': 0.9,
            'knowledgeSource': knowledge_base_option,
            'processingTime': round(time.time() - start_time, 2),
            'timestamp': datetime.now(timezone.utc).isoformat() + 'Z'
        })
        
    except Exception as e:
        print(f"❌ Error: {str(e)}")
        return jsonify({
            'success': False,
            'error': str(e)
        }), 500

# ================================
# SPECIALIZED XLSX/CSV PROCESSING API
# ================================

@app.route("/api/v2/analyze-structured-document", methods=["POST", "OPTIONS"])
def analyze_structured_document():
    """
    Analyze CSV/XLSX files and identify question cells that need user responses.
    Returns structure with sheet names, questions, cell locations, and analysis.
    """
    if request.method == "OPTIONS":
        return "", 200
        
    if not VERTEX_AVAILABLE:
        return jsonify({"error": "Vertex AI service is not available"}), 503
    
    try:
        # Check if file is uploaded
        if 'file' not in request.files:
            return jsonify({"error": "No file uploaded"}), 400
            
        file = request.files['file']
        if file.filename == '':
            return jsonify({"error": "No file selected"}), 400
        
        # Validate file type
        file_ext = file.filename.lower().split('.')[-1] if '.' in file.filename else ''
        if file_ext not in ['csv', 'xlsx', 'xls']:
            return jsonify({"error": "Only CSV, XLSX, and XLS files are supported"}), 400
        
        print(f"📊 Analyzing structured document: {file.filename}")
        
        # Save uploaded file temporarily
        temp_filename = f"temp_{uuid.uuid4().hex[:8]}_{file.filename}"
        temp_filepath = os.path.join(UPLOAD_FOLDER, temp_filename)
        file.save(temp_filepath)
        
        # Process file based on type
        if file_ext == 'csv':
            analysis_result = analyze_csv_structure(temp_filepath, file.filename)
        else:  # xlsx or xls
            analysis_result = analyze_xlsx_structure(temp_filepath, file.filename)
        
        # Clean up temporary file
        try:
            os.remove(temp_filepath)
        except:
            pass
        
        return jsonify({
            "success": True,
            "filename": file.filename,
            "analysis": analysis_result,
            "timestamp": datetime.now().isoformat()
        }), 200
        
    except Exception as e:
        print(f"❌ Document analysis error: {str(e)}")
        return jsonify({
            "success": False,
            "error": f"Analysis failed: {str(e)}"
        }), 500

def analyze_xlsx_structure(filepath, filename):
    """Analyze XLSX structure and identify question cells"""
    try:
        import pandas as pd
        
        excel_file = pd.ExcelFile(filepath)
        sheet_names = excel_file.sheet_names
        
        analysis_result = {
            "filename": filename,
            "total_sheets": len(sheet_names),
            "sheets": []
        }
        
        # Get document context from AI
        document_context = analyze_document_with_direct_ai(filepath, filename, [])
        
        for sheet_name in sheet_names:
            df = pd.read_excel(filepath, sheet_name=sheet_name)
            
            # Identify question cells
            question_cells = identify_question_cells_detailed(df, sheet_name)
            
            sheet_analysis = {
                "sheet_name": sheet_name,
                "rows": len(df),
                "columns": len(df.columns),
                "column_names": df.columns.tolist(),
                "questions_found": len(question_cells),
                "questions": question_cells
            }
            
            analysis_result["sheets"].append(sheet_analysis)
        
        analysis_result["document_analysis"] = document_context.get("analysis", "")
        
        return analysis_result
        
    except Exception as e:
        raise Exception(f"XLSX analysis error: {str(e)}")

def analyze_csv_structure(filepath, filename):
    """Analyze CSV structure and identify question cells"""
    try:
        import pandas as pd
        
        df = pd.read_csv(filepath)
        
        # Get document context from AI  
        document_context = analyze_document_with_direct_ai(filepath, filename, [])
        
        # Identify question cells
        question_cells = identify_question_cells_detailed(df, "Sheet1")
        
        analysis_result = {
            "filename": filename,
            "total_sheets": 1,
            "sheets": [{
                "sheet_name": "Sheet1",
                "rows": len(df),
                "columns": len(df.columns),
                "column_names": df.columns.tolist(),
                "questions_found": len(question_cells),
                "questions": question_cells
            }],
            "document_analysis": document_context.get("analysis", "")
        }
        
        return analysis_result
        
    except Exception as e:
        raise Exception(f"CSV analysis error: {str(e)}")

def identify_question_cells_detailed(df, sheet_name):
    """
    Identify question cells and return detailed information about each question
    and where responses should be placed
    """
    question_cells = []
    
    # Question patterns
    question_patterns = [
        r'\?$',  # Ends with question mark
        r'^\s*(?:what|how|when|where|why|which|who|describe|explain|provide|list)',
        r'(?:requirement|required|needed|must|should|shall)',
        r'(?:fill\s+in|complete|answer|respond|provide\s+details)',
        r'(?:TBD|TBC|TODO|PENDING)',
        r'(?:enter|input|specify|indicate|state)',
    ]
    
    for row_idx, row in df.iterrows():
        for col_idx, cell_value in enumerate(row):
            cell_str = str(cell_value).strip() if not pd.isna(cell_value) else ""
            
            # Check if this cell contains a question
            is_question = any(re.search(pattern, cell_str.lower()) for pattern in question_patterns) if cell_str else False
            
            if is_question and len(cell_str) > 10:  # Filter out very short matches
                
                # Find the best response cell location
                response_location = find_response_cell_location(df, row_idx, col_idx)
                
                question_info = {
                    "question_text": cell_str,
                    "question_cell": {
                        "row": row_idx + 1,  # 1-based indexing
                        "column": df.columns[col_idx],
                        "column_index": col_idx + 1,  # 1-based indexing
                        "cell_reference": f"{chr(65 + col_idx)}{row_idx + 1}"
                    },
                    "response_cell": response_location,
                    "question_type": categorize_question(cell_str),
                    "sheet_name": sheet_name
                }
                
                question_cells.append(question_info)
    
    return question_cells

def find_response_cell_location(df, question_row, question_col):
    """Find the best location for the response to a question"""
    
    # Strategy 1: Check right cell (same row, next column)
    if question_col + 1 < len(df.columns):
        right_cell_value = df.iloc[question_row, question_col + 1]
        if pd.isna(right_cell_value) or str(right_cell_value).strip() == "":
            return {
                "row": question_row + 1,
                "column": df.columns[question_col + 1],
                "column_index": question_col + 2,
                "cell_reference": f"{chr(65 + question_col + 1)}{question_row + 1}",
                "location_type": "right_adjacent"
            }
    
    # Strategy 2: Check cell below (next row, same column)
    if question_row + 1 < len(df):
        below_cell_value = df.iloc[question_row + 1, question_col]
        if pd.isna(below_cell_value) or str(below_cell_value).strip() == "":
            return {
                "row": question_row + 2,
                "column": df.columns[question_col],
                "column_index": question_col + 1,
                "cell_reference": f"{chr(65 + question_col)}{question_row + 2}",
                "location_type": "below_adjacent"
            }
    
    # Strategy 3: Check diagonal (next row, next column)
    if question_row + 1 < len(df) and question_col + 1 < len(df.columns):
        diagonal_cell_value = df.iloc[question_row + 1, question_col + 1]
        if pd.isna(diagonal_cell_value) or str(diagonal_cell_value).strip() == "":
            return {
                "row": question_row + 2,
                "column": df.columns[question_col + 1],
                "column_index": question_col + 2,
                "cell_reference": f"{chr(65 + question_col + 1)}{question_row + 2}",
                "location_type": "diagonal_adjacent"
            }
    
    # Default: suggest right adjacent even if not empty
    if question_col + 1 < len(df.columns):
        return {
            "row": question_row + 1,
            "column": df.columns[question_col + 1],
            "column_index": question_col + 2,
            "cell_reference": f"{chr(65 + question_col + 1)}{question_row + 1}",
            "location_type": "right_suggested"
        }
    
    return None

def categorize_question(question_text):
    """Categorize the type of question based on content"""
    question_lower = question_text.lower()
    
    if any(word in question_lower for word in ['date', 'when', 'timeline', 'deadline']):
        return "date"
    elif any(word in question_lower for word in ['number', 'count', 'quantity', 'amount', 'how many']):
        return "numeric"
    elif any(word in question_lower for word in ['yes', 'no', 'true', 'false', 'confirm']):
        return "boolean"
    elif any(word in question_lower for word in ['describe', 'explain', 'detail', 'how']):
        return "descriptive"
    elif any(word in question_lower for word in ['list', 'enumerate', 'items']):
        return "list"
    else:
        return "general"

@app.route("/api/v2/process-structured-document", methods=["POST", "OPTIONS"])
def process_structured_document():
    """
    Process CSV/XLSX files with multiple sheets, identify empty cells adjacent to questions,
    use direct Vertex AI file upload for analysis, and reconstruct document with answers.
    """
    if request.method == "OPTIONS":
        return "", 200
        
    if not VERTEX_AVAILABLE:
        return jsonify({"error": "Vertex AI service is not available"}), 503
    
    try:
        # Check if file is uploaded
        if 'file' not in request.files:
            return jsonify({"error": "No file uploaded"}), 400
            
        file = request.files['file']
        if file.filename == '':
            return jsonify({"error": "No file selected"}), 400
        
        # Validate file type
        file_ext = file.filename.lower().split('.')[-1] if '.' in file.filename else ''
        if file_ext not in ['csv', 'xlsx', 'xls']:
            return jsonify({"error": "Only CSV, XLSX, and XLS files are supported"}), 400
        
        print(f"📊 Processing structured document: {file.filename}")
        print(f"🤖 Using direct Vertex AI analysis with embeddings")
        
        # Save uploaded file temporarily
        file_id = str(uuid.uuid4())
        temp_filename = f"{file_id}_{file.filename}"
        temp_filepath = os.path.join(UPLOAD_FOLDER, temp_filename)
        file.save(temp_filepath)
        
        # Process the file based on type
        if file_ext == 'csv':
            analysis_result = process_csv_with_direct_ai(temp_filepath, file.filename)
        else:  # xlsx or xls
            analysis_result = process_xlsx_with_direct_ai(temp_filepath, file.filename)
        
        # Clean up temporary file
        if os.path.exists(temp_filepath):
            os.remove(temp_filepath)
        
        return jsonify({
            "file_id": file_id,
            "filename": file.filename,
            "file_type": file_ext,
            "processing_method": "direct_ai_upload",
            **analysis_result
        }), 200
        
    except Exception as e:
        error_msg = str(e)
        print(f"❌ Structured Document Processing Error: {error_msg}")
        traceback.print_exc()
        return jsonify({"error": "Internal server error", "details": error_msg}), 500

def process_csv_with_direct_ai(filepath, filename):
    """Process CSV file with direct AI upload and focus on empty cells adjacent to questions"""
    try:
        # Read CSV file
        df = pd.read_csv(filepath)
        
        print(f"📊 CSV Analysis: {len(df)} rows, {len(df.columns)} columns")
        
        # Identify empty cells adjacent to questions/requirements
        target_cells = identify_empty_answer_cells(df, "main")
        
        # Upload file directly to AI and get document analysis
        document_context = analyze_document_with_direct_ai(filepath, filename, target_cells)
        
        # Generate responses for each target cell
        filled_responses = generate_responses_for_target_cells(target_cells, document_context, df, "main")
        
        # Create processed version
        processed_df = fill_target_cells_with_responses(df, filled_responses)
        
        # Save processed file
        output_filename = f"processed_{uuid.uuid4().hex[:8]}_{filename}"
        output_path = os.path.join(UPLOAD_FOLDER, output_filename)
        processed_df.to_csv(output_path, index=False)
        
        return {
            "total_rows": len(df),
            "total_columns": len(df.columns),
            "columns": df.columns.tolist(),
            "sheets_analyzed": [{"sheet_name": "main", "rows": len(df), "columns": len(df.columns)}],
            "target_cells_identified": len(target_cells),
            "cells_filled": len(filled_responses),
            "output_file": output_filename,
            "cell_analysis": target_cells[:10],
            "responses": filled_responses[:10]
        }
        
    except Exception as e:
        raise Exception(f"CSV processing error: {str(e)}")

def process_xlsx_with_direct_ai(filepath, filename):
    """Process XLSX file with multiple sheets using direct AI analysis"""
    try:
        # Read all sheets
        excel_file = pd.ExcelFile(filepath)
        sheet_names = excel_file.sheet_names
        
        print(f"📊 XLSX Analysis: {len(sheet_names)} sheets")
        
        analysis_result = {
            "total_sheets": len(sheet_names),
            "sheets_analyzed": []
        }
        
        all_target_cells = []
        all_responses = []
        
        # Upload file directly to AI and get document analysis
        document_context = analyze_document_with_direct_ai(filepath, filename, [])
        
        # Process each sheet
        for sheet_name in sheet_names:
            df = pd.read_excel(filepath, sheet_name=sheet_name)
            
            sheet_info = {
                "sheet_name": sheet_name,
                "rows": len(df),
                "columns": len(df.columns),
                "column_names": df.columns.tolist()
            }
            analysis_result["sheets_analyzed"].append(sheet_info)
            
            # Identify empty cells adjacent to questions
            target_cells = identify_empty_answer_cells(df, sheet_name)
            all_target_cells.extend(target_cells)
            
            # Generate responses for this sheet
            sheet_responses = generate_responses_for_target_cells(target_cells, document_context, df, sheet_name)
            all_responses.extend(sheet_responses)
        
        # Save processed file with all sheets
        output_filename = f"processed_{uuid.uuid4().hex[:8]}_{filename}"
        output_path = os.path.join(UPLOAD_FOLDER, output_filename)
        
        with pd.ExcelWriter(output_path, engine='openpyxl') as writer:
            for sheet_name in sheet_names:
                df = pd.read_excel(filepath, sheet_name=sheet_name)
                
                # Apply responses to this sheet
                sheet_responses = [r for r in all_responses if r.get('sheet_name') == sheet_name]
                processed_df = fill_target_cells_with_responses(df, sheet_responses)
                
                processed_df.to_excel(writer, sheet_name=sheet_name, index=False)
        
        analysis_result.update({
            "target_cells_identified": len(all_target_cells),
            "cells_filled": len(all_responses),
            "output_file": output_filename,
            "cell_analysis": all_target_cells[:10],
            "responses": all_responses[:10]
        })
        
        return analysis_result
        
    except Exception as e:
        raise Exception(f"XLSX processing error: {str(e)}")

def identify_empty_answer_cells(df, sheet_name):
    """
    Focus on empty cells that are adjacent to questions/requirements.
    Look for empty cells next to questions based on column headings.
    """
    target_cells = []
    
    # Question patterns to identify question cells
    question_patterns = [
        r'\?$',  # Ends with question mark
        r'^\s*(?:what|how|when|where|why|which|who|describe|explain|provide|list)',
        r'(?:requirement|required|needed|must|should|shall)',
        r'(?:fill\s+in|complete|answer|respond|provide\s+details)',
        r'(?:TBD|TBC|TODO|PENDING)',
    ]
    
    # Iterate through each cell
    for row_idx, row in df.iterrows():
        for col_idx, cell_value in enumerate(row):
            cell_str = str(cell_value).strip() if not pd.isna(cell_value) else ""
            
            # Check if this cell contains a question/requirement
            is_question = any(re.search(pattern, cell_str.lower()) for pattern in question_patterns) if cell_str else False
            
            if is_question:
                # Look for empty cells adjacent to this question cell
                # Check right cell (same row, next column)
                if col_idx + 1 < len(df.columns):
                    next_cell_value = df.iloc[row_idx, col_idx + 1]
                    if pd.isna(next_cell_value) or str(next_cell_value).strip() == "":
                        # This is an empty cell next to a question - target for answering
                        column_header = df.columns[col_idx + 1] if col_idx + 1 < len(df.columns) else f"Column_{col_idx + 1}"
                        
                        target_cells.append({
                            "sheet_name": sheet_name,
                            "target_row": row_idx,
                            "target_column": col_idx + 1,
                            "target_column_name": column_header,
                            "question_row": row_idx,
                            "question_column": col_idx,
                            "question_content": cell_str,
                            "context": {
                                "column_header": column_header,
                                "row_context": get_row_context(df, row_idx),
                                "question_type": classify_question_type(cell_str)
                            }
                        })
                
                # Also check cells below the question if they're in an "answer" column
                column_name = df.columns[col_idx].lower() if col_idx < len(df.columns) else ""
                if any(word in column_name for word in ["question", "requirement"]):
                    # Look for answer column nearby
                    for check_col in range(col_idx + 1, min(col_idx + 3, len(df.columns))):
                        check_col_name = df.columns[check_col].lower()
                        if any(word in check_col_name for word in ["answer", "response", "value", "detail"]):
                            check_cell_value = df.iloc[row_idx, check_col]
                            if pd.isna(check_cell_value) or str(check_cell_value).strip() == "":
                                target_cells.append({
                                    "sheet_name": sheet_name,
                                    "target_row": row_idx,
                                    "target_column": check_col,
                                    "target_column_name": df.columns[check_col],
                                    "question_row": row_idx,
                                    "question_column": col_idx,
                                    "question_content": cell_str,
                                    "context": {
                                        "column_header": df.columns[check_col],
                                        "row_context": get_row_context(df, row_idx),
                                        "question_type": classify_question_type(cell_str)
                                    }
                                })
    
    print(f"🎯 Found {len(target_cells)} empty cells adjacent to questions in {sheet_name}")
    return target_cells

def get_row_context(df, row_idx):
    """Get context from the entire row"""
    row_context = []
    row_data = df.iloc[row_idx]
    
    for col_idx, value in enumerate(row_data):
        if not pd.isna(value) and str(value).strip():
            col_name = df.columns[col_idx] if col_idx < len(df.columns) else f"Column_{col_idx}"
            row_context.append(f"{col_name}: {str(value)[:100]}")
    
    return " | ".join(row_context)

def classify_question_type(content):
    """Classify the type of question for better response generation"""
    content_lower = content.lower()
    
    if any(word in content_lower for word in ["what", "describe", "explain"]):
        return "descriptive"
    elif any(word in content_lower for word in ["how", "process", "procedure"]):
        return "procedural"
    elif any(word in content_lower for word in ["when", "date", "time"]):
        return "temporal"
    elif any(word in content_lower for word in ["where", "location"]):
        return "locational"
    elif any(word in content_lower for word in ["why", "reason"]):
        return "causal"
    elif any(word in content_lower for word in ["requirement", "must", "shall", "required"]):
        return "requirement"
    elif content_lower.endswith("?"):
        return "direct_question"
    else:
        return "general"

def analyze_document_with_direct_ai(filepath, filename, target_cells):
    """
    Upload document directly to Vertex AI and get comprehensive analysis
    using embeddings and AI understanding
    """
    try:
        from vertexai.generative_models import GenerativeModel, Part, SafetySetting
        import pandas as pd
        
        print(f"🤖 Processing {filename} for Vertex AI analysis...")
        
        # Determine file extension
        file_ext = filename.lower().split('.')[-1] if '.' in filename else ''
        
        # For XLSX/XLS files, convert to text format since Vertex AI doesn't support these MIME types
        if file_ext in ['xlsx', 'xls']:
            print(f"📊 Converting {file_ext.upper()} to text format for AI analysis...")
            
            # Read Excel file and convert all sheets to text
            excel_file = pd.ExcelFile(filepath)
            text_content = f"EXCEL FILE: {filename}\n{'='*50}\n\n"
            
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(filepath, sheet_name=sheet_name)
                text_content += f"SHEET: {sheet_name}\n{'-'*30}\n"
                text_content += f"Columns: {', '.join(df.columns.tolist())}\n\n"
                
                # Convert DataFrame to readable text format
                for idx, row in df.iterrows():
                    row_text = []
                    for col_name, value in row.items():
                        if pd.notna(value) and str(value).strip():
                            row_text.append(f"{col_name}: {value}")
                    if row_text:
                        text_content += f"Row {idx + 1}: {' | '.join(row_text)}\n"
                
                text_content += "\n\n"
            
            # Create text part for AI analysis
            file_part = Part.from_text(text_content)
            
        elif file_ext == 'csv':
            # For CSV files, we can upload directly
            with open(filepath, "rb") as f:
                file_data = f.read()
            file_part = Part.from_data(data=file_data, mime_type='text/csv')
            
        else:
            # For other file types, try direct upload
            with open(filepath, "rb") as f:
                file_data = f.read()
            mime_type_map = {
                'txt': 'text/plain',
                'pdf': 'application/pdf',
                'docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
            }
            mime_type = mime_type_map.get(file_ext, 'application/octet-stream')
            file_part = Part.from_data(data=file_data, mime_type=mime_type)
        
        # Create comprehensive analysis prompt
        analysis_prompt = f"""
Analyze this {file_ext.upper()} document comprehensively. This document contains questions, requirements, or forms that need to be filled out.

Please provide:

1. DOCUMENT STRUCTURE ANALYSIS:
   - Overview of the document's purpose and content
   - Key sections or sheets identified
   - Types of questions or requirements present

2. CONTENT UNDERSTANDING:
   - Main topics and themes covered
   - Important information that could be used to answer questions
   - Key data points, requirements, or specifications mentioned

3. CONTEXTUAL INFORMATION:
   - Business context or domain this document relates to
   - Any specific terminology or concepts used
   - Relationships between different sections

4. ANSWERABLE CONTENT:
   - What types of questions this document could help answer
   - Key facts, figures, procedures, or requirements that are explicitly stated
   - Any implicit information that could be inferred

Focus on extracting factual information that could be used to answer questions or fill in requirements found in forms or questionnaires.
"""
        
        model = GenerativeModel("gemini-2.0-flash")
        
        # Generate comprehensive document analysis
        response = model.generate_content(
            [file_part, Part.from_text(analysis_prompt)],
            generation_config={
                "max_output_tokens": 8192,
                "temperature": 0.3,  # Lower temperature for more factual analysis
                "top_p": 0.95,
            },
            safety_settings=[
                SafetySetting(
                    category=SafetySetting.HarmCategory.HARM_CATEGORY_HATE_SPEECH,
                    threshold=SafetySetting.HarmBlockThreshold.OFF
                ),
                SafetySetting(
                    category=SafetySetting.HarmCategory.HARM_CATEGORY_DANGEROUS_CONTENT,
                    threshold=SafetySetting.HarmBlockThreshold.OFF
                ),
                SafetySetting(
                    category=SafetySetting.HarmCategory.HARM_CATEGORY_SEXUALLY_EXPLICIT,
                    threshold=SafetySetting.HarmBlockThreshold.OFF
                ),
                SafetySetting(
                    category=SafetySetting.HarmCategory.HARM_CATEGORY_HARASSMENT,
                    threshold=SafetySetting.HarmBlockThreshold.OFF
                ),
            ],
            stream=False
        )
        
        document_analysis = response.text.strip()
        print(f"✅ Document analysis completed: {len(document_analysis)} characters")
        
        return {
            "filename": filename,
            "analysis": document_analysis,
            "target_cells_context": len(target_cells),
            "processing_method": "direct_vertex_ai_upload"
        }
        
    except Exception as e:
        print(f"❌ Direct AI analysis error: {str(e)}")
        return {
            "filename": filename,
            "analysis": f"Error analyzing document: {str(e)}",
            "target_cells_context": len(target_cells),
            "processing_method": "error"
        }

def generate_responses_for_target_cells(target_cells, document_context, df, sheet_name):
    """
    Generate AI responses for each target cell using the document context
    """
    filled_responses = []
    
    for cell in target_cells:
        try:
            from vertexai.generative_models import GenerativeModel, Part, SafetySetting
            
            # Create specific query for this cell
            question_content = cell["question_content"]
            column_header = cell["target_column_name"]
            row_context = cell["context"]["row_context"]
            question_type = cell["context"]["question_type"]
            
            # Build contextual prompt
            prompt = f"""
Based on the document analysis provided below, please answer the specific question.

DOCUMENT CONTEXT:
{document_context['analysis']}

SPECIFIC QUESTION DETAILS:
- Question: {question_content}
- Expected Answer Column: {column_header}
- Row Context: {row_context}
- Question Type: {question_type}
- Sheet: {sheet_name}

INSTRUCTION:
Please provide a specific, concise answer to fill in the empty cell in the "{column_header}" column for this question. 
Base your answer ONLY on the information available in the document analysis above.
If the document doesn't contain relevant information, respond with "Information not available in document".

Answer:
"""
            
            model = GenerativeModel("gemini-2.0-flash")
            
            response = model.generate_content(
                [Part.from_text(prompt)],
                generation_config={
                    "max_output_tokens": 1000,  # Shorter responses for cell content
                    "temperature": 0.3,
                    "top_p": 0.9,
                },
                safety_settings=[
                    SafetySetting(
                        category=SafetySetting.HarmCategory.HARM_CATEGORY_HATE_SPEECH,
                        threshold=SafetySetting.HarmBlockThreshold.OFF
                    ),
                    SafetySetting(
                        category=SafetySetting.HarmCategory.HARM_CATEGORY_DANGEROUS_CONTENT,
                        threshold=SafetySetting.HarmBlockThreshold.OFF
                    ),
                    SafetySetting(
                        category=SafetySetting.HarmCategory.HARM_CATEGORY_SEXUALLY_EXPLICIT,
                        threshold=SafetySetting.HarmBlockThreshold.OFF
                    ),
                    SafetySetting(
                        category=SafetySetting.HarmCategory.HARM_CATEGORY_HARASSMENT,
                        threshold=SafetySetting.HarmBlockThreshold.OFF
                    ),
                ],
                stream=False
            )
            
            ai_answer = response.text.strip()
            
            filled_responses.append({
                "sheet_name": sheet_name,
                "target_row": cell["target_row"],
                "target_column": cell["target_column"],
                "target_column_name": column_header,
                "question_content": question_content,
                "ai_response": ai_answer,
                "question_type": question_type,
                "row_context": row_context
            })
            
            print(f"✅ Generated response for {sheet_name} [{cell['target_row']},{cell['target_column']}]: {ai_answer[:50]}...")
            
        except Exception as e:
            print(f"❌ Error generating response for cell [{cell['target_row']},{cell['target_column']}]: {str(e)}")
            filled_responses.append({
                "sheet_name": sheet_name,
                "target_row": cell["target_row"],
                "target_column": cell["target_column"],
                "target_column_name": cell["target_column_name"],
                "question_content": cell["question_content"],
                "ai_response": f"Error: {str(e)}",
                "error": True
            })
    
    return filled_responses

def fill_target_cells_with_responses(df, responses):
    """
    Fill the target cells in DataFrame with AI-generated responses
    """
    processed_df = df.copy()
    
    for response in responses:
        if "error" in response and response["error"]:
            continue
            
        row_idx = response["target_row"]
        col_idx = response["target_column"]
        
        if row_idx < len(processed_df) and col_idx < len(processed_df.columns):
            ai_response = response["ai_response"]
            # Limit response length for readability
            if len(ai_response) > 500:
                ai_response = ai_response[:497] + "..."
            processed_df.iloc[row_idx, col_idx] = ai_response
            
    return processed_df

# All old functions have been replaced with new AI-based direct upload functions

@app.route("/api/v2/download-processed-document/<filename>", methods=["GET", "OPTIONS"])
def download_processed_document(filename):
    """Download the processed document"""
    if request.method == "OPTIONS":
        return "", 200
    
    try:
        file_path = os.path.join(UPLOAD_FOLDER, filename)
        if not os.path.exists(file_path):
            return jsonify({"error": "File not found"}), 404
        
        # Return file for download
        return send_file(
            file_path,
            as_attachment=True,
            download_name=filename,
            mimetype='application/octet-stream'
        )
        
    except Exception as e:
        print(f"❌ Download error: {str(e)}")
        return jsonify({"error": "Download failed"}), 500


@app.route("/api/prototype/sheet-analyzer", methods=["POST", "OPTIONS"])
def prototype_sheet_analyzer():
    """
    Prototype endpoint for analyzing Excel/CSV files sheet by sheet.
    Identifies questions, answer types, and column mappings per sheet.
    """
    if request.method == "OPTIONS":
        return "", 200
    
    if not VERTEX_AVAILABLE:
        return jsonify({"error": "Vertex AI service is not available"}), 503
    
    try:
        # Handle file upload
        if "file" not in request.files:
            return jsonify({"error": "No file provided"}), 400
        
        file = request.files["file"]
        if file.filename == "":
            return jsonify({"error": "No file selected"}), 400
        
        # Get file extension
        filename = file.filename.lower()
        file_ext = filename.split('.')[-1]
        
        if file_ext not in ['csv', 'xlsx', 'xls']:
            return jsonify({"error": "Only CSV and Excel files are supported"}), 400
        
        # Save temporary file
        temp_filename = f"temp_{uuid.uuid4()}.{file_ext}"
        temp_filepath = os.path.join("uploads", temp_filename)
        os.makedirs("uploads", exist_ok=True)
        file.save(temp_filepath)
        
        try:
            # Process based on file type
            if file_ext == 'csv':
                result = analyze_csv_for_prototype(temp_filepath, file.filename)
            else:  # xlsx or xls
                result = analyze_excel_for_prototype(temp_filepath, file.filename)
            
            return jsonify(result)
            
        finally:
            # Clean up temp file
            if os.path.exists(temp_filepath):
                os.remove(temp_filepath)
                
    except Exception as e:
        print(f"❌ Prototype analyzer error: {str(e)}")
        return jsonify({"error": f"Processing failed: {str(e)}"}), 500


def analyze_csv_for_prototype(filepath, filename):
    """Analyze CSV file for prototype - identify questions and answer types"""
    import pandas as pd
    
    try:
        df = pd.read_csv(filepath)
        
        # Analyze the single sheet
        sheet_analysis = analyze_sheet_for_questions(df, "CSV_Sheet", filepath, filename)
        
        return {
            "filename": filename,
            "file_type": "csv",
            "total_sheets": 1,
            "sheets": [sheet_analysis]
        }
        
    except Exception as e:
        return {
            "filename": filename,
            "file_type": "csv",
            "error": f"CSV analysis failed: {str(e)}"
        }


def analyze_excel_for_prototype(filepath, filename):
    """Analyze Excel file for prototype - process each sheet separately"""
    import pandas as pd
    
    try:
        excel_file = pd.ExcelFile(filepath)
        sheet_names = excel_file.sheet_names
        
        sheets_analysis = []
        
        for sheet_name in sheet_names:
            try:
                df = pd.read_excel(filepath, sheet_name=sheet_name)
                sheet_analysis = analyze_sheet_for_questions(df, sheet_name, filepath, filename)
                sheets_analysis.append(sheet_analysis)
            except Exception as e:
                sheets_analysis.append({
                    "sheet_name": sheet_name,
                    "error": f"Sheet analysis failed: {str(e)}"
                })
        
        return {
            "filename": filename,
            "file_type": "excel",
            "total_sheets": len(sheet_names),
            "sheets": sheets_analysis
        }
        
    except Exception as e:
        return {
            "filename": filename,
            "file_type": "excel",
            "error": f"Excel analysis failed: {str(e)}"
        }


def analyze_sheet_for_questions(df, sheet_name, filepath, filename):
    """Analyze a single sheet/dataframe to identify questions and answer types"""
    from vertexai.generative_models import GenerativeModel, Part, SafetySetting
    
    try:
        # Convert dataframe to text representation for LLM analysis
        sheet_text = f"Sheet Name: {sheet_name}\n\n"
        sheet_text += f"Headers: {list(df.columns)}\n\n"
        sheet_text += "Sample Data (first 10 rows):\n"
        sheet_text += df.head(10).to_string(index=True)
        
        # Create analysis prompt
        analysis_prompt = f"""
Analyze this spreadsheet sheet and determine:

1. QUESTION DETECTION:
   - Does this sheet contain questions that need answers?
   - If yes, what are the specific questions?

2. QUESTION EXTRACTION:
   - List all questions found in the sheet
   - Identify which columns/headers represent questions
   - Extract the exact question text

3. ANSWER TYPE CLASSIFICATION:
   - For each question, determine the expected answer type:
     * YES_NO: Yes/No or True/False responses
     * MULTIPLE_CHOICE: Selection from predefined options
     * SHORT_TEXT: Brief text responses (1-2 words)
     * LONG_TEXT: Detailed text responses (sentences/paragraphs)
     * NUMBER: Numeric values
     * DATE: Date values
     * COMPLETED_STATUS: Completed/Not Completed, Done/Pending, etc.
     * RATING: Scale ratings (1-5, 1-10, etc.)
     * OTHER: Any other specific type

4. COLUMN MAPPING:
   - Which columns contain questions vs answers
   - Which columns are related to each question
   - Multiple answer columns for single questions (if any)

Sheet Data:
{sheet_text}

Please respond in JSON format:
{{
    "has_questions": boolean,
    "total_questions": number,
    "questions": [
        {{
            "question_text": "exact question text",
            "question_column": "column name containing question",
            "answer_columns": ["column name(s) for answers"],
            "answer_type": "answer type from list above",
            "is_multiple_answers": boolean,
            "sample_answers": ["sample answer values if available"]
        }}
    ],
    "sheet_purpose": "brief description of what this sheet is for"
}}
"""

        # Use Vertex AI to analyze
        model = GenerativeModel("gemini-2.0-flash")
        
        response = model.generate_content(
            [Part.from_text(analysis_prompt)],
            generation_config={
                "temperature": 0.1,
                "top_p": 0.8,
                "max_output_tokens": 2048,
            },
            safety_settings=[
                SafetySetting(
                    category=SafetySetting.HarmCategory.HARM_CATEGORY_DANGEROUS_CONTENT,
                    threshold=SafetySetting.HarmBlockThreshold.OFF
                ),
                SafetySetting(
                    category=SafetySetting.HarmCategory.HARM_CATEGORY_HATE_SPEECH,
                    threshold=SafetySetting.HarmBlockThreshold.OFF
                ),
                SafetySetting(
                    category=SafetySetting.HarmCategory.HARM_CATEGORY_SEXUALLY_EXPLICIT,
                    threshold=SafetySetting.HarmBlockThreshold.OFF
                ),
                SafetySetting(
                    category=SafetySetting.HarmCategory.HARM_CATEGORY_HARASSMENT,
                    threshold=SafetySetting.HarmBlockThreshold.OFF
                ),
            ],
            stream=False
        )
        
        # Parse AI response
        ai_analysis = response.text.strip()
        
        # Try to extract JSON from AI response
        try:
            import json
            import re
            
            # Find JSON in the response
            json_match = re.search(r'\{.*\}', ai_analysis, re.DOTALL)
            if json_match:
                analysis_json = json.loads(json_match.group())
            else:
                # Fallback if no JSON found
                analysis_json = {
                    "has_questions": False,
                    "total_questions": 0,
                    "questions": [],
                    "sheet_purpose": "Analysis failed - no structured response"
                }
        except:
            analysis_json = {
                "has_questions": False,
                "total_questions": 0,
                "questions": [],
                "sheet_purpose": "JSON parsing failed"
            }
        
        # Add sheet metadata
        result = {
            "sheet_name": sheet_name,
            "total_columns": len(df.columns),
            "total_rows": len(df),
            "columns": list(df.columns),
            "analysis": analysis_json,
            "raw_ai_response": ai_analysis
        }
        
        return result
        
    except Exception as e:
        return {
            "sheet_name": sheet_name,
            "error": f"Sheet analysis failed: {str(e)}",
            "total_columns": len(df.columns) if 'df' in locals() else 0,
            "total_rows": len(df) if 'df' in locals() else 0
        }

# ================================
# PHASE 1: ENHANCED ENDPOINTS WITH QUEUE SUPPORT
# ================================

def upload_with_queue():
    """Handle upload using queue system"""
    try:
        if 'file' not in request.files:
            return jsonify({"error": "No file uploaded"}), 400
        
        file = request.files['file']
        if file.filename == '':
            return jsonify({"error": "No file selected"}), 400
        
        # Extract parameters
        org_id = request.form.get('org_id', '')
        user_id = request.form.get('user_id', '')
        file_id = request.form.get('file_id', str(uuid.uuid4()))
        
        print(f"🚀 Queue-based upload: {file.filename} (FileID: {file_id})")
        
        # Save file temporarily
        filename = file.filename
        save_path = os.path.abspath(os.path.join(UPLOAD_FOLDER, f"{file_id}_{filename}"))
        file.save(save_path)
        
        # Prepare file data for queue
        file_data = {
            "file_path": save_path,
            "filename": filename,
            "file_id": file_id
        }
        
        upload_params = {
            "org_id": org_id,
            "user_id": user_id,
            "file_id": file_id,
            "filename": filename
        }
        
        # Submit to queue
        task = process_file_upload.delay(file_data, upload_params)
        
        # Set initial progress in Redis
        redis_manager.set_progress(
            task.id, "queued", 0, "File uploaded, queued for processing", filename
        )
        
        return jsonify({
            "success": True,
            "message": "File queued for processing",
            "task_id": task.id,
            "file_id": file_id,
            "queue_mode": True,
            "status_endpoint": f"/task-status/{task.id}"
        })
        
    except Exception as e:
        error_msg = str(e)
        print(f"❌ Queue upload error: {error_msg}")
        traceback.print_exc()
        return jsonify({"error": "Upload failed", "details": error_msg}), 500

@app.route("/task-status/<task_id>", methods=["GET"])
def get_task_status_enhanced(task_id):
    """Get comprehensive task status - works with both queue and traditional progress"""
    try:
        # First check if it's a queue task
        if QUEUE_AVAILABLE and redis_manager.is_connected():
            progress_data = redis_manager.get_progress(task_id)
            if progress_data:
                # Get additional Celery status
                status_info = get_task_status(task_id)
                return jsonify({
                    "success": True,
                    "task_id": task_id,
                    "queue_mode": True,
                    "status": progress_data.get('status', 'unknown'),
                    "progress": progress_data.get('progress', 0),
                    "stage": progress_data.get('stage', ''),
                    "filename": progress_data.get('filename', ''),
                    "updated_at": progress_data.get('updated_at', ''),
                    "celery_info": status_info
                })
        
        # Fall back to traditional progress tracking
        with progress_lock:
            if task_id in upload_progress:
                progress_data = upload_progress[task_id]
                return jsonify({
                    "success": True,
                    "task_id": task_id,
                    "queue_mode": False,
                    **progress_data
                })
        
        return jsonify({
            "success": False,
            "error": "Task not found",
            "task_id": task_id
        }), 404
        
    except Exception as e:
        return jsonify({
            "error": "Failed to get task status",
            "details": str(e)
        }), 500

# Helper functions for queue integration
def process_file_upload_internal(file_data, upload_params, task_id):
    """
    Internal function to process file uploads
    This preserves existing functionality while being callable from tasks
    """
    try:
        file_path = file_data["file_path"]
        filename = file_data["filename"] 
        file_id = file_data["file_id"]
        
        org_id = upload_params["org_id"]
        user_id = upload_params["user_id"]
        
        # Update progress
        if QUEUE_AVAILABLE:
            redis_manager.set_progress(
                task_id, "processing", 25, "Extracting content", filename
            )
        
        # Get file extension
        file_ext = filename.split('.')[-1].lower() if '.' in filename else ''
        
        # Extract content using existing function
        chunks = parse_and_chunk(file_path, file_ext, chunk_size=50, max_chunks=500)
        
        if not chunks:
            raise Exception("No content extracted from file")
        
        # Update progress
        if QUEUE_AVAILABLE:
            redis_manager.set_progress(
                task_id, "embedding", 50, "Generating embeddings", filename
            )
        
        # Generate embeddings using existing function
        embeddings = embed_chunks(chunks, upload_id=file_id, org_id=org_id, filename=filename)
        
        if not embeddings:
            raise Exception("Failed to generate embeddings")
        
        # Store in Firestore using existing functionality
        if QUEUE_AVAILABLE:
            redis_manager.set_progress(
                task_id, "storing", 75, "Storing in database", filename
            )
        
        # Use existing Firestore storage logic
        store_embeddings_in_firestore(embeddings, chunks, file_id, org_id, filename)
        
        # Send webhook notification
        if user_id and QUEUE_AVAILABLE:
            webhook_task = send_webhook.delay(
                f"{os.environ.get('BACKEND_API_URL', 'http://localhost:8080')}/api/v2/files/webhook/ai-status",
                {
                    "fileId": file_id,
                    "userId": user_id,
                    "status": "completed",
                    "embeddingComplete": True,
                    "timestamp": datetime.now().isoformat(),
                    "source": "flask_ai_v3"
                }
            )
            print(f"📤 Webhook queued: {webhook_task.id}")
        
        # Clean up temporary file
        if os.path.exists(file_path):
            os.remove(file_path)
            print(f"🗑️ Cleaned up: {file_path}")
        
        return {
            "success": True,
            "file_id": file_id,
            "chunks_processed": len(chunks),
            "embeddings_generated": len(embeddings),
            "filename": filename
        }
        
    except Exception as e:
        # Clean up on error
        if 'file_path' in locals() and os.path.exists(file_path):
            os.remove(file_path)
        raise e

def store_embeddings_in_firestore(embeddings, chunks, file_id, org_id, filename):
    """Store embeddings in Firestore using existing logic"""
    global db
    
    if not FIREBASE_AVAILABLE or not db:
        raise Exception("Firestore not available")
    
    try:
        collection_ref = db.collection("document_embeddings")
        
        # Use batch operations for better performance
        batch = db.batch()
        batch_count = 0
        
        for i, (embedding, chunk) in enumerate(zip(embeddings, chunks)):
            doc_data = {
                "file_id": file_id,
                "org_id": org_id,
                "filename": filename,
                "chunk_index": i,
                "chunk_text": chunk,
                "embedding": embedding,
                "timestamp": datetime.now(),
                "created_at": datetime.now().isoformat(),
                "source": "flask_ai_v3"
            }
            
            doc_ref = collection_ref.document()
            batch.set(doc_ref, doc_data)
            batch_count += 1
            
            # Commit in batches of 450 (Firestore limit is 500)
            if batch_count >= 450:
                batch.commit()
                print(f"📦 Committed batch of {batch_count} embeddings")
                batch = db.batch()
                batch_count = 0
        
        # Commit remaining items
        if batch_count > 0:
            batch.commit()
            print(f"📦 Committed final batch of {batch_count} embeddings")
        
        print(f"✅ Stored {len(embeddings)} embeddings for {filename}")
        
    except Exception as e:
        print(f"❌ Firestore storage error: {e}")
        raise e

@app.route("/api/v3/chat", methods=["POST", "OPTIONS"])
def chat_v3():
    """V3 Chat endpoint - Simple query with smart defaults"""
    if request.method == "OPTIONS":
        return "", 200
    
    try:
        # Get request data
        data = request.get_json()
        if not data:
            return jsonify({"error": "No JSON data provided"}), 400
        
        query = data.get("query")
        org_id = data.get("orgId")  # Get orgId from request
        conversation_history = data.get("conversationHistory", "")  # Get conversation history
        
        print(f"💬 V3 Chat query: '{query[:100]}...'")
        print(f"🏢 V3 Chat orgId: {org_id}")
        print(f"📜 V3 Chat conversation history: {conversation_history[:100]}..." if conversation_history else "📜 No conversation history")
        
        if not query:
            return jsonify({"error": "Query is required"}), 400
        
        # Use provided orgId or default to None (search all orgs)
        user_id = "api-user"  # Default user
        file_ids = None  # Search all files
        max_results = 5  # Good default
        include_context = True  # Always include context for chat
        
        # Import retrieval and LLM systems
        from retrieval_system import search_documents, get_context
        from llm_integration import generate_rag_answer
        
        # Optimized: Single search call for both context and results
        print("🎯 Getting context for query...")
        search_results = search_documents(
            query=query,
            org_id=org_id,
            file_ids=file_ids,
            top_k=max_results
        )
        
        # Build context from search results (faster than separate context call)
        context_result = {
            "context": "\n\n".join([result["text"][:500] for result in search_results[:3]]),  # Use top 3 results, 500 chars each
            "sources": list(set([result["file_id"] for result in search_results])),
            "total_chunks": len(search_results),
            "avg_similarity": sum([result["similarity_score"] for result in search_results]) / len(search_results) if search_results else 0,
            "chunks_metadata": search_results
        }
        
        # Generate LLM answer using retrieved context
        print("🤖 Generating LLM answer...")
        
        # Prepare enhanced context with conversation history
        enhanced_context = context_result["context"]
        if conversation_history:
            enhanced_context = f"Previous conversation:\n{conversation_history}\n\nRelevant context:\n{context_result['context']}"
        
        llm_result = generate_rag_answer(
            query=query,
            context=enhanced_context,
            max_length=1024,
            temperature=0.7
        )
        
        if llm_result["success"]:
            # Clean the LLM response to extract only the answer
            raw_answer = llm_result["answer"]
            cleaned_answer = raw_answer
            
            # Try to extract just the assistant's answer from the full response
            if "assistant" in raw_answer and raw_answer.count("assistant") > 0:
                # Find the last occurrence of 'assistant' and extract what comes after
                parts = raw_answer.split("assistant")
                if len(parts) > 1:
                    cleaned_answer = parts[-1].strip()
            
            # Remove any remaining system/user prefixes and common artifacts
            for prefix in ["system\n", "user\n", "Answer:\n", "Answer:", "\nassistant\n"]:
                if cleaned_answer.startswith(prefix):
                    cleaned_answer = cleaned_answer[len(prefix):].strip()
            
            # Remove trailing artifacts
            for suffix in ["<|im_end|>", "<|endoftext|>"]:
                if cleaned_answer.endswith(suffix):
                    cleaned_answer = cleaned_answer[:-len(suffix)].strip()
            
            print(f"🧹 Cleaned answer: {cleaned_answer[:100]}...")
            
            response = {
                "query": query,
                "orgId": org_id,
                "conversationHistory": conversation_history,
                "answer": cleaned_answer,
                "sources": context_result["sources"],
                "search_results": search_results,
                "context_metadata": {
                    "sources": context_result["sources"],
                    "total_chunks": context_result["total_chunks"],
                    "avg_similarity": context_result["avg_similarity"],
                    "chunks_metadata": context_result["chunks_metadata"]
                },
                "llm_metadata": {
                    "model": llm_result["model"],
                    "context_used": llm_result["context_used"],
                    "context_length": llm_result["context_length"]
                },
                "retrieval_method": "vector_similarity",
                "embedding_model": "custom_hf_endpoint",
                "pipeline": "complete_rag"
            }
        else:
            # Fallback if LLM fails - return retrieval results
            print("⚠️ LLM failed, returning retrieval results only")
            response = {
                "query": query,
                "orgId": org_id,
                "conversationHistory": conversation_history,
                "answer": f"I found {len(search_results)} relevant results, but couldn't generate a complete answer. Please check the search results below.",
                "sources": context_result["sources"],
                "search_results": search_results,
                "context_metadata": {
                    "sources": context_result["sources"],
                    "total_chunks": context_result["total_chunks"],
                    "avg_similarity": context_result["avg_similarity"],
                    "chunks_metadata": context_result["chunks_metadata"]
                },
                "llm_error": llm_result.get("error", "Unknown LLM error"),
                "retrieval_method": "vector_similarity",
                "embedding_model": "custom_hf_endpoint",
                "pipeline": "retrieval_only"
            }
        
        print(f"✅ V3 Chat response: {len(search_results) if search_results else 0} results")
        return jsonify(response), 200
        
    except Exception as e:
        print(f"❌ V3 Chat error: {str(e)}")
        return jsonify({"error": f"Chat failed: {str(e)}"}), 500

@app.route("/api/v3/search", methods=["POST", "OPTIONS"])
def search_v3():
    """V3 Search endpoint - Simple query with optimal defaults"""
    if request.method == "OPTIONS":
        return "", 200
    
    try:
        # Get request data
        data = request.get_json()
        if not data:
            return jsonify({"error": "No JSON data provided"}), 400
        
        query = data.get("query")
        
        print(f"🔍 V3 Search: '{query[:100]}...'")
        
        if not query:
            return jsonify({"error": "Query is required"}), 400
        
        # Optimized defaults for fast and accurate results
        org_id = None  # Search all organizations
        file_ids = None  # Search all files
        top_k = 5  # Fewer results for faster response
        similarity_threshold = 0.7  # Higher threshold for better quality
        
        # Import and run search
        from retrieval_system import NeonDBRetrieval
        import asyncio
        
        retrieval = NeonDBRetrieval()
        results = asyncio.run(retrieval.search_similar_chunks(
            query=query,
            org_id=org_id,
            file_ids=file_ids,
            top_k=top_k,
            similarity_threshold=similarity_threshold
        ))
        
        response = {
            "query": query,
            "results": results,
            "total_results": len(results),
            "search_params": {
                "top_k": top_k,
                "similarity_threshold": similarity_threshold,
                "org_id": org_id,
                "file_ids": file_ids
            },
            "retrieval_method": "vector_similarity",
            "embedding_model": "custom_hf_endpoint"
        }
        
        print(f"✅ V3 Search response: {len(results)} results")
        return jsonify(response), 200
        
    except Exception as e:
        print(f"❌ V3 Search error: {str(e)}")
        return jsonify({"error": f"Search failed: {str(e)}"}), 500


# Keep-alive service for Render server
def keep_alive_service():
    """Background service that pings the server to keep it alive"""
    global keep_alive_running
    
    # Get configuration
    interval = int(os.environ.get("KEEP_ALIVE_INTERVAL", "840"))  # 14 minutes default
    port = int(os.environ.get("PORT", "8002"))
    base_url = f"http://localhost:{port}"
    
    print(f"🔄 Keep-alive service started (interval: {interval}s)")
    
    while keep_alive_running:
        try:
            # Wait for the interval
            time.sleep(interval)
            
            if not keep_alive_running:
                break
                
            # Make a request to our own health endpoint
            response = requests.get(f"{base_url}/health", timeout=30)
            if response.status_code == 200:
                print(f"✅ Keep-alive ping successful at {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
            else:
                print(f"⚠️ Keep-alive ping returned status {response.status_code}")
                
        except requests.exceptions.RequestException as e:
            print(f"⚠️ Keep-alive ping failed: {str(e)}")
        except Exception as e:
            print(f"❌ Keep-alive service error: {str(e)}")
    
    print("🛑 Keep-alive service stopped")


def start_keep_alive_service():
    """Start the keep-alive service in a background thread"""
    global keep_alive_thread, keep_alive_running
    
    # Check if keep-alive is enabled
    enabled = os.environ.get("KEEP_ALIVE_ENABLED", "false").lower() == "true"
    debug_mode = os.environ.get("DEBUG", "0") == "1"
    
    if not enabled:
        print("⏸️ Keep-alive service disabled (set KEEP_ALIVE_ENABLED=true to enable)")
        return
    
    if debug_mode:
        print("⏸️ Keep-alive service disabled in debug mode")
        return
        
    if keep_alive_thread and keep_alive_thread.is_alive():
        print("⚠️ Keep-alive service already running")
        return
    
    keep_alive_running = True
    keep_alive_thread = Thread(target=keep_alive_service, daemon=True)
    keep_alive_thread.start()
    print("🚀 Keep-alive service started in background")


def stop_keep_alive_service():
    """Stop the keep-alive service"""
    global keep_alive_running
    keep_alive_running = False
    print("🛑 Keep-alive service stopping...")


if __name__ == "__main__":
    port = int(os.environ.get("PORT", 8002))
    debug_mode = os.environ.get("DEBUG", "0") == "1"
    
    print(f"🚀 CareAI API v2.3.1 - Enhanced with Prototype Sheet Analyzer")
    print(f"📋 Features: FileId management, Auto status sync, Multimodal question extraction, Direct file upload, Sheet Analyzer")
    print(f"🎯 Supported: Documents (PDF, DOCX, TXT, CSV, JSON, XLSX) + Images (PNG, JPG, GIF, WEBP)")
    print(f"🔬 New: /api/prototype/sheet-analyzer - Analyze Excel/CSV sheets for questions and answer types")
    
    if debug_mode:
        print(f"🔍 Starting Flask development server on port {port}")
        app.run(host="0.0.0.0", port=port, debug=True)
    else:
        threads = int(os.environ.get("WAITRESS_THREADS", "8"))
        print(f"🚀 Starting server with Waitress on port {port} with {threads} threads")
        
        # Start keep-alive service for production
        start_keep_alive_service()
        
        try:
            serve(app, host="0.0.0.0", port=port, threads=threads)
        except KeyboardInterrupt:
            print("\n🛑 Shutting down server...")
            stop_keep_alive_service()
        except Exception as e:
            print(f"❌ Server error: {e}")
            stop_keep_alive_service()